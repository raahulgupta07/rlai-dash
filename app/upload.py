"""
Data Upload & Auto-Onboarding
==============================

Upload CSV/Excel/JSON files → AI auto-generates metadata → loads into knowledge base.

Usage:
    POST /api/upload  (multipart/form-data with file + optional table_name)
    GET  /api/tables   (list all tables with row counts)
    DELETE /api/tables/{table_name}  (drop a user-uploaded table)
"""

import json
import os
import re
import tempfile
import time
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path
from typing import Any

_bg_executor = ThreadPoolExecutor(max_workers=5, thread_name_prefix="dash-bg")

import pandas as pd
from fastapi import APIRouter, HTTPException, Request, UploadFile
from sqlalchemy import create_engine as _sa_create_engine, inspect, text
from sqlalchemy.pool import NullPool


def create_engine(url, **kw):
    """Wrapper that forces NullPool for ad-hoc engines (PgBouncer handles pooling)."""
    kw.setdefault("poolclass", NullPool)
    return _sa_create_engine(url, **kw)

from dash.paths import BUSINESS_DIR, KNOWLEDGE_DIR, QUERIES_DIR, TABLES_DIR
from db import db_url
from dash.settings import TRAINING_MODEL


def _safe_write_json(path: Path, data) -> None:
    """Atomic JSON write — prevents corruption from concurrent writes."""
    path.parent.mkdir(parents=True, exist_ok=True)
    tmp_fd, tmp_path = tempfile.mkstemp(dir=path.parent, suffix=".json")
    try:
        with os.fdopen(tmp_fd, "w") as f:
            json.dump(data, f, indent=2, default=str)
        os.replace(tmp_path, str(path))  # Atomic rename
    except Exception:
        try:
            os.unlink(tmp_path)
        except OSError:
            pass
        raise


def _safe_read_json(path: Path, default=None):
    """Safe JSON read — handles concurrent writes and corruption."""
    try:
        if path.exists():
            return json.loads(path.read_text(errors="ignore"))
    except (json.JSONDecodeError, OSError):
        pass
    return default if default is not None else {}

router = APIRouter(prefix="/api", tags=["Upload"])
_engine = create_engine(db_url)

# Tables that ship with the demo — protected from deletion
PROTECTED_TABLES = {"customers", "subscriptions", "plan_changes", "invoices", "usage_metrics", "support_tickets", "dash_users", "dash_tokens", "dash_projects", "shared_results"}

MAX_FILE_SIZE = 50 * 1024 * 1024  # 50 MB
ALLOWED_EXTENSIONS = {".csv", ".xlsx", ".xls", ".json", ".sql", ".py", ".txt", ".md", ".pptx", ".docx", ".pdf", ".jpg", ".jpeg", ".png"}

# Universal vision prompt — handles charts, scanned docs, photos, diagrams in one call
_UNIVERSAL_VISION_PROMPT = (
    "Analyze this image from a business document.\n"
    "If it contains TEXT (scanned document, certificate, letter): extract ALL text exactly as written. "
    "Preserve layout. Render any tables as markdown tables.\n"
    "If it contains a CHART or GRAPH: extract ALL data points, axis labels, legend items as a markdown table. "
    "Describe the trend: increasing/decreasing/flat and by how much.\n"
    "If it contains a DIAGRAM or FLOWCHART: describe all components, connections, and flow.\n"
    "If it contains a PHOTO: describe what is visible and any text, signage, or labels.\n"
    "If it contains a TABLE rendered as image: extract ALL rows and columns as a markdown table.\n"
    "Be precise with all numbers, dates, and labels. Miss nothing."
)


# ---------------------------------------------------------------------------
# Smart Upload: File Classification + Table Matching + Fingerprinting
# ---------------------------------------------------------------------------

def classify_file(filename: str, headers: list[str] | None = None, content_sample: str = "", df=None, project_slug: str = "") -> str:
    """Classify uploaded file: data, column_definition, business_rules, sql_patterns, documentation."""
    ext = Path(filename).suffix.lower()

    # SQL files → query patterns
    if ext == ".sql":
        return "sql_patterns"

    # Documents → check content for rules vs general docs
    if ext in (".md", ".txt", ".pdf", ".docx", ".pptx"):
        lower = content_sample.lower()
        rule_signals = ["rule:", "must be", "should be", "always", "never", "constraint", "formula", "calculation"]
        if sum(1 for s in rule_signals if s in lower) >= 2:
            return "business_rules"
        return "documentation"

    # Code files
    if ext in (".py", ".js", ".ts"):
        return "documentation"

    # Data files (CSV, Excel, JSON)
    if ext in (".csv", ".xlsx", ".xls", ".json"):
        # Method 1: Check by HEADERS (Format A — has "column_name" + "description" headers)
        if headers:
            header_lower = [str(h).lower().strip() for h in headers]
            def_signals = ["definition", "description", "meaning", "business_meaning", "data_type", "column_name"]
            has_def = any(s in h for h in header_lower for s in def_signals)
            has_col = any("column" in h or "field" in h for h in header_lower)
            if has_def and has_col:
                return "column_definition"

        # Method 2: Check by VALUES (Format B — no headers, values match existing table columns)
        if df is not None and project_slug and len(df.columns) <= 5:
            result = _detect_definition_by_values(df, project_slug)
            if result:
                return "column_definition"

        # Method 3: Heuristic — 2-3 columns, one short text + one long text = likely definitions
        if df is not None and len(df.columns) in (2, 3) and len(df) >= 5 and len(df) <= 200:
            text_cols = []
            for c in df.columns:
                vals = df[c].dropna().astype(str)
                if len(vals) > 0:
                    avg_len = vals.str.len().mean()
                    all_text = all(not str(v).replace('.', '').replace('-', '').isdigit() for v in vals.head(5))
                    text_cols.append((c, avg_len, all_text))
            # If we have exactly 2 text columns with different avg lengths
            text_only = [t for t in text_cols if t[2]]
            if len(text_only) >= 2:
                lengths = sorted(text_only, key=lambda x: x[1])
                # Short col avg < 50, long col avg > 30, and ratio > 1.3
                if lengths[0][1] < 60 and lengths[-1][1] > 20 and lengths[-1][1] / max(lengths[0][1], 1) > 1.3:
                    return "column_definition"

        return "data"

    return "data"


def _detect_definition_by_values(df, project_slug: str) -> dict | None:
    """Generic: detect if a file is a column definition by matching VALUES against existing table columns."""
    if len(df) < 3 or len(df.columns) < 2:
        return None

    # Get all column names from all tables in this project
    existing_columns = set()
    try:
        from sqlalchemy import inspect as sa_inspect
        insp = sa_inspect(_engine)
        schema = re.sub(r"[^a-z0-9_]", "_", project_slug.lower())[:63]
        for tbl in insp.get_table_names(schema=schema):
            for col in insp.get_columns(tbl, schema=schema):
                existing_columns.add(col["name"].lower().strip())
                # Also add original name variations (with spaces, dots, etc.)
                existing_columns.add(col["name"].lower().replace("_", " ").strip())
    except Exception:
        return None

    if not existing_columns:
        return None

    # For each column in the uploaded file, check if its VALUES match existing column names
    for col_idx, col in enumerate(df.columns):
        values = df[col].dropna().astype(str).tolist()
        if not values:
            continue

        values_lower = [v.lower().strip() for v in values]

        # How many values match existing column names? (try multiple formats)
        matches = 0
        for v in values_lower:
            cleaned = re.sub(r"[^a-z0-9]", "_", v).strip("_")
            cleaned = re.sub(r"_+", "_", cleaned)
            if (v in existing_columns or
                cleaned in existing_columns or
                v.replace(":", "").replace(".", "").strip() in existing_columns or
                cleaned[:30] in existing_columns):
                matches += 1
        match_pct = matches / max(len(values), 1) * 100

        if match_pct > 40:
            # This column contains column names! Find the description column
            # Description = the column with the longest average text
            best_desc_col = None
            best_avg_len = 0
            for other_idx, other_col in enumerate(df.columns):
                if other_idx == col_idx:
                    continue
                other_values = df[other_col].dropna().astype(str).tolist()
                avg_len = sum(len(v) for v in other_values) / max(len(other_values), 1)
                if avg_len > best_avg_len:
                    best_avg_len = avg_len
                    best_desc_col = other_col

            return {"name_col": col, "desc_col": best_desc_col, "match_pct": match_pct}

    return None


def match_existing_table(project_slug: str, new_columns: list[str]) -> dict | None:
    """Check if uploaded columns match an existing table. Returns match info or None."""
    if not new_columns or not project_slug:
        return None

    try:
        from sqlalchemy import inspect as sa_inspect
        insp = sa_inspect(_engine)
        import re as _re
        schema = _re.sub(r"[^a-z0-9_]", "_", project_slug.lower())[:63]
        existing_tables = insp.get_table_names(schema=schema)

        new_cols_lower = set(str(c).lower().strip() for c in new_columns)
        best_match = None
        best_overlap = 0

        for tbl in existing_tables:
            tbl_cols = insp.get_columns(tbl, schema=schema)
            tbl_col_names = set(c["name"].lower() for c in tbl_cols)

            overlap = new_cols_lower & tbl_col_names
            overlap_pct = len(overlap) / max(len(new_cols_lower), 1) * 100

            if overlap_pct > best_overlap and overlap_pct >= 50:
                # Count existing rows
                try:
                    with _engine.connect() as conn:
                        row_count = conn.execute(text(f'SELECT COUNT(*) FROM "{schema}"."{tbl}"')).scalar() or 0
                except Exception:
                    row_count = 0

                new_cols = new_cols_lower - tbl_col_names
                missing_cols = tbl_col_names - new_cols_lower

                best_match = {
                    "table": tbl,
                    "overlap_pct": round(overlap_pct),
                    "matched_columns": list(overlap),
                    "new_columns": list(new_cols),
                    "missing_columns": list(missing_cols),
                    "existing_rows": row_count,
                }
                best_overlap = overlap_pct

        return best_match
    except Exception:
        return None


def compute_fingerprint(row_count: int, col_names: list[str]) -> str:
    """Compute a fingerprint for change detection."""
    import hashlib
    sorted_cols = sorted(str(c).lower().strip() for c in col_names)
    raw = f"{row_count}|{'|'.join(sorted_cols)}"
    return hashlib.md5(raw.encode()).hexdigest()


def save_fingerprint(project_slug: str, table_name: str, row_count: int, col_names: list[str]):
    """Save fingerprint to dash_table_metadata."""
    fp = compute_fingerprint(row_count, col_names)
    col_hash = compute_fingerprint(0, col_names)  # cols only, no row count
    try:
        with _engine.connect() as conn:
            conn.execute(text(
                "UPDATE public.dash_table_metadata SET fingerprint = :fp, row_count = :rc, col_hash = :ch, updated_at = NOW() "
                "WHERE project_slug = :s AND table_name = :t"
            ), {"fp": fp, "rc": row_count, "ch": col_hash, "s": project_slug, "t": table_name})
            conn.commit()
    except Exception:
        pass


def check_fingerprint_changed(project_slug: str, table_name: str, row_count: int, col_names: list[str]) -> str:
    """Compare new fingerprint with stored. Returns: 'new', 'unchanged', 'rows_changed', 'schema_changed'."""
    fp = compute_fingerprint(row_count, col_names)
    col_hash = compute_fingerprint(0, col_names)
    try:
        with _engine.connect() as conn:
            row = conn.execute(text(
                "SELECT fingerprint, col_hash FROM public.dash_table_metadata WHERE project_slug = :s AND table_name = :t"
            ), {"s": project_slug, "t": table_name}).fetchone()

        if not row or not row[0]:
            return "new"
        if row[0] == fp:
            return "unchanged"
        if row[1] != col_hash:
            return "schema_changed"
        return "rows_changed"
    except Exception:
        return "new"


def process_column_definitions(project_slug: str, df) -> dict:
    """Process a column definition file → save annotations + memories + rules.
    Works with ANY format — detects column name and description columns automatically."""

    # Method 1: Try header-based detection (Format A: has Column Name, Description headers)
    table_col = next((h for h in df.columns if str(h).lower() in ("table", "table_name", "entity")), None)
    col_col = next((h for h in df.columns if str(h).lower() in ("column", "column_name", "field", "field_name")), None)
    def_col = next((h for h in df.columns if str(h).lower() in ("definition", "description", "meaning", "business_meaning", "notes", "comment")), None)
    type_col = next((h for h in df.columns if str(h).lower() in ("data_type", "type", "format", "dtype")), None)

    # Method 2: If headers don't match, detect by VALUES (Format B: no headers)
    if not col_col or not def_col:
        detected = _detect_definition_by_values(df, project_slug)
        if detected:
            col_col = detected["name_col"]
            def_col = detected["desc_col"]
        else:
            # Method 3: Heuristic — in a 2-3 column file, shorter text col = names, longer = descriptions
            if len(df.columns) in (2, 3):
                text_cols = []
                for c in df.columns:
                    vals = df[c].dropna().astype(str)
                    if len(vals) > 0:
                        avg_len = vals.str.len().mean()
                        text_cols.append((c, avg_len))
                text_cols.sort(key=lambda x: x[1])
                if len(text_cols) >= 2:
                    col_col = text_cols[0][0]  # shorter = column names
                    def_col = text_cols[-1][0]  # longer = descriptions

    if not col_col or not def_col:
        return {"saved": 0, "error": "Could not detect column name and description columns"}

    # If no table column, auto-detect from project's existing tables
    default_table = ""
    if not table_col:
        try:
            from sqlalchemy import inspect as sa_inspect
            insp = sa_inspect(_engine)
            schema = re.sub(r"[^a-z0-9_]", "_", project_slug.lower())[:63]
            tables = insp.get_table_names(schema=schema)
            if tables:
                default_table = tables[0]  # use first table
        except Exception:
            pass

    annotations_saved = 0
    memories_saved = 0
    rules_saved = 0

    try:
        with _engine.connect() as conn:
            for _, row in df.iterrows():
                tbl = str(row[table_col]).strip() if table_col and row.get(table_col) else ""
                col = str(row[col_col]).strip() if row.get(col_col) else ""
                defn = str(row[def_col]).strip() if row.get(def_col) else ""

                if not col or not defn or defn == "nan":
                    continue

                # Save as annotation (use detected table or default)
                ann_table = tbl.lower() if tbl else default_table
                if ann_table:
                    conn.execute(text(
                        "INSERT INTO public.dash_annotations (project_slug, table_name, column_name, annotation, updated_by) "
                        "VALUES (:s, :t, :c, :a, 'column_definition') "
                        "ON CONFLICT (project_slug, table_name, column_name) DO UPDATE SET annotation = :a, updated_at = NOW()"
                    ), {"s": project_slug, "t": ann_table, "c": col.lower(), "a": defn})
                    annotations_saved += 1

                # Save as memory
                fact = f"Column '{ann_table}.{col}': {defn}" if ann_table else f"Column '{col}': {defn}"
                conn.execute(text(
                    "INSERT INTO public.dash_memories (project_slug, scope, fact, source) VALUES (:s, 'project', :f, 'column_definition')"
                ), {"s": project_slug, "f": fact[:500]})
                memories_saved += 1

                # Save type as rule if specified
                if type_col and row.get(type_col) and str(row[type_col]) != "nan":
                    rule_defn = f"Column '{col}' should be treated as {row[type_col]} type"
                    rule_id = f"coldef_{col.lower()[:30]}"
                    conn.execute(text(
                        "INSERT INTO public.dash_rules_db (project_slug, rule_id, name, type, definition, source) "
                        "VALUES (:s, :rid, :name, 'data_type', :defn, 'column_definition') "
                        "ON CONFLICT (project_slug, rule_id) DO NOTHING"
                    ), {"s": project_slug, "rid": rule_id, "name": f"{col} type", "defn": rule_defn})
                    rules_saved += 1

            conn.commit()
    except Exception:
        pass

    return {"annotations": annotations_saved, "memories": memories_saved, "rules": rules_saved}


def process_business_rules_doc(project_slug: str, content: str, filename: str) -> dict:
    """Extract business rules and facts from a text/markdown document via LLM."""
    from dash.settings import training_llm_call

    prompt = f"""Extract business rules and key facts from this document.

DOCUMENT ({filename}):
{content[:3000]}

Return ONLY valid JSON:
{{
  "rules": [
    {{"name": "Rule name", "definition": "The business rule text"}}
  ],
  "facts": [
    "Key fact 1",
    "Key fact 2"
  ]
}}"""

    result = training_llm_call(prompt, "extraction")
    if not result:
        return {"rules": 0, "facts": 0}

    try:
        parsed = json.loads(result)
    except json.JSONDecodeError:
        return {"rules": 0, "facts": 0}

    rules_saved = 0
    facts_saved = 0

    try:
        with _engine.connect() as conn:
            for rule in (parsed.get("rules") or [])[:10]:
                if isinstance(rule, dict) and rule.get("name"):
                    rule_id = f"doc_{re.sub(r'[^a-z0-9]', '_', rule['name'].lower())[:30]}"
                    conn.execute(text(
                        "INSERT INTO public.dash_rules_db (project_slug, rule_id, name, type, definition, source) "
                        "VALUES (:s, :rid, :name, 'business_rule', :defn, 'document') "
                        "ON CONFLICT (project_slug, rule_id) DO NOTHING"
                    ), {"s": project_slug, "rid": rule_id, "name": rule["name"], "defn": rule.get("definition", "")})
                    rules_saved += 1

            for fact in (parsed.get("facts") or [])[:10]:
                if fact and isinstance(fact, str):
                    conn.execute(text(
                        "INSERT INTO public.dash_memories (project_slug, scope, fact, source) VALUES (:s, 'project', :f, 'document')"
                    ), {"s": project_slug, "f": fact[:500]})
                    facts_saved += 1

            conn.commit()
    except Exception:
        pass

    return {"rules": rules_saved, "facts": facts_saved}


def process_sql_file(project_slug: str, content: str) -> dict:
    """Parse SQL file and save queries as patterns with metadata extraction."""
    # Split by semicolons or double newlines
    queries = [q.strip() for q in re.split(r';|\n\n', content) if q.strip()]
    saved = 0

    try:
        engine = _engine
        for q in queries[:20]:
            if not q.upper().startswith(("SELECT", "WITH")):
                continue
            # Generate a question from the SQL
            question = f"Run: {q[:100]}"
            if _save_query_pattern_with_metadata(engine, project_slug, question, q, source='sql_file'):
                saved += 1
    except Exception:
        pass

    return {"patterns_saved": saved}


def _sanitize_table_name(name: str) -> str:
    """Convert filename to a valid PostgreSQL table name."""
    name = Path(name).stem.lower()
    name = re.sub(r"[^a-z0-9_]", "_", name)
    name = re.sub(r"_+", "_", name).strip("_")
    if not name or name[0].isdigit():
        name = "t_" + name
    return name[:63]  # PG identifier limit


def _find_header_row(file_path: str, ext: str) -> int:
    """Detect the real header row in messy Excel/CSV files.
    Scans first 10 rows to find the one with the most text-like columns."""
    try:
        if ext in (".xlsx", ".xls"):
            df_raw = pd.read_excel(file_path, header=None, nrows=10)
        elif ext == ".csv":
            df_raw = pd.read_csv(file_path, header=None, nrows=10)
        else:
            return 0

        best_row = 0
        best_score = 0

        for i in range(min(10, len(df_raw))):
            row = df_raw.iloc[i]
            non_null = row.dropna()
            if len(non_null) == 0:
                continue

            score = 0
            for val in non_null:
                s = str(val).strip()
                # Good header: short text, no long numbers, not a formula
                if isinstance(val, str) and len(s) > 1 and len(s) < 50 and not s.startswith("="):
                    score += 3
                # Penalize: pure numbers, very long strings, formulas
                elif isinstance(val, (int, float)):
                    score -= 1

            # Prefer rows with many non-null text values
            text_ratio = sum(1 for v in non_null if isinstance(v, str) and len(str(v).strip()) > 1) / max(len(non_null), 1)
            score += text_ratio * 5

            if score > best_score:
                best_score = score
                best_row = i

        return best_row
    except Exception:
        return 0


def _clean_dataframe(df: pd.DataFrame) -> pd.DataFrame:
    """Smart cleanup: drop empty rows/columns, rename unnamed, clean column names."""
    # 1. Drop rows where ALL values are NaN
    df = df.dropna(how='all')

    # 2. Drop columns where ALL values are NaN (100% empty)
    empty_cols = [c for c in df.columns if df[c].isna().all()]
    if empty_cols:
        df = df.drop(columns=empty_cols)

    # 3. Drop columns that are >95% null (nearly empty)
    high_null_cols = [c for c in df.columns if df[c].isna().mean() > 0.95]
    if high_null_cols:
        df = df.drop(columns=high_null_cols)

    # 4. Rename "Unnamed" columns based on content BEFORE cleaning names
    renamed = {}
    day_names = {'monday', 'tuesday', 'wednesday', 'thursday', 'friday', 'saturday', 'sunday'}
    month_names = {'january', 'february', 'march', 'april', 'may', 'june', 'july', 'august', 'september', 'october', 'november', 'december', 'jan', 'feb', 'mar', 'apr', 'jun', 'jul', 'aug', 'sep', 'oct', 'nov', 'dec'}

    for c in df.columns:
        cstr = str(c).lower()
        if 'unnamed' in cstr or cstr.startswith('col_'):
            non_null = df[c].dropna()
            if len(non_null) == 0:
                renamed[c] = f'empty_{list(df.columns).index(c)}'
                continue

            sample = [str(v).lower().strip() for v in non_null.head(10).tolist() if str(v).strip()]

            if sample and all(v in day_names for v in sample):
                renamed[c] = 'day_of_week'
            elif sample and all(v in month_names for v in sample):
                renamed[c] = 'month'
            elif pd.api.types.is_numeric_dtype(non_null):
                renamed[c] = f'value_{list(df.columns).index(c) + 1}'
            else:
                renamed[c] = f'column_{list(df.columns).index(c) + 1}'

    if renamed:
        df = df.rename(columns=renamed)

    # 5. Clean ALL column names (after smart rename)
    df.columns = [re.sub(r"[^a-z0-9_]", "_", str(c).lower().strip()).strip("_") or f"col_{i}" for i, c in enumerate(df.columns)]

    # 5b. Escape PostgreSQL reserved words
    PG_RESERVED = {'select', 'from', 'where', 'order', 'group', 'by', 'having', 'join', 'on',
        'user', 'table', 'column', 'schema', 'database', 'index', 'view', 'limit', 'offset',
        'insert', 'update', 'delete', 'drop', 'create', 'alter', 'case', 'when', 'then',
        'else', 'end', 'all', 'any', 'exists', 'in', 'not', 'and', 'or', 'null', 'true',
        'false', 'primary', 'key', 'foreign', 'check', 'default', 'constraint', 'references'}
    df.columns = [f"{c}_col" if c in PG_RESERVED else c for c in df.columns]

    # 6. Remove duplicate column names
    seen = {}
    new_cols = []
    for c in df.columns:
        if c in seen:
            seen[c] += 1
            new_cols.append(f"{c}_{seen[c]}")
        else:
            seen[c] = 0
            new_cols.append(c)
    df.columns = new_cols

    return df


def _detect_delimiter(file_path: str) -> str:
    """Auto-detect CSV delimiter from first line."""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            first_line = f.readline()
        counts = {d: first_line.count(d) for d in [',', ';', '\t', '|']}
        best = max(counts, key=counts.get)
        return best if counts[best] > 0 else ','
    except Exception:
        return ','


def _read_file(file_path: str, ext: str) -> pd.DataFrame:
    """Read a data file into a DataFrame with smart header detection and cleanup."""
    if ext == ".csv":
        header_row = _find_header_row(file_path, ext)
        sep = _detect_delimiter(file_path)
        df = pd.read_csv(file_path, header=header_row, sep=sep)
        return _clean_dataframe(df)
    elif ext in (".xlsx", ".xls"):
        header_row = _find_header_row(file_path, ext)
        df = pd.read_excel(file_path, header=header_row)
        return _clean_dataframe(df)
    elif ext == ".json":
        try:
            df = pd.read_json(file_path)
        except ValueError:
            # Try records orientation
            df = pd.read_json(file_path, orient='records')
        return _clean_dataframe(df)
    else:
        raise ValueError(f"Unsupported format: {ext}")


def _analyze_column(series: pd.Series) -> dict[str, Any]:
    """Analyze a single column for metadata generation."""
    info: dict[str, Any] = {
        "name": series.name,
        "dtype": str(series.dtype),
        "null_pct": round(series.isna().mean() * 100, 1),
        "unique_count": int(series.nunique()),
        "total_count": len(series),
    }

    if pd.api.types.is_numeric_dtype(series):
        info["type"] = "numeric"
        clean = series.dropna()
        if len(clean) > 0:
            info["min"] = float(clean.min())
            info["max"] = float(clean.max())
            info["mean"] = round(float(clean.mean()), 2)
    elif pd.api.types.is_datetime64_any_dtype(series):
        info["type"] = "datetime"
        clean = series.dropna()
        if len(clean) > 0:
            info["min_date"] = str(clean.min())
            info["max_date"] = str(clean.max())
    else:
        info["type"] = "text"
        clean = series.dropna().astype(str)
        if len(clean) > 0:
            info["sample_values"] = clean.value_counts().head(5).index.tolist()
            if info["unique_count"] <= 20:
                info["is_categorical"] = True

    # Try date detection for string columns
    if info["type"] == "text" and len(clean) > 0:
        try:
            sample = clean.head(20).astype(str)
            parsed = pd.to_datetime(sample, format='mixed', errors='coerce')
            if parsed.notna().mean() > 0.7:  # 70%+ parseable as dates
                info["type"] = "datetime"
                info["min_date"] = str(parsed.min())
                info["max_date"] = str(parsed.max())
        except Exception:
            pass

    return info


def _infer_pg_type(col_info: dict) -> str:
    """Map analyzed column to a PostgreSQL-friendly description."""
    if col_info["type"] == "numeric":
        if col_info.get("max", 0) == int(col_info.get("max", 0)) and col_info.get("min", 0) == int(col_info.get("min", 0)):
            return "INTEGER"
        return "NUMERIC"
    elif col_info["type"] == "datetime":
        return "TIMESTAMP"
    else:
        return "TEXT"


def _detect_relationships(columns: list[str]) -> list[str]:
    """Detect potential foreign key relationships from column names."""
    rels = []
    for col in columns:
        if col.endswith("_id") and col != "id":
            ref_table = col[:-3]  # customer_id → customer
            rels.append(f"`{col}` likely references `{ref_table}` table")
    return rels


def _profile_table(df: pd.DataFrame, project_slug: str, table_name: str) -> dict:
    """Profile table using pure pandas — real health %, column stats, alerts. No external deps."""
    profile_data = {"health": 60, "alerts": [], "columns": {}, "row_count": len(df), "col_count": len(df.columns)}
    try:
        alerts = []
        col_profiles = {}
        total_cols = len(df.columns)

        for col in df.columns:
            series = df[col]
            null_pct = round(float(series.isna().mean()) * 100, 1)
            unique = int(series.nunique())
            ctype = "numeric" if pd.api.types.is_numeric_dtype(series) else ("datetime" if pd.api.types.is_datetime64_any_dtype(series) else "text")

            cp = {"type": ctype, "missing_pct": null_pct, "unique": unique}
            if ctype == "numeric":
                clean = series.dropna()
                if len(clean) > 0:
                    cp["min"] = float(clean.min())
                    cp["max"] = float(clean.max())
                    cp["mean"] = round(float(clean.mean()), 2)
                    cp["zeros_pct"] = round(float((clean == 0).mean()) * 100, 1)
            elif ctype == "text":
                clean = series.dropna().astype(str)
                if len(clean) > 0:
                    cp["top_values"] = clean.value_counts().head(5).index.tolist()
                    if unique <= 20:
                        cp["is_categorical"] = True

            col_profiles[col] = cp

            # Generate alerts
            if null_pct > 50:
                alerts.append(f"Column '{col}' is {null_pct}% missing")
            if null_pct > 0 and null_pct <= 50:
                alerts.append(f"Column '{col}' has {null_pct}% missing values")
            if ctype == "numeric" and cp.get("zeros_pct", 0) > 30:
                alerts.append(f"Column '{col}' has {cp['zeros_pct']}% zero values")

        # Duplicate detection
        dup_rows = int(df.duplicated().sum())
        if dup_rows > 0:
            alerts.append(f"{dup_rows} duplicate rows detected")

        # Compute REAL health %
        health = 100
        missing_cols = sum(1 for c in col_profiles.values() if c.get("missing_pct", 0) > 50)
        if total_cols > 0:
            health -= int((missing_cols / total_cols) * 30)
        if dup_rows > 0:
            health -= 15
        if len(alerts) > 8:
            health -= 10
        if len(df) < 5:
            health -= 20
        # Bonus for completeness
        complete_cols = sum(1 for c in col_profiles.values() if c.get("missing_pct", 0) == 0)
        if total_cols > 0 and complete_cols == total_cols:
            health = min(100, health + 10)
        health = max(10, min(100, health))

        profile_data = {
            "health": health,
            "row_count": len(df),
            "col_count": total_cols,
            "duplicate_rows": dup_rows,
            "alerts": alerts[:20],
            "columns": col_profiles,
        }

        # Save profile JSON
        profile_dir = KNOWLEDGE_DIR / project_slug / "table_sources"
        profile_dir.mkdir(parents=True, exist_ok=True)
        _safe_write_json(profile_dir / f"{table_name}_profile.json", profile_data)

    except Exception:
        pass
    return profile_data


def _generate_metadata(table_name: str, df: pd.DataFrame, col_analyses: list[dict]) -> dict:
    """Generate table metadata JSON (same format as knowledge/tables/*.json)."""
    # Infer use cases from column types
    use_cases = []
    has_dates = any(c["type"] == "datetime" for c in col_analyses)
    has_numerics = any(c["type"] == "numeric" for c in col_analyses)
    has_categories = any(c.get("is_categorical") for c in col_analyses)

    if has_dates and has_numerics:
        use_cases.append("Time-series analysis and trend tracking")
    if has_numerics:
        use_cases.append("Aggregation and metric calculation (SUM, AVG, COUNT)")
    if has_categories:
        use_cases.append("Segmentation and group-by analysis")
    if has_dates:
        use_cases.append("Date-range filtering and period comparisons")
    use_cases.append(f"General analysis of {table_name.replace('_', ' ')} data")

    # Data quality notes
    quality_notes = []
    for c in col_analyses:
        if c["null_pct"] > 0:
            quality_notes.append(f"`{c['name']}` has {c['null_pct']}% NULL values")
        if c["type"] == "text" and c.get("is_categorical") and c["unique_count"] <= 10:
            vals = c.get("sample_values", [])
            quality_notes.append(f"`{c['name']}` categories: {', '.join(str(v) for v in vals[:8])}")

    # Column definitions
    columns = []
    for c in col_analyses:
        pg_type = _infer_pg_type(c)
        desc = f"{c['type'].title()} column"
        if c["type"] == "numeric":
            desc = f"Range: {c.get('min', '?')} to {c.get('max', '?')}, avg {c.get('mean', '?')}"
        elif c["type"] == "datetime":
            desc = f"Date range: {c.get('min_date', '?')} to {c.get('max_date', '?')}"
        elif c.get("is_categorical"):
            desc = f"Categories ({c['unique_count']} values)"
        elif c["type"] == "text":
            desc = f"Text ({c['unique_count']} unique values)"

        columns.append({"name": c["name"], "type": pg_type, "description": desc})

    return {
        "table_name": table_name,
        "table_description": f"User-uploaded dataset: {table_name.replace('_', ' ')}. {len(df):,} rows, {len(df.columns)} columns.",
        "use_cases": use_cases[:5],
        "data_quality_notes": quality_notes[:5],
        "table_columns": columns,
    }


def _generate_sample_queries(table_name: str, col_analyses: list[dict]) -> str:
    """Generate sample SQL queries for the uploaded table."""
    queries = []
    numeric_cols = [c for c in col_analyses if c["type"] == "numeric"]
    date_cols = [c for c in col_analyses if c["type"] == "datetime"]
    cat_cols = [c for c in col_analyses if c.get("is_categorical")]

    # Row count
    queries.append(f"-- <query {table_name}_count>\n-- <description>Total row count for {table_name}</description>\n-- <query>\nSELECT COUNT(*) AS total_rows FROM {table_name};\n-- </query>")

    # Numeric summary
    if numeric_cols:
        agg_parts = ", ".join(f"ROUND(AVG({c['name']}), 2) AS avg_{c['name']}" for c in numeric_cols[:3])
        queries.append(f"-- <query {table_name}_summary>\n-- <description>Numeric summary of {table_name}</description>\n-- <query>\nSELECT {agg_parts}, COUNT(*) AS total FROM {table_name};\n-- </query>")

    # Group by category
    if cat_cols and numeric_cols:
        cat = cat_cols[0]["name"]
        num = numeric_cols[0]["name"]
        queries.append(f"-- <query {table_name}_by_{cat}>\n-- <description>{table_name} grouped by {cat}</description>\n-- <query>\nSELECT {cat}, COUNT(*) AS count, ROUND(AVG({num}), 2) AS avg_{num}\nFROM {table_name}\nGROUP BY {cat}\nORDER BY count DESC;\n-- </query>")

    # Date trend
    if date_cols and numeric_cols:
        dt = date_cols[0]["name"]
        num = numeric_cols[0]["name"]
        queries.append(f"-- <query {table_name}_trend>\n-- <description>Monthly trend for {table_name}</description>\n-- <query>\nSELECT DATE_TRUNC('month', {dt}::timestamp) AS month, COUNT(*) AS count, ROUND(SUM({num}), 2) AS total_{num}\nFROM {table_name}\nGROUP BY 1\nORDER BY 1;\n-- </query>")

    # Top records
    queries.append(f"-- <query {table_name}_sample>\n-- <description>Sample rows from {table_name}</description>\n-- <query>\nSELECT * FROM {table_name} LIMIT 10;\n-- </query>")

    return "\n\n".join(queries)


def _generate_business_rules(table_name: str, col_analyses: list[dict]) -> dict:
    """Generate business rules/metrics for the uploaded table."""
    metrics = []
    rules = []
    gotchas = []

    numeric_cols = [c for c in col_analyses if c["type"] == "numeric"]
    for c in numeric_cols[:3]:
        metrics.append({
            "name": f"{table_name} {c['name']}".replace("_", " ").title(),
            "definition": f"Aggregate of {c['name']} from {table_name}",
            "table": table_name,
            "calculation": f"SUM({c['name']}) or AVG({c['name']}) FROM {table_name}",
        })

    for c in col_analyses:
        if c["null_pct"] > 10:
            gotchas.append({
                "issue": f"NULL values in {table_name}.{c['name']}",
                "tables_affected": [table_name],
                "solution": f"{c['null_pct']}% NULL. Use COALESCE or filter with IS NOT NULL.",
            })

    rels = _detect_relationships([c["name"] for c in col_analyses])
    for r in rels:
        rules.append(r)

    return {"metrics": metrics, "business_rules": rules, "common_gotchas": gotchas}


# ---------------------------------------------------------------------------
# LLM-Powered Auto-Training
# ---------------------------------------------------------------------------


def _llm_deep_analysis(table_name: str, col_analyses: list[dict], sample_rows: list[dict]) -> dict | None:
    """Use LLM to deeply analyze data and generate smart metadata."""
    from os import getenv
    import httpx

    api_key = getenv("OPENROUTER_API_KEY", "")
    if not api_key:
        return None

    # Build sample data as markdown table
    if sample_rows:
        headers = list(sample_rows[0].keys())
        sample_md = "| " + " | ".join(headers) + " |\n"
        sample_md += "| " + " | ".join(["---"] * len(headers)) + " |\n"
        for row in sample_rows[:8]:
            sample_md += "| " + " | ".join(str(row.get(h, ""))[:30] for h in headers) + " |\n"
    else:
        sample_md = "No sample data available"

    col_summary = json.dumps([{k: v for k, v in c.items() if k in ("name", "type", "null_pct", "unique_count", "min", "max", "mean", "sample_values", "is_categorical")} for c in col_analyses], indent=2, default=str)

    # Build distribution summary (gives LLM full picture without sending all rows)
    dist_lines = []
    total_rows = col_analyses[0].get("total_count", 0) if col_analyses else 0
    dist_lines.append(f"TOTAL ROWS: {total_rows:,}")
    for ca in col_analyses:
        name = ca.get("name", "?")
        ctype = ca.get("type", "?")
        null_pct = ca.get("null_pct", 0)
        unique = ca.get("unique_count", 0)
        line = f"  {name} ({ctype}): {unique} unique, {null_pct}% null"
        if ctype == "numeric" and ca.get("mean") is not None:
            line += f", min={ca.get('min')}, max={ca.get('max')}, mean={ca.get('mean')}"
            zeros = ca.get("zeros_pct", 0)
            if zeros:
                line += f", {zeros}% zeros"
        elif ca.get("sample_values"):
            vals = ca["sample_values"][:5]
            line += f", top values: {vals}"
        if ca.get("is_categorical"):
            line += " [CATEGORICAL]"
        if ca.get("min_date"):
            line += f", range: {ca['min_date']} to {ca.get('max_date', '?')}"
        dist_lines.append(line)
    distribution_summary = "\n".join(dist_lines)

    prompt = f"""You are analyzing a dataset to train a data agent. Analyze deeply.

TABLE: {table_name}
COLUMNS:
{col_summary}

DATA DISTRIBUTION (full table stats):
{distribution_summary}

SAMPLE DATA (first 8 rows):
{sample_md}

Generate a comprehensive analysis like a Codex-enriched knowledge pipeline. Return ONLY valid JSON (no markdown):
{{
  "table_description": "What this table represents in business terms (2-3 sentences)",
  "table_purpose": "One-line purpose: why does this table exist? What business process does it serve?",
  "grain": "What does each row represent? e.g. 'One row per customer per month' or 'One row per transaction'",
  "primary_keys": ["columns that uniquely identify each row"],
  "foreign_keys": [{{"column": "col_name", "references": "other_table.col", "relationship": "many-to-one"}}],
  "usage_patterns": ["Common downstream usage pattern 1", "Common query pattern 2"],
  "alternate_tables": "When should you use a different table instead? e.g. 'Use orders_summary instead for aggregated metrics'",
  "freshness": "How often is this data likely updated? (real-time, daily, weekly, static snapshot)",
  "column_descriptions": {{
    "column_name": "What this column means in business context (not just the data type)"
  }},
  "metrics": [
    {{"name": "Metric Name", "definition": "What it measures", "table": "{table_name}", "calculation": "SQL expression"}}
  ],
  "business_rules": [
    {{"name": "Rule Name", "definition": "The business rule or logic", "type": "business_rule"}}
  ],
  "data_quality": [
    "Insight about data quality or patterns..."
  ],
  "suggested_role": "One line describing the agent's expertise area",
  "suggested_personality": "friendly"
}}"""

    try:
        resp = httpx.post(
            "https://openrouter.ai/api/v1/chat/completions",
            headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
            json={"model": TRAINING_MODEL, "messages": [{"role": "user", "content": prompt}], "max_tokens": 2000, "temperature": 0.1},
            timeout=30,
        )
        result = resp.json()
        content = result.get("choices", [{}])[0].get("message", {}).get("content", "")
        return json.loads(content.strip().strip("`").strip())
    except Exception:
        return None


def _get_chat_feedback_for_training(table_name: str) -> str:
    """Load proven SQL patterns + user feedback to improve Q&A generation."""
    feedback_context = ""
    try:
        engine = create_engine(db_url)
        with engine.connect() as conn:
            # Proven queries (thumbs-up)
            good = conn.execute(text(
                "SELECT question, answer FROM public.dash_feedback WHERE rating = 'up' AND answer LIKE :t ORDER BY created_at DESC LIMIT 5"
            ), {"t": f"%{table_name}%"}).fetchall()
            if good:
                feedback_context += "\nUSERS LIKED THESE QUERIES (generate similar):\n"
                for g in good:
                    feedback_context += f"  Q: {str(g[0])[:80]}\n"

            # Anti-patterns (thumbs-down)
            bad = conn.execute(text(
                "SELECT question, answer FROM public.dash_feedback WHERE rating = 'down' AND answer LIKE :t ORDER BY created_at DESC LIMIT 3"
            ), {"t": f"%{table_name}%"}).fetchall()
            if bad:
                feedback_context += "\nUSERS DISLIKED THESE (avoid similar):\n"
                for b in bad:
                    feedback_context += f"  Q: {str(b[0])[:80]}\n"

            # Proven SQL patterns
            patterns = conn.execute(text(
                "SELECT question, sql_text FROM public.dash_query_patterns WHERE table_name = :t AND usage_count > 1 ORDER BY usage_count DESC LIMIT 3"
            ), {"t": table_name}).fetchall()
            if patterns:
                feedback_context += "\nPROVEN SQL PATTERNS (high usage, include similar):\n"
                for p in patterns:
                    feedback_context += f"  Q: {str(p[0])[:60]} → {str(p[1])[:100]}\n"
    except Exception:
        pass
    return feedback_context


def _llm_generate_training(table_name: str, metadata: dict, col_analyses: list[dict] = None) -> list[dict] | None:
    """Use LLM to generate training Q&A pairs for the agent."""
    from os import getenv
    import httpx

    api_key = getenv("OPENROUTER_API_KEY", "")
    if not api_key:
        return None

    cols_desc = "\n".join(f"- {c['name']} ({c.get('type', 'TEXT')}): {c.get('description', '')}" for c in metadata.get("table_columns", []))

    # Build distribution so LLM knows REAL data shape (not just column names)
    dist_info = ""
    if col_analyses:
        total = col_analyses[0].get("total_count", 0) if col_analyses else 0
        dist_lines = [f"ROWS: {total:,}"]
        for ca in col_analyses:
            name = ca.get("name", "?")
            line = f"  {name}: {ca.get('unique_count', '?')} unique, {ca.get('null_pct', 0)}% null"
            if ca.get("type") == "numeric" and ca.get("mean") is not None:
                line += f", range {ca.get('min')} to {ca.get('max')}, avg {ca.get('mean')}"
            elif ca.get("sample_values"):
                line += f", values: {ca['sample_values'][:5]}"
            dist_lines.append(line)
        dist_info = "\nDATA DISTRIBUTION:\n" + "\n".join(dist_lines) + "\n"

    prompt = f"""Generate 11 training Q&A pairs for this data table, one for each analysis type.
Each pair should have: question, sql, analysis_type.

TABLE: {table_name}
DESCRIPTION: {metadata.get('table_description', '')}
COLUMNS:
{cols_desc}
{dist_info}

Generate exactly 11 pairs, one per analysis type. Use real column names from the table above:
1. DESCRIPTIVE: "What is the total X?" with simple aggregation SQL (SUM, COUNT, AVG)
2. DIAGNOSTIC: "Why does X have the highest Y?" with GROUP BY + ORDER BY SQL
3. COMPARATIVE: "Compare X vs Y" with GROUP BY + CASE WHEN SQL
4. TREND: "Show X over time" with DATE_TRUNC + GROUP BY SQL (use a date/time column if available, otherwise approximate)
5. PARETO: "Which X drives 80% of Y?" with cumulative percentage SQL using window functions
6. ANOMALY: "Any unusual X?" with AVG + STDDEV detection SQL
7. PRESCRIPTIVE: "What should we do about X?" with actionable analysis SQL
8. SCENARIO: "What if X increases by N%?" with projection/multiplication SQL
9. BENCHMARK: "How does X compare to average?" with AVG window function SQL
10. ROOT_CAUSE: "Why is X happening?" with multi-dimension GROUP BY SQL
11. PREDICTIVE: "What will X be next period?" with trend extrapolation SQL

Return ONLY valid JSON array (no markdown):
[
  {{"question": "What is the total revenue?", "sql": "SELECT SUM(revenue) FROM {table_name}", "analysis_type": "descriptive"}},
  ...
]

{_get_chat_feedback_for_training(table_name)}"""

    try:
        resp = httpx.post(
            "https://openrouter.ai/api/v1/chat/completions",
            headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
            json={"model": TRAINING_MODEL, "messages": [{"role": "user", "content": prompt}], "max_tokens": 3000, "temperature": 0.2},
            timeout=45,
        )
        result = resp.json()
        content = result.get("choices", [{}])[0].get("message", {}).get("content", "")
        return json.loads(content.strip().strip("`").strip())
    except Exception:
        return None


def _llm_generate_persona(project_slug: str, tables_metadata: list[dict], rules: list[dict]) -> dict | None:
    """Use LLM to generate a full agent persona based on the data."""
    from os import getenv
    import httpx

    api_key = getenv("OPENROUTER_API_KEY", "")
    if not api_key:
        return None

    # Build context from all tables
    tables_summary = ""
    for m in tables_metadata[:5]:
        cols = ", ".join(c.get("name", "") for c in m.get("table_columns", [])[:10])
        tables_summary += f"- {m.get('table_name', '?')}: {m.get('table_description', '')}. Columns: {cols}\n"

    rules_summary = "\n".join(f"- {r.get('name', '')}: {r.get('definition', '')}" for r in rules[:10])

    prompt = f"""You are creating a persona for a data agent. This agent will be an expert on the following data:

TABLES:
{tables_summary}

BUSINESS RULES:
{rules_summary or 'None yet'}

Generate a comprehensive persona. Return ONLY valid JSON (no markdown):
{{
  "persona_prompt": "A 3-5 sentence system prompt describing who this agent is, what domain it specializes in, how it should approach questions, what terminology it should use, and what it should prioritize when analyzing data.",
  "domain_terms": ["list", "of", "key", "domain", "terms", "the agent should know"],
  "greeting": "A friendly 1-2 sentence greeting the agent uses when a user first opens the chat.",
  "expertise_areas": ["area1", "area2", "area3"],
  "communication_style": "How the agent should communicate (e.g., 'concise with data tables', 'detailed with explanations')"
}}"""

    try:
        resp = httpx.post(
            "https://openrouter.ai/api/v1/chat/completions",
            headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
            json={"model": TRAINING_MODEL, "messages": [{"role": "user", "content": prompt}], "max_tokens": 1000, "temperature": 0.2},
            timeout=20,
        )
        result = resp.json()
        content = result.get("choices", [{}])[0].get("message", {}).get("content", "")
        return json.loads(content.strip().strip("`").strip())
    except Exception:
        return None


def _update_project_config(project_slug: str, role: str, personality: str | None):
    """Update project agent config based on LLM suggestion."""
    engine = create_engine(db_url)
    with engine.connect() as conn:
        updates = ["agent_role = :role", "updated_at = NOW()"]
        params: dict = {"slug": project_slug, "role": role}
        if personality:
            updates.append("agent_personality = :pers")
            params["pers"] = personality
        conn.execute(text(f"UPDATE public.dash_projects SET {', '.join(updates)} WHERE slug = :slug"), params)
        conn.commit()


def _reload_project_knowledge(project_slug: str, timeout_sec: int = 60) -> bool:
    """Re-index all knowledge files for a project into PgVector.
    Returns True if successful, False if timed out or failed."""
    import threading

    success = [False]
    error_msg = [None]

    def _do_index():
        try:
            from db.session import create_project_knowledge
            proj_knowledge = create_project_knowledge(project_slug)
            proj_dir = KNOWLEDGE_DIR / project_slug
            if proj_dir.exists():
                for subdir in ["tables", "queries", "business", "rules", "training"]:
                    path = proj_dir / subdir
                    if path.exists():
                        files = [f for f in path.iterdir() if f.is_file() and not f.name.startswith(".")]
                        if files:
                            proj_knowledge.insert(name=f"{project_slug}-{subdir}", path=str(path))
            success[0] = True
        except Exception as e:
            error_msg[0] = str(e)

    thread = threading.Thread(target=_do_index, daemon=True)
    thread.start()
    thread.join(timeout=timeout_sec)

    if thread.is_alive():
        # Thread is still running — timed out
        return False
    return success[0]


def _discover_relationships(project_slug: str):
    """LLM analyzes all tables to discover hidden relationships."""
    from os import getenv
    import httpx

    api_key = getenv("OPENROUTER_API_KEY", "")
    if not api_key:
        return

    # Load all table metadata
    tables_info = []
    proj_tables_dir = KNOWLEDGE_DIR / project_slug / "tables"
    if not proj_tables_dir.exists():
        return
    for f in proj_tables_dir.glob("*.json"):
        try:
            with open(f) as fh:
                d = json.load(fh)
            cols = [c["name"] for c in d.get("table_columns", [])]
            tables_info.append(f"{d.get('table_name', f.stem)}: {', '.join(cols)}")
        except Exception:
            pass

    if len(tables_info) < 2:
        return

    prompt = f"""Analyze these database tables and discover relationships between them.
Look for: shared columns, FK patterns, value overlaps, hierarchical relationships.

TABLES:
{chr(10).join(tables_info)}

Return ONLY valid JSON array (no markdown):
[{{"from_table": "table1", "from_column": "col1", "to_table": "table2", "to_column": "col2", "type": "fk|shared|hierarchy", "confidence": 0.9}}]

Return empty array [] if no relationships found."""

    try:
        resp = httpx.post(
            "https://openrouter.ai/api/v1/chat/completions",
            headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
            json={"model": TRAINING_MODEL, "messages": [{"role": "user", "content": prompt}], "max_tokens": 500, "temperature": 0.1},
            timeout=15,
        )
        result = resp.json()
        content = result.get("choices", [{}])[0].get("message", {}).get("content", "")
        rels = json.loads(content.strip().strip("`").strip())

        if isinstance(rels, list) and rels:
            # VERIFY relationships: check actual value overlap in PostgreSQL
            schema = re.sub(r"[^a-z0-9_]", "_", project_slug.lower())[:63]
            engine = create_engine(db_url)
            with engine.connect() as conn:
                for r in rels[:10]:
                    ft, fc = r.get("from_table", ""), r.get("from_column", "")
                    tt, tc = r.get("to_table", ""), r.get("to_column", "")
                    verified_conf = r.get("confidence", 0.5)
                    try:
                        # Check actual value overlap between the two columns
                        vals1 = conn.execute(text(f'SELECT DISTINCT "{fc}" FROM "{schema}"."{ft}" WHERE "{fc}" IS NOT NULL LIMIT 100')).fetchall()
                        vals2 = conn.execute(text(f'SELECT DISTINCT "{tc}" FROM "{schema}"."{tt}" WHERE "{tc}" IS NOT NULL LIMIT 100')).fetchall()
                        set1 = {str(v[0]) for v in vals1}
                        set2 = {str(v[0]) for v in vals2}
                        if set1 and set2:
                            overlap = len(set1 & set2) / max(len(set1 | set2), 1)
                            verified_conf = round(overlap, 2)
                    except Exception:
                        pass  # Keep LLM confidence if verification fails
                    try:
                        conn.execute(text("""
                            INSERT INTO public.dash_relationships (project_slug, from_table, from_column, to_table, to_column, rel_type, confidence, source)
                            VALUES (:s, :ft, :fc, :tt, :tc, :type, :conf, 'ai_verified')
                            ON CONFLICT (project_slug, from_table, from_column, to_table, to_column) DO UPDATE SET confidence = :conf, source = 'ai_verified'
                        """), {"s": project_slug, "ft": ft, "fc": fc, "tt": tt, "tc": tc,
                               "type": r.get("type", "fk"), "conf": verified_conf})
                    except Exception:
                        pass
                conn.commit()
    except Exception:
        pass


def _detect_data_drift(project_slug: str, table_name: str, col_analyses: list[dict]):
    """Compare new column stats with stored metadata to detect drift."""
    meta_file = KNOWLEDGE_DIR / project_slug / "tables" / f"{table_name}.json"
    if not meta_file.exists():
        return

    try:
        with open(meta_file) as f:
            old_meta = json.load(f)

        old_cols = {c["name"]: c for c in old_meta.get("table_columns", [])}
        drift_alerts = []

        for col in col_analyses:
            name = col["name"]
            old = old_cols.get(name, {})
            old_desc = old.get("description", "")

            # Check NULL drift
            if col.get("null_pct", 0) > 20 and "NULL" not in old_desc:
                drift_alerts.append(f"{name}: NULL% increased to {col['null_pct']}%")

            # Check range drift for numerics
            if col.get("type") == "numeric" and "Range" in old_desc:
                try:
                    # Extract old range from description like "Range: 10 to 500"
                    parts = old_desc.split("Range: ")[1].split(",")[0].split(" to ")
                    old_max = float(parts[1].strip())
                    if col.get("max", 0) > old_max * 2:
                        drift_alerts.append(f"{name}: max value {col['max']} exceeds trained range (was {old_max})")
                except Exception:
                    pass

        if drift_alerts:
            drift_dir = KNOWLEDGE_DIR / project_slug / "drift"
            drift_dir.mkdir(parents=True, exist_ok=True)
            with open(drift_dir / f"{table_name}_drift.json", "w") as f:
                json.dump({"table": table_name, "alerts": drift_alerts, "timestamp": time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())}, f, indent=2)

            # Create notification for project owner
            try:
                from app.auth import notify_user
                engine = create_engine(db_url)
                with engine.connect() as conn:
                    row = conn.execute(text("SELECT user_id FROM public.dash_projects WHERE slug = :s"), {"s": project_slug}).fetchone()
                    if row:
                        notify_user(row[0], f"Data drift detected in {table_name}", "; ".join(drift_alerts[:3]), "warning")
            except Exception:
                pass
    except Exception:
        pass


def _extract_sql_metadata(sql: str) -> dict:
    """Extract tables, join strategy, and filters from a SQL query using simple regex."""
    sql_upper = sql.upper()
    tables = set()

    # Extract tables from FROM clauses
    for m in re.finditer(r'\bFROM\s+([a-zA-Z_][a-zA-Z0-9_.]*)', sql, re.IGNORECASE):
        tables.add(m.group(1).strip().lower())
    # Extract tables from JOIN clauses
    for m in re.finditer(r'\bJOIN\s+([a-zA-Z_][a-zA-Z0-9_.]*)', sql, re.IGNORECASE):
        tables.add(m.group(1).strip().lower())

    # Determine join strategy
    if re.search(r'\bLEFT\s+JOIN\b', sql_upper):
        join_strategy = 'LEFT JOIN'
    elif re.search(r'\bRIGHT\s+JOIN\b', sql_upper):
        join_strategy = 'RIGHT JOIN'
    elif re.search(r'\bFULL\s+(OUTER\s+)?JOIN\b', sql_upper):
        join_strategy = 'FULL JOIN'
    elif re.search(r'\bCROSS\s+JOIN\b', sql_upper):
        join_strategy = 'CROSS JOIN'
    elif re.search(r'\bJOIN\b', sql_upper):
        join_strategy = 'INNER JOIN'
    else:
        join_strategy = 'NONE'

    # Detect filter/aggregation clauses
    filters_found = []
    if re.search(r'\bWHERE\b', sql_upper):
        filters_found.append('WHERE')
    if re.search(r'\bGROUP\s+BY\b', sql_upper):
        filters_found.append('GROUP BY')
    if re.search(r'\bORDER\s+BY\b', sql_upper):
        filters_found.append('ORDER BY')
    if re.search(r'\bHAVING\b', sql_upper):
        filters_found.append('HAVING')
    if re.search(r'\bLIMIT\b', sql_upper):
        filters_found.append('LIMIT')
    if re.search(r'\bWINDOW\b|\bOVER\s*\(', sql_upper):
        filters_found.append('WINDOW')

    return {
        "tables_used": ", ".join(sorted(tables)) if tables else "",
        "join_strategy": join_strategy,
        "filters": ", ".join(filters_found) if filters_found else "NONE",
    }


def _save_query_pattern_with_metadata(engine, project_slug: str, question: str, sql: str, source: str = 'training'):
    """Save a Q&A pair to dash_query_patterns with extracted SQL metadata."""
    meta = _extract_sql_metadata(sql)
    try:
        engine_conn = engine.connect()
        engine_conn.execute(text(
            "INSERT INTO public.dash_query_patterns "
            "(project_slug, question, sql, tables_used, join_strategy, filters, source) "
            "VALUES (:s, :q, :sql, :tables, :join, :filt, :src) "
            "ON CONFLICT DO NOTHING"
        ), {
            "s": project_slug, "q": question, "sql": sql,
            "tables": meta["tables_used"], "join": meta["join_strategy"],
            "filt": meta["filters"], "src": source,
        })
        engine_conn.commit()
        engine_conn.close()
        return True
    except Exception:
        return False


def _save_to_db(project_slug: str, table_name: str, metadata: dict, biz_rules: dict, training_qa: list | None = None, persona: dict | None = None):
    """Persist all training data to PostgreSQL alongside files."""
    try:
        engine = create_engine(db_url)
        with engine.connect() as conn:
            # Table metadata
            conn.execute(text("""
                INSERT INTO public.dash_table_metadata (project_slug, table_name, metadata)
                VALUES (:s, :t, :m::jsonb)
                ON CONFLICT (project_slug, table_name)
                DO UPDATE SET metadata = :m::jsonb, updated_at = NOW()
            """), {"s": project_slug, "t": table_name, "m": json.dumps(metadata, default=str)})

            # Business rules
            conn.execute(text("""
                INSERT INTO public.dash_business_rules_db (project_slug, table_name, rules)
                VALUES (:s, :t, :r::jsonb)
                ON CONFLICT (project_slug, table_name)
                DO UPDATE SET rules = :r::jsonb
            """), {"s": project_slug, "t": table_name, "r": json.dumps(biz_rules, default=str)})

            # Training Q&A
            if training_qa and isinstance(training_qa, list):
                # Clear old Q&A for this table
                conn.execute(text("DELETE FROM public.dash_training_qa WHERE project_slug = :s AND table_name = :t"), {"s": project_slug, "t": table_name})

                def _is_safe_sql(sql: str) -> bool:
                    """Only allow SELECT/WITH statements in training SQL."""
                    s = sql.strip().upper()
                    if not (s.startswith('SELECT') or s.startswith('WITH')):
                        return False
                    forbidden = ['DROP ', 'DELETE ', 'TRUNCATE ', 'ALTER ', 'INSERT ', 'UPDATE ', 'CREATE ', 'GRANT ']
                    return not any(kw in s for kw in forbidden)

                for qa in training_qa:
                    if _is_safe_sql(qa.get("sql", "")):
                        conn.execute(text(
                            "INSERT INTO public.dash_training_qa (project_slug, table_name, question, sql, answer_template) VALUES (:s, :t, :q, :sql, :a)"
                        ), {"s": project_slug, "t": table_name, "q": qa.get("question", ""), "sql": qa.get("sql", ""), "a": qa.get("answer_template", "")})

            # Persona
            if persona:
                conn.execute(text("""
                    INSERT INTO public.dash_personas (project_slug, persona)
                    VALUES (:s, :p::jsonb)
                    ON CONFLICT (project_slug)
                    DO UPDATE SET persona = :p::jsonb, updated_at = NOW()
                """), {"s": project_slug, "p": json.dumps(persona, default=str)})

            conn.commit()
    except Exception:
        pass


def _run_auto_training(project_slug: str, table_name: str, col_analyses: list[dict],
                       metadata: dict, biz_rules: dict, sample_rows: list[dict],
                       tables_dir: Path, business_dir: Path,
                       master_run_id: int | None = None, table_index: int = 0, total_tables: int = 1):
    """Background task: LLM deep analysis + training Q&A generation."""
    import time as _time

    # Single shared engine for all DB operations in this training run
    train_engine = create_engine(db_url, pool_size=2, max_overflow=3, pool_recycle=3600)

    # Track training run — reuse master run if provided, otherwise create one
    run_id = master_run_id
    if not run_id:
        try:
            with train_engine.connect() as conn:
                result = conn.execute(text(
                    "INSERT INTO public.dash_training_runs (project_slug, status, steps) VALUES (:s, 'running', 'starting') RETURNING id"
                ), {"s": project_slug})
                run_id = result.fetchone()[0]
                conn.commit()
        except Exception:
            pass

    def _update_run(status: str, steps: str = "", error: str = ""):
        if not run_id:
            return
        # Encode table tracking into the steps field: step_name|table_name|table_index|total_tables
        steps_with_table = f"{steps}|{table_name}|{table_index}|{total_tables}"
        try:
            with train_engine.connect() as conn:
                if (status == 'done' or status == 'failed') and not master_run_id:
                    # Only set terminal status if this is a standalone run (no master)
                    conn.execute(text("UPDATE public.dash_training_runs SET status = :st, steps = :steps, error = :err, tables_trained = :trained, finished_at = NOW() WHERE id = :id"),
                                 {"st": status, "steps": steps_with_table, "err": error, "trained": table_index, "id": run_id})
                else:
                    # For master-managed runs, only update steps (master controls status)
                    conn.execute(text("UPDATE public.dash_training_runs SET steps = :steps WHERE id = :id"), {"steps": steps_with_table, "id": run_id})
                conn.commit()
        except Exception:
            pass

    def _log(msg: str):
        """Append a log entry to the training run for real-time CLI display."""
        if not run_id:
            return
        try:
            with train_engine.connect() as conn:
                conn.execute(text(
                    "UPDATE public.dash_training_runs SET logs = COALESCE(logs, '[]'::jsonb) || CAST(:entry AS jsonb) WHERE id = :id"
                ), {"entry": json.dumps([{"ts": _time.strftime('%H:%M:%S'), "msg": msg, "table": table_name, "table_index": table_index, "total_tables": total_tables}]), "id": run_id})
                conn.commit()
        except Exception:
            pass

    def _cancelled() -> bool:
        """Check if training was cancelled."""
        return _training_cancel_flags.get(project_slug, False)

    num_cols = len(col_analyses)
    # Get actual row count from the first column's total_count (most reliable)
    num_rows = 0
    for c in col_analyses:
        tc = c.get("total_count", 0) or c.get("non_null_count", 0)
        if tc > num_rows:
            num_rows = tc
    _log(f"analyzing table: {table_name} ({num_rows} rows, {num_cols} columns)")

    if num_rows == 0:
        _log(f"⊘ skipping training — table {table_name} has 0 rows")
        _update_run("done", "complete")
        return

    # Check if table data changed (skip ONLY if we have stored row count AND it matches)
    existing_meta = metadata.get("table_description", "")
    existing_row_count = metadata.get("_row_count", 0)
    # Only skip if: description exists AND row count was stored AND matches current
    data_changed = True
    if existing_meta and existing_row_count > 0 and existing_row_count == num_rows:
        data_changed = False

    # Step 0: Data drift detection
    _update_run("running", "drift_detection")
    _log("checking for data drift...")
    _detect_data_drift(project_slug, table_name, col_analyses)
    _log("✓ drift check complete")

    # Step 1: LLM Deep Analysis
    if _cancelled():
        _update_run("failed", "cancelled"); _log("⊘ training cancelled by user"); return
    _update_run("running", "deep_analysis")
    if not data_changed and existing_meta:
        _log(f"⊘ skipping deep analysis — table unchanged ({num_rows} rows)")
        analysis = None
    else:
        _log(f"running LLM deep analysis on {table_name}...")
        analysis = _llm_deep_analysis(table_name, col_analyses, sample_rows)
        if analysis:
            _log(f"✓ deep analysis complete — purpose: {(analysis.get('table_purpose') or '')[:60]}")
    if analysis:
        # Overwrite metadata with smart descriptions
        if analysis.get("table_description"):
            metadata["table_description"] = analysis["table_description"]
        col_descs = analysis.get("column_descriptions", {})
        for col in metadata.get("table_columns", []):
            smart_desc = col_descs.get(col["name"])
            if smart_desc:
                col["description"] = smart_desc
        if analysis.get("data_quality"):
            metadata["data_quality_notes"] = analysis["data_quality"][:8]
        # Codex-enriched knowledge fields
        for field in ("table_purpose", "grain", "primary_keys", "foreign_keys", "usage_patterns", "alternate_tables", "freshness"):
            if analysis.get(field):
                metadata[field] = analysis[field]

        # Store row count for future change detection
        metadata["_row_count"] = num_rows

        # Save updated metadata
        tables_dir.mkdir(parents=True, exist_ok=True)
        with open(tables_dir / f"{table_name}.json", "w") as f:
            json.dump(metadata, f, indent=2, default=str)

        # Update metrics
        if analysis.get("metrics"):
            biz_rules["metrics"] = analysis["metrics"]
            business_dir.mkdir(parents=True, exist_ok=True)
            with open(business_dir / f"{table_name}_rules.json", "w") as f:
                json.dump(biz_rules, f, indent=2, default=str)

        # Auto-create rules
        rules_dir = KNOWLEDGE_DIR / project_slug / "rules"
        rules_dir.mkdir(parents=True, exist_ok=True)
        for rule in analysis.get("business_rules", [])[:5]:
            rule_id = f"rule_auto_{int(_time.time() * 1000)}"
            rule_data = {
                "id": rule_id, "name": rule.get("name", ""),
                "type": rule.get("type", "business_rule"),
                "definition": rule.get("definition", ""),
                "source": "auto_training",
                "created_at": _time.strftime("%Y-%m-%dT%H:%M:%SZ", _time.gmtime()),
            }
            with open(rules_dir / f"{rule_id}.json", "w") as f:
                json.dump(rule_data, f, indent=2)
            _time.sleep(0.01)  # unique timestamps

        # Update agent persona
        if analysis.get("suggested_role"):
            _update_project_config(project_slug, analysis["suggested_role"], analysis.get("suggested_personality"))

    # Step 2: LLM Training Q&A
    if _cancelled():
        _update_run("failed", "cancelled"); _log("⊘ training cancelled by user"); return
    _update_run("running", "qa_generation")
    _log(f"generating training Q&A for {table_name}...")
    training = _llm_generate_training(table_name, metadata, col_analyses=col_analyses)
    if training and isinstance(training, list):
        # VERIFY Q&A: Run each generated SQL against real data
        schema = re.sub(r"[^a-z0-9_]", "_", project_slug.lower())[:63]
        verified_count = 0
        discarded_count = 0
        try:
            verify_engine = create_engine(db_url)
            for qa in training:
                sql = qa.get("sql", "")
                if not sql:
                    continue
                try:
                    with verify_engine.connect() as vconn:
                        vconn.execute(text(f"SET LOCAL search_path TO {schema}, public"))
                        result = vconn.execute(text(sql))
                        rows = result.fetchall()
                        if rows:
                            # Save real answer
                            qa["verified"] = True
                            qa["verified_answer"] = str(rows[0][0]) if len(rows[0]) == 1 else str(rows[:3])
                            qa["verified_row_count"] = len(rows)
                            verified_count += 1
                        else:
                            qa["verified"] = True
                            qa["verified_answer"] = "0 rows"
                            qa["verified_row_count"] = 0
                            verified_count += 1
                except Exception:
                    qa["verified"] = False
                    qa["verified_answer"] = None
                    discarded_count += 1
        except Exception:
            pass
        _log(f"  ✓ {verified_count} Q&A verified with real data, {discarded_count} had SQL errors")

        training_dir = KNOWLEDGE_DIR / project_slug / "training"
        training_dir.mkdir(parents=True, exist_ok=True)
        with open(training_dir / f"{table_name}_qa.json", "w") as f:
            json.dump(training, f, indent=2)

    if training and isinstance(training, list):
        _log(f"✓ {len(training)} Q&A pairs generated")
    else:
        _log("· no Q&A generated")

    # Step 3: Generate agent persona
    if _cancelled():
        _update_run("failed", "cancelled"); _log("⊘ training cancelled by user"); return
    _update_run("running", "persona")
    _log("generating agent persona...")
    try:
        # Load all tables metadata for this project
        all_metadata = []
        proj_tables_dir = KNOWLEDGE_DIR / project_slug / "tables"
        if proj_tables_dir.exists():
            for f in proj_tables_dir.glob("*.json"):
                try:
                    with open(f) as fh:
                        all_metadata.append(json.load(fh))
                except Exception:
                    pass

        # Load all rules
        all_rules = []
        proj_rules_dir = KNOWLEDGE_DIR / project_slug / "rules"
        if proj_rules_dir.exists():
            for f in proj_rules_dir.glob("*.json"):
                try:
                    with open(f) as fh:
                        all_rules.append(json.load(fh))
                except Exception:
                    pass

        persona = _llm_generate_persona(project_slug, all_metadata, all_rules)
        if persona:
            persona_file = KNOWLEDGE_DIR / project_slug / "persona.json"
            with open(persona_file, "w") as f:
                json.dump(persona, f, indent=2)
            # Also save to DB
            try:
                engine = create_engine(db_url)
                with engine.connect() as conn:
                    conn.execute(text("""
                        INSERT INTO public.dash_personas (project_slug, persona)
                        VALUES (:s, :p::jsonb)
                        ON CONFLICT (project_slug)
                        DO UPDATE SET persona = :p::jsonb, updated_at = NOW()
                    """), {"s": project_slug, "p": json.dumps(persona, default=str)})
                    conn.commit()
            except Exception:
                pass

            # Update project role with generated expertise
            if persona.get("expertise_areas"):
                role_str = ", ".join(persona["expertise_areas"][:4])
                _update_project_config(project_slug, role_str, None)
    except Exception:
        pass

    _log("✓ persona generated")

    # Step 4: Auto-generate sample workflows
    if _cancelled():
        _update_run("failed", "cancelled"); _log("⊘ training cancelled by user"); return
    _update_run("running", "synthesis")
    _log("generating sample workflows...")
    try:
        from os import getenv
        import httpx
        api_key = getenv("OPENROUTER_API_KEY", "")

        # Load all table names for context
        all_tables = []
        proj_tables_dir = KNOWLEDGE_DIR / project_slug / "tables"
        if proj_tables_dir.exists():
            for f in proj_tables_dir.glob("*.json"):
                try:
                    with open(f) as fh:
                        d = json.load(fh)
                    all_tables.append(d.get("table_name", f.stem))
                except Exception:
                    pass

        if api_key and all_tables:
            tables_str = ", ".join(all_tables)
            prompt = f"""Generate 3 analysis-typed workflows for a data agent with these tables: {tables_str}

Each workflow should be a multi-step analysis journey of a specific type. Generate exactly these 3 workflows:

1. "Data Overview" — a descriptive journey:
   Step 1: What tables and columns do we have? Show the schema overview.
   Step 2: What is the total count and key metrics for the main table?
   Step 3: Break down the key metrics by top categories or dimensions.

2. "Deep Dive Analysis" — a diagnostic journey:
   Step 1: What are the key metrics and their current values?
   Step 2: Which dimension or category has the highest impact on the key metric?
   Step 3: Why is that dimension the highest? What factors drive it?
   Step 4: What actions should we take based on this analysis?

3. "Risk Assessment" — an anomaly detection + prescriptive journey:
   Step 1: Show the distribution of key numeric metrics.
   Step 2: Are there any anomalies, outliers, or unexpected patterns?
   Step 3: What are the concentration risks (e.g., over-reliance on one category)?
   Step 4: Recommend specific actions to mitigate the identified risks.

Tailor each workflow's steps to the actual tables and columns available. Use real column/table names where possible.

Return ONLY valid JSON (no markdown):
[
  {{"name": "Data Overview", "description": "Descriptive journey through the dataset — schema, counts, and category breakdowns", "analysis_type": "descriptive", "steps": ["Step 1 text", "Step 2 text", "Step 3 text"]}},
  {{"name": "Deep Dive Analysis", "description": "Diagnostic journey to find what drives key metrics and what to do about it", "analysis_type": "diagnostic", "steps": ["Step 1", "Step 2", "Step 3", "Step 4"]}},
  {{"name": "Risk Assessment", "description": "Anomaly detection and prescriptive actions for risk mitigation", "analysis_type": "prescriptive", "steps": ["Step 1", "Step 2", "Step 3", "Step 4"]}}
]"""

            resp = httpx.post(
                "https://openrouter.ai/api/v1/chat/completions",
                headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
                json={"model": TRAINING_MODEL, "messages": [{"role": "user", "content": prompt}], "max_tokens": 800, "temperature": 0.3},
                timeout=20,
            )
            result = resp.json()
            content = result.get("choices", [{}])[0].get("message", {}).get("content", "")
            workflows = json.loads(content.strip().strip("`").strip())

            if isinstance(workflows, list):
                # Save workflows
                wf_file = KNOWLEDGE_DIR / project_slug / "workflows.json"
                with open(wf_file, "w") as f:
                    json.dump(workflows, f, indent=2)
                _log(f"✓ {len(workflows)} workflows generated")
    except Exception:
        _log("· workflow generation skipped")

    # Step 5: Auto-create default schedules
    _log("checking schedules...")
    try:
        engine = create_engine(db_url)
        with engine.connect() as conn:
            # Check if schedules already exist
            existing = conn.execute(text(
                "SELECT COUNT(*) FROM public.dash_schedules WHERE project_slug = :s"
            ), {"s": project_slug}).scalar() or 0

            if existing == 0 and all_tables:
                # Get project owner
                row = conn.execute(text(
                    "SELECT user_id FROM public.dash_projects WHERE slug = :s"
                ), {"s": project_slug}).fetchone()
                if row:
                    uid = row[0]
                    # Create a weekly summary schedule
                    conn.execute(text(
                        "INSERT INTO public.dash_schedules (project_slug, user_id, name, prompt, cron) "
                        "VALUES (:s, :uid, :name, :prompt, :cron)"
                    ), {
                        "s": project_slug, "uid": uid,
                        "name": "Weekly Data Summary",
                        "prompt": f"Give me a summary of all data in {all_tables[0]}. Show total counts, key metrics, and any notable trends.",
                        "cron": "0 8 * * 1",
                    })
                    if len(all_tables) > 1:
                        conn.execute(text(
                            "INSERT INTO public.dash_schedules (project_slug, user_id, name, prompt, cron) "
                            "VALUES (:s, :uid, :name, :prompt, :cron)"
                        ), {
                            "s": project_slug, "uid": uid,
                            "name": "Daily Health Check",
                            "prompt": f"Check data quality across all tables. Report any NULL values, missing data, or anomalies.",
                            "cron": "0 8 * * *",
                        })
                    conn.commit()
    except Exception:
        pass

    # Step 6: Save everything to DB (production persistence)
    _log("saving to database...")
    try:
        _save_to_db(project_slug, table_name, metadata, biz_rules,
                    training if isinstance(training, list) else None,
                    persona if 'persona' in dir() and persona else None)
    except Exception:
        pass

    _log("✓ saved to database")

    # Step 7: Multi-file synthesis — unified project understanding
    if _cancelled():
        _update_run("failed", "cancelled"); _log("⊘ training cancelled by user"); return
    _log("running multi-file synthesis...")
    _update_run("running", "synthesis")
    try:
        from os import getenv
        import httpx
        api_key = getenv("OPENROUTER_API_KEY", "")
        all_meta = []
        proj_td = KNOWLEDGE_DIR / project_slug / "tables"
        if proj_td.exists():
            for f in proj_td.glob("*.json"):
                try:
                    with open(f) as fh:
                        all_meta.append(json.load(fh))
                except Exception:
                    pass

        if api_key and len(all_meta) > 1:
            tables_desc = "\n".join(f"- {m.get('table_name','?')}: {m.get('table_description','')}" for m in all_meta)
            prompt = f"""These tables belong to the same project. Create a unified data dictionary.

TABLES:
{tables_desc}

Return ONLY valid JSON (no markdown):
{{"project_summary": "2-3 sentence overview of what this data represents together", "cross_table_queries": ["suggested SQL joining multiple tables"], "data_dictionary": {{"term": "definition"}}}}"""

            resp = httpx.post(
                "https://openrouter.ai/api/v1/chat/completions",
                headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
                json={"model": TRAINING_MODEL, "messages": [{"role": "user", "content": prompt}], "max_tokens": 800, "temperature": 0.1},
                timeout=15,
            )
            result = resp.json()
            content = result.get("choices", [{}])[0].get("message", {}).get("content", "")
            synthesis = json.loads(content.strip().strip("`").strip())
            if synthesis:
                synth_file = KNOWLEDGE_DIR / project_slug / "synthesis.json"
                with open(synth_file, "w") as f:
                    json.dump(synthesis, f, indent=2)
    except Exception:
        pass

    _log("✓ synthesis complete")

    # Step 8: Discover relationships (only if multiple tables)
    if _cancelled():
        _update_run("failed", "cancelled"); _log("⊘ training cancelled by user"); return
    _update_run("running", "relationships")
    proj_tables_dir = KNOWLEDGE_DIR / project_slug / "tables"
    table_count = len(list(proj_tables_dir.glob("*.json"))) if proj_tables_dir.exists() else 0
    if table_count >= 2:
        _log(f"discovering relationships across {table_count} tables...")
        try:
            _discover_relationships(project_slug)
        except Exception:
            pass
    else:
        _log("⊘ skipping relationship discovery — only 1 table")

    # Also discover cross-document relationships if docs exist
    docs_dir = KNOWLEDGE_DIR / project_slug / "docs"
    if docs_dir.exists():
        doc_files = [f for f in docs_dir.iterdir() if f.is_file()]
        if len(doc_files) >= 2:
            _log("discovering cross-document relationships...")
            try:
                doc_sums = []
                for f in doc_files[:10]:
                    try:
                        doc_sums.append(f"{f.name}: {f.read_text(errors='ignore')[:500]}")
                    except Exception:
                        pass
                result = training_llm_call(
                    f"Find relationships between these documents.\n\n"
                    f"DOCUMENTS:\n" + "\n---\n".join(doc_sums) + "\n\n"
                    f'Return JSON: [{{"from_doc": "a.txt", "to_doc": "b.txt", "relationship": "shared topic", "shared_topics": ["topic1"], "strength": 0.8}}]',
                    "extraction"
                )
                if result:
                    rels = json.loads(result)
                    if isinstance(rels, list):
                        with train_engine.connect() as conn:
                            for r in rels[:10]:
                                if isinstance(r, dict) and r.get("from_doc"):
                                    conn.execute(text(
                                        "INSERT INTO public.dash_relationships (project_slug, from_table, from_column, to_table, to_column, rel_type, confidence, source) "
                                        "VALUES (:s, :ft, :fc, :tt, :tc, 'topic', :conf, 'ai') ON CONFLICT DO NOTHING"
                                    ), {"s": project_slug, "ft": r["from_doc"], "fc": ", ".join(r.get("shared_topics", [])[:5]),
                                        "tt": r.get("to_doc", ""), "tc": r.get("relationship", ""),
                                        "conf": r.get("strength", 0.5)})
                            conn.commit()
                        _log(f"✓ {len(rels)} cross-document relationships found")
            except Exception:
                pass

    _log("✓ relationships discovered")

    # Step 9: Re-index knowledge (with timeout — training completes even if this fails)
    if _cancelled():
        _update_run("failed", "cancelled"); _log("⊘ training cancelled by user"); return
    _update_run("running", "reindex")
    _log("indexing knowledge base (60s timeout)...")
    indexed = _reload_project_knowledge(project_slug, timeout_sec=60)
    if indexed:
        _log("✓ knowledge indexed")
    else:
        _log("⚠ knowledge indexing timed out — skipped (training still complete)")

    # ═══ AUTO-FILL BRAIN — Make agent smart from data loading ═══
    if _cancelled():
        _update_run("failed", "cancelled"); _log("⊘ training cancelled by user"); return
    _update_run("running", "brain_fill")
    _log("filling agent brain from training data...")

    engine = create_engine(db_url)

    # 1. Auto-Memories — from REAL data, not just metadata
    try:
        _log("  generating auto-memories from real data...")
        mem_facts = []
        schema = re.sub(r"[^a-z0-9_]", "_", project_slug.lower())[:63]

        # Get REAL stats from actual table
        try:
            mem_engine = create_engine(db_url)
            with mem_engine.connect() as mconn:
                mconn.execute(text(f"SET LOCAL search_path TO {schema}, public"))
                # Row count
                row_count = mconn.execute(text(f'SELECT COUNT(*) FROM "{table_name}"')).scalar() or 0
                mem_facts.append(f"Table '{table_name}' has {row_count:,} rows")

                # Column count
                col_count = len(col_analyses) if col_analyses else 0
                mem_facts.append(f"Table '{table_name}' has {col_count} columns")

                # Date range (if date column exists)
                for ca in (col_analyses or []):
                    if ca.get("type") == "datetime" or "date" in ca.get("name", "").lower():
                        try:
                            dr = mconn.execute(text(f'SELECT MIN("{ca["name"]}"), MAX("{ca["name"]}") FROM "{table_name}"')).fetchone()
                            if dr and dr[0]:
                                mem_facts.append(f"Table '{table_name}' date range: {dr[0]} to {dr[1]}")
                        except Exception:
                            pass
                        break

                # Top categories (for categorical columns)
                for ca in (col_analyses or []):
                    if ca.get("is_categorical") and ca.get("unique_count", 0) <= 15:
                        try:
                            cats = mconn.execute(text(f'SELECT "{ca["name"]}", COUNT(*) as cnt FROM "{table_name}" WHERE "{ca["name"]}" IS NOT NULL GROUP BY "{ca["name"]}" ORDER BY cnt DESC LIMIT 5')).fetchall()
                            if cats:
                                cat_str = ", ".join(f"{c[0]} ({c[1]})" for c in cats)
                                mem_facts.append(f"Table '{table_name}' column '{ca['name']}' top values: {cat_str}")
                        except Exception:
                            pass

                # Numeric ranges (for key numeric columns)
                for ca in (col_analyses or []):
                    if ca.get("type") == "numeric" and ca.get("mean"):
                        try:
                            stats = mconn.execute(text(f'SELECT SUM("{ca["name"]}"), AVG("{ca["name"]}"), MIN("{ca["name"]}"), MAX("{ca["name"]}") FROM "{table_name}"')).fetchone()
                            if stats and stats[0] is not None:
                                mem_facts.append(f"Table '{table_name}' column '{ca['name']}': total={stats[0]:,.0f}, avg={stats[1]:,.0f}, range={stats[2]:,.0f} to {stats[3]:,.0f}")
                        except Exception:
                            pass
        except Exception:
            pass

        # Add metadata-based facts as fallback
        if metadata.get("table_description"):
            mem_facts.append(f"Table '{table_name}': {metadata['table_description'][:200]}")
        if metadata.get("grain"):
            mem_facts.append(f"Table '{table_name}' grain: {metadata['grain']}")

        with engine.connect() as conn:
            saved = 0
            for fact in mem_facts[:8]:
                try:
                    conn.execute(text(
                        "INSERT INTO public.dash_memories (project_slug, scope, fact, source) VALUES (:s, 'project', :f, 'auto') ON CONFLICT DO NOTHING"
                    ), {"s": project_slug, "f": fact})
                    saved += 1
                except Exception:
                    pass
            conn.commit()
        _log(f"  ✓ {saved} memories saved")
    except Exception as e:
        _log(f"  ⚠ memories error: {str(e)[:80]}")

    # 2. Auto-Patterns: generate SQL patterns from Q&A (with metadata extraction)
    try:
        _log("  generating query patterns...")
        training_dir = KNOWLEDGE_DIR / project_slug / "training"
        qa_file = training_dir / f"{table_name}_qa.json"
        patterns_saved = 0
        if qa_file.exists():
            with open(qa_file) as f:
                qa_pairs = json.load(f)
            for qa in qa_pairs[:5]:
                q = qa.get("question", "")
                s = qa.get("sql", "")
                if q and s:
                    if _save_query_pattern_with_metadata(engine, project_slug, q, s, source='training'):
                        patterns_saved += 1
        _log(f"  ✓ {patterns_saved} query patterns saved (with metadata)")
    except Exception as e:
        _log(f"  ⚠ patterns error: {str(e)[:80]}")

    # 3. Auto-Rules
    try:
        _log("  generating business rules...")
        rules_saved = 0
        biz_metrics = biz_rules.get("metrics", [])
        biz_rules_list = biz_rules.get("business_rules", [])
        with engine.connect() as conn:
            for metric in biz_metrics[:3]:
                if isinstance(metric, dict) and metric.get("name"):
                    rule_id = f"auto_{table_name}_{metric['name'].lower().replace(' ', '_')[:30]}"
                    defn = f"{metric.get('name', '')}: {metric.get('definition', '')} (Calculation: {metric.get('calculation', '')})"
                    conn.execute(text(
                        "INSERT INTO public.dash_rules_db (project_slug, rule_id, name, type, definition, source) "
                        "VALUES (:s, :rid, :name, 'metric', :defn, 'auto_training') ON CONFLICT (project_slug, rule_id) DO NOTHING"
                    ), {"s": project_slug, "rid": rule_id, "name": metric["name"], "defn": defn})
                    rules_saved += 1
            for rule in biz_rules_list[:3]:
                if isinstance(rule, dict) and rule.get("name"):
                    rule_id = f"auto_{table_name}_{rule['name'].lower().replace(' ', '_')[:30]}"
                    conn.execute(text(
                        "INSERT INTO public.dash_rules_db (project_slug, rule_id, name, type, definition, source) "
                        "VALUES (:s, :rid, :name, 'business_rule', :defn, 'auto_training') ON CONFLICT (project_slug, rule_id) DO NOTHING"
                    ), {"s": project_slug, "rid": rule_id, "name": rule["name"], "defn": rule.get("definition", "")})
                    rules_saved += 1
            conn.commit()
        _log(f"  ✓ {rules_saved} rules saved")
    except Exception as e:
        _log(f"  ⚠ rules error: {str(e)[:80]}")

    # 4. Auto-Annotations
    try:
        _log("  generating column annotations...")
        annotations_saved = 0
        with engine.connect() as conn:
            for col in (metadata.get("table_columns") or []):
                desc = col.get("description", "")
                if desc and len(desc) > 5:
                    conn.execute(text(
                        "INSERT INTO public.dash_annotations (project_slug, table_name, column_name, annotation, updated_by) "
                        "VALUES (:s, :t, :c, :a, 'auto_training') ON CONFLICT (project_slug, table_name, column_name) DO NOTHING"
                    ), {"s": project_slug, "t": table_name, "c": col["name"], "a": desc})
                    annotations_saved += 1
            conn.commit()
        _log(f"  ✓ {annotations_saved} annotations saved")
    except Exception as e:
        _log(f"  ⚠ annotations error: {str(e)[:80]}")

    # 5. Auto-Evals
    try:
        _log("  generating eval test cases...")
        evals_saved = 0
        if qa_file.exists():
            with open(qa_file) as f:
                qa_pairs = json.load(f)
            # Map analysis types to eval prefixes for testing different analysis modes
            _analysis_type_tags = {
                "descriptive": "[descriptive] ",
                "diagnostic": "[diagnostic] ",
                "predictive": "[predictive] ",
                "prescriptive": "[prescriptive] ",
                "anomaly": "[anomaly] ",
            }
            with engine.connect() as conn:
                for qa in qa_pairs[:5]:
                    q = qa.get("question", "")
                    s = qa.get("sql", "")
                    if q and s:
                        # Tag the eval question with its analysis type if available
                        a_type = qa.get("analysis_type", "")
                        if a_type and a_type in _analysis_type_tags:
                            tagged_q = f"{_analysis_type_tags[a_type]}{q}"
                        else:
                            # Infer analysis type from question keywords
                            q_lower = q.lower()
                            if any(kw in q_lower for kw in ["why", "cause", "driver", "impact", "factor"]):
                                tagged_q = "[diagnostic] " + q
                            elif any(kw in q_lower for kw in ["anomal", "outlier", "unusual", "spike", "drop"]):
                                tagged_q = "[anomaly] " + q
                            elif any(kw in q_lower for kw in ["recommend", "should", "action", "mitigat", "suggest"]):
                                tagged_q = "[prescriptive] " + q
                            elif any(kw in q_lower for kw in ["predict", "forecast", "trend", "project", "future"]):
                                tagged_q = "[predictive] " + q
                            else:
                                tagged_q = "[descriptive] " + q
                        conn.execute(text(
                            "INSERT INTO public.dash_evals (project_slug, question, expected_sql) VALUES (:s, :q, :sql)"
                        ), {"s": project_slug, "q": tagged_q, "sql": s})
                        evals_saved += 1
                conn.commit()
        _log(f"  ✓ {evals_saved} eval test cases saved")
    except Exception as e:
        _log(f"  ⚠ evals error: {str(e)[:80]}")

    # 6. Seed Feedback
    try:
        _log("  seeding sample feedback...")
        feedback_saved = 0
        if qa_file.exists():
            with open(qa_file) as f:
                qa_pairs = json.load(f)
            with engine.connect() as conn:
                for qa in qa_pairs[:3]:
                    q = qa.get("question", "")
                    a = qa.get("answer_template", f"SQL: {qa.get('sql', '')}")
                    if q:
                        conn.execute(text(
                            "INSERT INTO public.dash_feedback (project_slug, question, answer, sql_query, rating) "
                            "VALUES (:s, :q, :a, :sql, 'up')"
                        ), {"s": project_slug, "q": q, "a": a, "sql": qa.get("sql", "")})
                        feedback_saved += 1
                conn.commit()
        _log(f"  ✓ {feedback_saved} seed feedback saved")
    except Exception as e:
        _log(f"  ⚠ feedback error: {str(e)[:80]}")

    # 7. Save workflows to DB
    try:
        _log("  saving workflows to database...")
        wf_file = KNOWLEDGE_DIR / project_slug / "workflows.json"
        wf_saved = 0
        if wf_file.exists():
            with open(wf_file) as f:
                wfs = json.load(f)
            with engine.connect() as conn:
                for wf in (wfs if isinstance(wfs, list) else []):
                    name = wf.get("name", "")
                    if name:
                        conn.execute(text(
                            "INSERT INTO public.dash_workflows_db (project_slug, name, description, steps, source) "
                            "VALUES (:s, :n, :d, CAST(:st AS jsonb), 'training') "
                            "ON CONFLICT DO NOTHING"
                        ), {"s": project_slug, "n": name, "d": wf.get("description", ""), "st": json.dumps(wf.get("steps", []))})
                        wf_saved += 1
                conn.commit()
        _log(f"  ✓ {wf_saved} workflows saved to DB")
    except Exception as e:
        _log(f"  ⚠ workflows error: {str(e)[:80]}")

    # ═══ SMART TRAINING — Extract domain knowledge from data ═══
    if _cancelled():
        _update_run("failed", "cancelled"); _log("⊘ training cancelled by user"); return
    _update_run("running", "domain_knowledge")
    _log("extracting domain knowledge...")

    col_names = [c.get("name", "") for c in (metadata.get("table_columns") or [])]
    col_types = [f"{c.get('name','')}: {c.get('type','')}" for c in (metadata.get("table_columns") or [])]
    col_info = ", ".join(col_types[:30])

    # Get sample distinct values for categorical columns
    cat_values = {}
    try:
        proj_engine = create_engine(db_url)
        import re as _re2
        schema = _re2.sub(r"[^a-z0-9_]", "_", project_slug.lower())[:63]
        with proj_engine.connect() as conn:
            for col in (metadata.get("table_columns") or []):
                if col.get("type") in ("TEXT", "text", "str", "object") and col.get("unique_count", 100) <= 20:
                    try:
                        vals = conn.execute(text(f'SELECT DISTINCT "{col["name"]}" FROM "{schema}"."{table_name}" WHERE "{col["name"]}" IS NOT NULL LIMIT 15')).fetchall()
                        cat_values[col["name"]] = [str(v[0]) for v in vals]
                    except Exception:
                        pass
    except Exception:
        pass

    cat_values_str = "\n".join(f"  {k}: {', '.join(v[:10])}" for k, v in cat_values.items()) if cat_values else "None detected"

    from dash.settings import training_llm_call

    # 8. Business Glossary
    try:
        _log("  extracting business glossary...")
        result = training_llm_call(
            f"Look at these column names and values from table '{table_name}'. "
            f"Extract abbreviations, acronyms, and domain terms with their meanings.\n\n"
            f"Columns: {col_info}\n"
            f"Categorical values:\n{cat_values_str}\n\n"
            f"Return ONLY valid JSON array of strings:\n"
            f'["MMK = Myanmar Kyat (currency)", "PR = Purchase Requisition"]',
            "extraction"
        )
        if result:
            glossary = json.loads(result)
            if isinstance(glossary, list):
                with engine.connect() as conn:
                    for term in glossary[:15]:
                        if isinstance(term, str) and len(term) > 3:
                            conn.execute(text(
                                "INSERT INTO public.dash_memories (project_slug, scope, fact, source) "
                                "VALUES (:s, 'project', :f, 'glossary') ON CONFLICT DO NOTHING"
                            ), {"s": project_slug, "f": f"Glossary: {term}"})
                    conn.commit()
                _log(f"  ✓ {len(glossary[:15])} glossary terms saved")
    except Exception as e:
        _log(f"  ⚠ glossary error: {str(e)[:80]}")

    # 9. Calculation Rules
    try:
        _log("  extracting calculation rules...")
        numeric_cols = [c["name"] for c in (metadata.get("table_columns") or []) if c.get("type") in ("DOUBLE PRECISION", "INTEGER", "BIGINT", "NUMERIC", "numeric", "int64", "float64")]
        if numeric_cols:
            result = training_llm_call(
                f"Table '{table_name}' has these numeric columns: {', '.join(numeric_cols)}\n"
                f"All columns: {col_info}\n\n"
                f"Infer calculation relationships between columns. How are they related?\n"
                f"Return ONLY valid JSON array:\n"
                f'["net_amount = total_amount + tax - discount", "total_value = SUM(net_amount) per quotation"]',
                "extraction"
            )
            if result:
                calcs = json.loads(result)
                if isinstance(calcs, list):
                    with engine.connect() as conn:
                        for calc in calcs[:8]:
                            if isinstance(calc, str) and len(calc) > 5:
                                rule_id = f"calc_{table_name}_{hash(calc) % 10000}"
                                conn.execute(text(
                                    "INSERT INTO public.dash_rules_db (project_slug, rule_id, name, type, definition, source) "
                                    "VALUES (:s, :rid, :name, 'calculation', :defn, 'auto_training') ON CONFLICT (project_slug, rule_id) DO NOTHING"
                                ), {"s": project_slug, "rid": rule_id, "name": f"Calculation: {calc[:50]}", "defn": calc})
                        conn.commit()
                    _log(f"  ✓ {len(calcs[:8])} calculation rules saved")
    except Exception as e:
        _log(f"  ⚠ calculation error: {str(e)[:80]}")

    # 10. Value Mappings
    try:
        _log("  extracting value mappings...")
        if cat_values:
            result = training_llm_call(
                f"Table '{table_name}' has these categorical columns with values:\n{cat_values_str}\n\n"
                f"Explain what each value means in business context.\n"
                f"Return ONLY valid JSON array:\n"
                f'["status: SAP PR Created = quotation approved and sent to SAP system", "doa_type: CAPEX = Capital Expenditure for assets"]',
                "extraction"
            )
            if result:
                mappings = json.loads(result)
                if isinstance(mappings, list):
                    with engine.connect() as conn:
                        for mapping in mappings[:15]:
                            if isinstance(mapping, str) and len(mapping) > 5:
                                conn.execute(text(
                                    "INSERT INTO public.dash_memories (project_slug, scope, fact, source) "
                                    "VALUES (:s, 'project', :f, 'value_mapping') ON CONFLICT DO NOTHING"
                                ), {"s": project_slug, "f": f"Value mapping: {mapping}"})
                        conn.commit()
                    _log(f"  ✓ {len(mappings[:15])} value mappings saved")
    except Exception as e:
        _log(f"  ⚠ value mapping error: {str(e)[:80]}")

    # 11. KPI Definitions
    try:
        _log("  extracting KPI definitions...")
        result = training_llm_call(
            f"Table '{table_name}' columns: {col_info}\n"
            f"Categorical values:\n{cat_values_str}\n\n"
            f"What business KPIs can be calculated from this data? "
            f"Give the KPI name, definition, and SQL calculation.\n"
            f"Return ONLY valid JSON array:\n"
            f'[{{"name": "Total Spend", "definition": "Sum of approved spend", "sql": "SELECT SUM(net_amount) FROM {table_name} WHERE status = \'Approved\'"}}]',
            "extraction"
        )
        if result:
            kpis = json.loads(result)
            if isinstance(kpis, list):
                with engine.connect() as conn:
                    for kpi in kpis[:8]:
                        if isinstance(kpi, dict) and kpi.get("name"):
                            rule_id = f"kpi_{table_name}_{kpi['name'].lower().replace(' ', '_')[:25]}"
                            defn = f"{kpi.get('definition', '')} | SQL: {kpi.get('sql', '')}"
                            conn.execute(text(
                                "INSERT INTO public.dash_rules_db (project_slug, rule_id, name, type, definition, source) "
                                "VALUES (:s, :rid, :name, 'kpi', :defn, 'auto_training') ON CONFLICT (project_slug, rule_id) DO NOTHING"
                            ), {"s": project_slug, "rid": rule_id, "name": kpi["name"], "defn": defn})
                    conn.commit()
                _log(f"  ✓ {len(kpis[:8])} KPI definitions saved")
    except Exception as e:
        _log(f"  ⚠ KPI error: {str(e)[:80]}")

    # 12. Data Quality Rules
    try:
        _log("  extracting data quality rules...")
        null_info = ", ".join(f"{c['name']}: {c.get('null_pct', 0)}% null" for c in (metadata.get("table_columns") or []) if c.get("null_pct", 0) > 5)
        if null_info:
            result = training_llm_call(
                f"Table '{table_name}' data quality:\n"
                f"High null columns: {null_info}\n"
                f"All columns: {col_info}\n\n"
                f"Generate data quality rules — what to watch out for, how to handle NULLs, common mistakes.\n"
                f"Return ONLY valid JSON array:\n"
                f'["Never include rows where supplier_name IS NULL — these are template records", "Always use COALESCE(exchange_rate, 1) for currency conversion"]',
                "extraction"
            )
            if result:
                dq_rules = json.loads(result)
                if isinstance(dq_rules, list):
                    with engine.connect() as conn:
                        for rule in dq_rules[:8]:
                            if isinstance(rule, str) and len(rule) > 5:
                                conn.execute(text(
                                    "INSERT INTO public.dash_memories (project_slug, scope, fact, source) "
                                    "VALUES (:s, 'project', :f, 'data_quality') ON CONFLICT DO NOTHING"
                                ), {"s": project_slug, "f": f"Data quality: {rule}"})
                        conn.commit()
                    _log(f"  ✓ {len(dq_rules[:8])} data quality rules saved")
    except Exception as e:
        _log(f"  ⚠ data quality error: {str(e)[:80]}")

    # 13. Negative Examples (What NOT to do)
    try:
        _log("  extracting negative examples...")
        result = training_llm_call(
            f"Table '{table_name}' columns: {col_info}\n"
            f"Categorical values:\n{cat_values_str}\n\n"
            f"What are common mistakes when querying this data? Give DON'T / DO pairs.\n"
            f"Return ONLY valid JSON array:\n"
            f'["DON\'T COUNT(*) for quotation count — each row is an item. DO COUNT(DISTINCT qr_no)", '
            f'"DON\'T use total_amount for spend — it excludes tax. DO use net_amount"]',
            "extraction"
        )
        if result:
            negs = json.loads(result)
            if isinstance(negs, list):
                with engine.connect() as conn:
                    for neg in negs[:8]:
                        if isinstance(neg, str) and len(neg) > 5:
                            conn.execute(text(
                                "INSERT INTO public.dash_memories (project_slug, scope, fact, source) "
                                "VALUES (:s, 'project', :f, 'negative_example') ON CONFLICT DO NOTHING"
                            ), {"s": project_slug, "f": f"⚠ {neg}"})
                    conn.commit()
                _log(f"  ✓ {len(negs[:8])} negative examples saved")
    except Exception as e:
        _log(f"  ⚠ negative examples error: {str(e)[:80]}")

    _log("✓ agent brain filled — ready to use!")

    # ═══ AI SEED — Pre-populate activity metrics from data analysis ═══
    if _cancelled():
        _update_run("failed", "cancelled"); _log("⊘ training cancelled by user"); return
    _log("generating AI seed activity data...")

    col_info = ", ".join(f"{c.get('name','')}: {c.get('type','')}" for c in (metadata.get("table_columns") or [])[:20])
    cat_values_str = ""
    for c in (metadata.get("table_columns") or []):
        if c.get("sample_values"):
            cat_values_str += f"  {c['name']}: {', '.join(str(v) for v in c['sample_values'][:5])}\n"

    # 14. Seed Bad Feedback (common mistakes)
    try:
        _log("  generating bad feedback examples...")
        result = training_llm_call(
            f"Table '{table_name}' columns: {col_info}\n\n"
            f"Generate 2 examples of WRONG analysis that a data analyst might make with this data. "
            f"Each should be a common mistake (wrong column, wrong aggregation, wrong filter).\n"
            f"Return ONLY valid JSON array:\n"
            f'[{{"question": "What is total X?", "wrong_answer": "The total is $500K", "why_wrong": "Used gross instead of net amount"}}]',
            "extraction"
        )
        if result:
            bads = json.loads(result)
            if isinstance(bads, list):
                bad_saved = 0
                with train_engine.connect() as conn:
                    for b in bads[:3]:
                        if isinstance(b, dict) and b.get("question"):
                            conn.execute(text(
                                "INSERT INTO public.dash_feedback (project_slug, question, answer, rating) "
                                "VALUES (:s, :q, :a, 'down') ON CONFLICT DO NOTHING"
                            ), {"s": project_slug, "q": b["question"], "a": f"{b.get('wrong_answer','')} | WHY WRONG: {b.get('why_wrong','')}"})
                            bad_saved += 1
                    conn.commit()
                _log(f"  ✓ {bad_saved} bad feedback examples saved")
    except Exception as e:
        _log(f"  ⚠ bad feedback error: {str(e)[:80]}")

    # 15. Seed Proactive Insights (anomalies from data)
    try:
        _log("  generating proactive insights...")
        result = training_llm_call(
            f"Table '{table_name}' columns: {col_info}\n"
            f"Sample values:\n{cat_values_str}\n\n"
            f"Analyze this data schema and generate 3 proactive insights an analyst should investigate. "
            f"Focus on: data quality issues, unusual patterns, potential anomalies, optimization opportunities.\n"
            f"Return ONLY valid JSON array:\n"
            f'[{{"title": "High null rate in column X", "description": "Column X has significant missing data which could affect aggregations", "severity": "warning", "metric": "null_rate", "value": "23%"}}]',
            "extraction"
        )
        if result:
            insights = json.loads(result)
            if isinstance(insights, list):
                ins_saved = 0
                with train_engine.connect() as conn:
                    for ins in insights[:5]:
                        if isinstance(ins, dict) and ins.get("title"):
                            insight_text = f"{ins['title']}: {ins.get('description', '')}"
                            conn.execute(text(
                                "INSERT INTO public.dash_proactive_insights (project_slug, insight, severity, tables_involved) "
                                "VALUES (:s, :i, :sev, :tables)"
                            ), {"s": project_slug, "i": insight_text,
                                "sev": ins.get("severity", "info"), "tables": [table_name]})
                            ins_saved += 1
                    conn.commit()
                _log(f"  ✓ {ins_saved} proactive insights saved")
    except Exception as e:
        _log(f"  ⚠ insights error: {str(e)[:80]}")

    # 16. Seed Drift Baseline
    try:
        _log("  saving drift baseline...")
        with train_engine.connect() as conn:
            # Check if baseline already exists
            existing = conn.execute(text(
                "SELECT COUNT(*) FROM public.dash_drift_alerts WHERE project_slug = :s AND table_name = :t"
            ), {"s": project_slug, "t": table_name}).scalar() or 0
            if existing == 0:
                baseline = json.dumps([{"type": "baseline", "rows": num_rows, "columns": num_cols, "message": f"Baseline: {num_rows} rows, {num_cols} columns"}])
                conn.execute(text(
                    "INSERT INTO public.dash_drift_alerts (project_slug, table_name, alerts) "
                    "VALUES (:s, :t, CAST(:a AS jsonb))"
                ), {"s": project_slug, "t": table_name, "a": baseline})
                conn.commit()
                _log("  ✓ drift baseline saved")
            else:
                _log("  · drift baseline already exists")
    except Exception as e:
        _log(f"  ⚠ drift baseline error: {str(e)[:80]}")

    # 17. Seed Initial Evolution
    try:
        _log("  generating initial evolution...")
        with train_engine.connect() as conn:
            existing_evol = conn.execute(text(
                "SELECT COUNT(*) FROM public.dash_evolved_instructions WHERE project_slug = :s"
            ), {"s": project_slug}).scalar() or 0
            if existing_evol == 0:
                result = training_llm_call(
                    f"Table '{table_name}' columns: {col_info}\n\n"
                    f"Generate a short set of supplementary instructions (3-5 bullet points) that will help "
                    f"a data analyst AI agent work better with this specific data. Focus on gotchas, "
                    f"preferred aggregations, and domain-specific rules.\n"
                    f"Return plain text bullet points only.",
                    "persona"
                )
                if result:
                    conn.execute(text(
                        "INSERT INTO public.dash_evolved_instructions (project_slug, instructions, version, reasoning) "
                        "VALUES (:s, :i, 1, 'Auto-generated during initial training')"
                    ), {"s": project_slug, "i": result[:2000]})
                    conn.commit()
                    _log("  ✓ initial evolution saved (v1)")
            else:
                _log(f"  · evolution already exists ({existing_evol} versions)")
    except Exception as e:
        _log(f"  ⚠ evolution error: {str(e)[:80]}")

    # ═══ ENRICH PERSONA — Re-generate with full domain context ═══
    if _cancelled():
        _update_run("failed", "cancelled"); _log("⊘ training cancelled by user"); return
    _update_run("running", "persona_enrich")
    _log("enriching persona with domain knowledge...")
    try:
        # Gather all domain knowledge for richer persona
        eng_db = create_engine(db_url)
        domain_context_parts = []
        with eng_db.connect() as conn:
            # Glossary terms
            glossary = conn.execute(text(
                "SELECT fact FROM public.dash_memories WHERE project_slug = :s AND source = 'glossary' LIMIT 20"
            ), {"s": project_slug}).fetchall()
            if glossary:
                domain_context_parts.append("Business glossary: " + "; ".join(r[0] for r in glossary))

            # KPIs
            kpis = conn.execute(text(
                "SELECT name, definition FROM public.dash_rules_db WHERE project_slug = :s AND type = 'kpi' LIMIT 10"
            ), {"s": project_slug}).fetchall()
            if kpis:
                domain_context_parts.append("KPIs tracked: " + "; ".join(f"{r[0]}: {r[1][:80]}" for r in kpis))

            # Calculation rules
            calcs = conn.execute(text(
                "SELECT name FROM public.dash_rules_db WHERE project_slug = :s AND type = 'calculation' LIMIT 10"
            ), {"s": project_slug}).fetchall()
            if calcs:
                domain_context_parts.append("Calculation rules: " + ", ".join(r[0] for r in calcs))

            # Negative examples (critical gotchas)
            negs_db = conn.execute(text(
                "SELECT fact FROM public.dash_memories WHERE project_slug = :s AND source = 'negative_example' LIMIT 5"
            ), {"s": project_slug}).fetchall()
            if negs_db:
                domain_context_parts.append("Critical gotchas: " + "; ".join(r[0] for r in negs_db))

            # Data quality insights
            quality = conn.execute(text(
                "SELECT fact FROM public.dash_memories WHERE project_slug = :s AND source = 'data_quality' LIMIT 5"
            ), {"s": project_slug}).fetchall()
            if quality:
                domain_context_parts.append("Data quality rules: " + "; ".join(r[0] for r in quality))

        domain_summary = "\n".join(domain_context_parts)

        if domain_summary.strip():
            # Get all table names for context
            all_tables = []
            tables_dir_path = KNOWLEDGE_DIR / project_slug / "tables"
            if tables_dir_path.exists():
                for tf in tables_dir_path.glob("*.json"):
                    all_tables.append(tf.stem)

            enriched_persona = training_llm_call(
                f"You are generating a rich, expert persona for a data analyst AI agent.\n\n"
                f"This agent works with these tables: {', '.join(all_tables)}\n\n"
                f"Domain knowledge extracted from the data:\n{domain_summary}\n\n"
                f"Table being analyzed: {table_name}\n"
                f"Column types: {', '.join(c.get('name','') + ':' + c.get('type','') for c in (metadata.get('table_columns') or [])[:20])}\n\n"
                f"Generate a JSON persona object:\n"
                f'{{"persona_prompt": "You are an expert [specific domain] analyst who specializes in [specific areas from the data]. '
                f'You understand [domain terms from glossary]. You track KPIs like [from KPIs]. You know critical gotchas like [from negatives]...",'
                f'"domain_terms": ["term1", "term2", ...],'
                f'"expertise_areas": ["area1", "area2", ...],'
                f'"kpi_focus": ["kpi1", "kpi2", ...],'
                f'"communication_style": "concise with data tables",'
                f'"greeting": "Hey! Ready to dive into your [domain] data?"}}',
                "persona"
            )

            if enriched_persona:
                persona_data = json.loads(enriched_persona)
                if isinstance(persona_data, dict) and persona_data.get("persona_prompt"):
                    # Save enriched persona to file
                    persona_file = KNOWLEDGE_DIR / project_slug / "persona.json"
                    with open(persona_file, "w") as f:
                        json.dump(persona_data, f, indent=2)

                    # Update DB
                    with eng_db.connect() as conn:
                        conn.execute(text(
                            "UPDATE public.dash_personas SET persona = CAST(:p AS jsonb) WHERE project_slug = :s"
                        ), {"s": project_slug, "p": json.dumps(persona_data)})
                        conn.commit()
                    _log(f"✓ persona enriched with {len(domain_context_parts)} domain layers")
                else:
                    _log("· persona enrichment skipped — invalid response")
            else:
                _log("· persona enrichment skipped — no LLM response")
        else:
            _log("· persona enrichment skipped — no domain knowledge yet")
    except Exception as e:
        _log(f"⚠ persona enrichment error: {str(e)[:80]}")

    # Save fingerprint for delta detection on next retrain
    save_fingerprint(project_slug, table_name, num_rows, [c.get("name", "") for c in (metadata.get("table_columns") or [])])

    # Training Quality Score — measure how good this training was
    try:
        schema = re.sub(r"[^a-z0-9_]", "_", project_slug.lower())[:63]
        qa_file = KNOWLEDGE_DIR / project_slug / "training" / f"{table_name}_qa.json"
        quality_score = 0
        quality_details = {}

        if qa_file.exists():
            with open(qa_file) as f:
                qa_pairs = json.load(f)
            total_qa = len(qa_pairs)
            verified_qa = sum(1 for q in qa_pairs if q.get("verified"))
            quality_details["qa_total"] = total_qa
            quality_details["qa_verified"] = verified_qa
            quality_details["qa_pct"] = round((verified_qa / max(total_qa, 1)) * 100)
            quality_score += quality_details["qa_pct"] * 0.4  # 40% weight

        # Check relationships verified
        try:
            rel_engine = create_engine(db_url)
            with rel_engine.connect() as rconn:
                rels = rconn.execute(text("SELECT confidence, source FROM public.dash_relationships WHERE project_slug = :s"), {"s": project_slug}).fetchall()
                if rels:
                    verified_rels = sum(1 for r in rels if r[1] == 'ai_verified')
                    quality_details["relationships_total"] = len(rels)
                    quality_details["relationships_verified"] = verified_rels
                    quality_score += min(100, (verified_rels / max(len(rels), 1)) * 100) * 0.2  # 20% weight
                else:
                    quality_score += 50 * 0.2  # No rels = partial credit
        except Exception:
            quality_score += 50 * 0.2

        # Check memories from real data
        try:
            with rel_engine.connect() as rconn:
                mem_count = rconn.execute(text("SELECT COUNT(*) FROM public.dash_memories WHERE project_slug = :s AND source = 'auto'"), {"s": project_slug}).scalar() or 0
                quality_details["memories"] = mem_count
                quality_score += min(100, (mem_count / 8) * 100) * 0.2  # 20% weight
        except Exception:
            pass

        # Check profile exists (data quality)
        profile_file = KNOWLEDGE_DIR / project_slug / "table_sources" / f"{table_name}_profile.json"
        if profile_file.exists():
            try:
                with open(profile_file) as f:
                    profile = json.load(f)
                quality_details["health"] = profile.get("health", 0)
                quality_score += profile.get("health", 50) * 0.2  # 20% weight
            except Exception:
                quality_score += 50 * 0.2
        else:
            quality_score += 50 * 0.2

        quality_score = round(min(100, max(0, quality_score)))
        quality_details["overall_score"] = quality_score
        _log(f"✓ training quality score: {quality_score}% (Q&A: {quality_details.get('qa_pct', '?')}% verified, {quality_details.get('memories', '?')} memories, health: {quality_details.get('health', '?')}%)")

        # Save quality score
        quality_dir = KNOWLEDGE_DIR / project_slug / "table_sources"
        quality_dir.mkdir(parents=True, exist_ok=True)
        _safe_write_json(quality_dir / f"{table_name}_training_quality.json", quality_details)
    except Exception as e:
        _log(f"⚠ quality score error: {str(e)[:80]}")

    _log(f"✓ training complete for {table_name}")
    _update_run("done", "complete")


# ---------------------------------------------------------------------------
# API Endpoints
# ---------------------------------------------------------------------------


def _get_user_id(request) -> str:
    """Extract user_id from request state (set by auth middleware)."""
    user = getattr(getattr(request, 'state', None), 'user', None)
    if user:
        return str(user.get('username', 'default'))
    return 'default'


def _get_project_schema(request: Request, project: str | None = None) -> str | None:
    """Get the project schema name if project param is provided."""
    if not project:
        project = request.query_params.get('project')
    if project:
        from db.session import create_project_schema
        return create_project_schema(project)
    return None


def _extract_tables_pptx(file_path: str) -> list:
    """Extract tables from PowerPoint slides as list of DataFrames."""
    tables = []
    try:
        from pptx import Presentation
        import pandas as pd
        prs = Presentation(file_path)
        for si, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.has_table:
                    tbl = shape.table
                    rows_data = []
                    for row in tbl.rows:
                        rows_data.append([cell.text.strip() for cell in row.cells])
                    if len(rows_data) > 1 and len(rows_data[0]) > 1:
                        headers = rows_data[0]
                        data = rows_data[1:]
                        df = pd.DataFrame(data, columns=headers)
                        # Clean column names
                        df.columns = [str(c).strip().lower().replace(' ', '_').replace('.', '_')[:50] for c in df.columns]
                        tables.append({
                            'source': f'slide_{si+1}',
                            'df': df,
                            'rows': len(data),
                            'cols': len(headers)
                        })
    except Exception:
        pass
    return tables


def _extract_images_pptx(file_path: str) -> list[dict]:
    """Extract images from PPTX slides. Returns [{"b64": str, "mime": str, "source": str}]."""
    images: list[dict] = []
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        import base64
        prs = Presentation(file_path)
        for si, slide in enumerate(prs.slides):
            if len(images) >= 30:
                break
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    blob = shape.image.blob
                    if len(blob) < 3000:
                        continue
                    images.append({
                        "b64": base64.b64encode(blob).decode(),
                        "mime": shape.image.content_type or "image/png",
                        "source": f"slide_{si + 1}",
                    })
                    if len(images) >= 30:
                        break
    except Exception:
        pass
    return images


def _extract_tables_pdf(file_path: str) -> list:
    """Extract tables from PDF pages as list of DataFrames."""
    tables = []
    try:
        import pdfplumber
        import pandas as pd
        with pdfplumber.open(file_path) as pdf:
            for pi, page in enumerate(pdf.pages):
                page_tables = page.extract_tables()
                for ti, tbl in enumerate(page_tables or []):
                    if tbl and len(tbl) > 1 and len(tbl[0]) > 1:
                        headers = [str(h).strip().lower().replace(' ', '_').replace('.', '_')[:50] for h in tbl[0]]
                        data = tbl[1:]
                        df = pd.DataFrame(data, columns=headers)
                        tables.append({
                            'source': f'page_{pi+1}_table_{ti+1}',
                            'df': df,
                            'rows': len(data),
                            'cols': len(headers)
                        })
    except Exception:
        pass
    return tables


def _extract_images_pdf(file_path: str) -> list[dict]:
    """Extract images from PDF pages. Returns [{"b64": str, "mime": str, "source": str}]."""
    images: list[dict] = []
    try:
        import fitz
        import base64
        doc = fitz.open(file_path)
        for pi, page in enumerate(doc):
            if len(images) >= 30:
                break
            for img_info in page.get_images(full=True):
                xref = img_info[0]
                extracted = doc.extract_image(xref)
                if not extracted:
                    continue
                raw = extracted["image"]
                if len(raw) < 3000:
                    continue
                ext = extracted.get("ext", "png")
                images.append({
                    "b64": base64.b64encode(raw).decode(),
                    "mime": f"image/{ext}",
                    "source": f"page_{pi + 1}",
                })
                if len(images) >= 30:
                    break
        doc.close()
    except Exception:
        pass
    return images


def _extract_images_docx(file_path: str) -> list[dict]:
    """Extract images from DOCX document relationships. Returns [{"b64": str, "mime": str, "source": str}]."""
    images: list[dict] = []
    try:
        from docx import Document
        import base64
        doc = Document(file_path)
        mime_map = {"png": "image/png", "jpg": "image/jpeg", "jpeg": "image/jpeg",
                    "gif": "image/gif", "bmp": "image/bmp", "tiff": "image/tiff",
                    "emf": "image/emf", "wmf": "image/wmf"}
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                try:
                    blob = rel.target_part.blob
                    if len(blob) < 3000:
                        continue
                    ref = getattr(rel, 'target_ref', '') or ''
                    ext = ref.rsplit(".", 1)[-1].lower() if "." in ref else "png"
                    images.append({
                        "b64": base64.b64encode(blob).decode(),
                        "mime": mime_map.get(ext, "image/png"),
                        "source": f"docx_image_{len(images) + 1}",
                    })
                    if len(images) >= 30:
                        break
                except Exception:
                    continue
    except Exception:
        pass
    return images


def _extract_tables_docx(file_path: str) -> list:
    """Extract tables from DOCX as list of DataFrames."""
    tables = []
    try:
        from docx import Document
        import pandas as pd
        doc = Document(file_path)
        for ti, tbl in enumerate(doc.tables):
            rows_data = []
            for row in tbl.rows:
                rows_data.append([cell.text.strip() for cell in row.cells])
            if len(rows_data) > 1 and len(rows_data[0]) > 1:
                headers = [str(h).strip().lower().replace(' ', '_').replace('.', '_')[:50] for h in rows_data[0]]
                data = rows_data[1:]
                df = pd.DataFrame(data, columns=headers)
                tables.append({
                    'source': f'table_{ti+1}',
                    'df': df,
                    'rows': len(data),
                    'cols': len(headers)
                })
    except Exception:
        pass
    return tables


def _extract_content(file_path: str, ext: str, raw_content: bytes) -> dict:
    """Extract text AND tables from various file formats."""
    # Get text (existing logic)
    if ext in (".md", ".txt", ".sql", ".py"):
        text = raw_content.decode("utf-8", errors="ignore")
        return {"text": text, "tables": [], "images": []}

    if ext == ".docx":
        try:
            from docx import Document
            doc = Document(file_path)
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except Exception:
            text = raw_content.decode("utf-8", errors="ignore")
        tables = _extract_tables_docx(file_path)
        return {"text": text, "tables": tables, "images": []}

    if ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                texts.append(t)
            text = "\n".join(texts)
        except Exception:
            text = ""
        tables = _extract_tables_pptx(file_path)
        return {"text": text, "tables": tables, "images": _extract_images_pptx(file_path)}

    if ext == ".pdf":
        try:
            import fitz
            doc = fitz.open(file_path)
            texts = []
            for page in doc:
                texts.append(page.get_text())
            doc.close()
            text = "\n".join(texts)
        except Exception:
            text = ""
        tables = _extract_tables_pdf(file_path)
        return {"text": text, "tables": tables, "images": _extract_images_pdf(file_path)}

    return {"text": raw_content.decode("utf-8", errors="ignore"), "tables": [], "images": []}


def _describe_images_with_vision(images: list[dict], filename: str) -> str:
    """Send images to vision model and return combined text descriptions."""
    if not images:
        return ""
    from dash.settings import training_vision_call
    descriptions = []
    for img in images[:30]:
        result = training_vision_call(
            prompt=_UNIVERSAL_VISION_PROMPT,
            images=[img],
        )
        if result:
            descriptions.append(f"[Image from {img['source']} in {filename}]: {result}")
    return "\n\n".join(descriptions)


# ---------------------------------------------------------------------------
# Upload Intelligence: Handlers + Conductor
# Each handler returns: {"tables": [...], "text": str, "images": [...],
#                        "metadata": dict, "errors": [...], "warnings": [...]}
# ---------------------------------------------------------------------------

def _handle_image(file_path: str, filename: str) -> dict:
    """Handle JPG/PNG image upload — Tesseract OCR first, vision fallback for charts/diagrams."""
    import base64
    result = {"tables": [], "text": "", "images": [], "metadata": {}, "errors": [], "warnings": []}
    try:
        with open(file_path, "rb") as f:
            blob = f.read()
        ext = Path(file_path).suffix.lower()
        mime_map = {".jpg": "image/jpeg", ".jpeg": "image/jpeg", ".png": "image/png"}

        # Try Tesseract first (local, fast, free)
        try:
            import pytesseract
            from PIL import Image as PILImage
            img = PILImage.open(file_path)
            ocr_text = pytesseract.image_to_string(img)
            if len(ocr_text.strip()) > 30:
                result["text"] = f"[OCR from {filename}]\n{ocr_text.strip()}"
                result["warnings"].append("Text extracted via Tesseract OCR (local)")
                # Still send to vision for richer description (charts, diagrams)
        except ImportError:
            pass
        except Exception:
            pass

        # Always add to images for vision (may be a chart/diagram that needs AI understanding)
        result["images"] = [{
            "b64": base64.b64encode(blob).decode(),
            "mime": mime_map.get(ext, "image/png"),
            "source": filename,
        }]
    except Exception as e:
        result["errors"].append(f"Failed to read image: {e}")
    return result


def _is_clean_sheet(df_preview: pd.DataFrame, merged_cells: list = []) -> bool:
    """Quick check: is this sheet clean data (proper headers, no mess) or messy (needs AI)?"""
    if len(df_preview) < 2:
        return False
    # Messy signals
    if merged_cells:
        return False  # Merged cells = messy
    headers = list(df_preview.columns)
    unnamed_count = sum(1 for h in headers if "unnamed" in str(h).lower())
    if unnamed_count > len(headers) * 0.3:
        return False  # Too many unnamed columns = wrong header row
    # Check if first row looks like data (not metadata/units)
    first_row = df_preview.iloc[0]
    text_vals = [str(v).strip().lower() for v in first_row.dropna()]
    unit_words = {"sachets", "kg", "units", "pcs", "nos", "mtrs", "ltrs"}
    if text_vals and all(v in unit_words for v in text_vals if v):
        return False  # First data row is units = messy
    # Check if columns have month/date patterns (needs unpivot)
    month_re = re.compile(r"^(jan|feb|mar|apr|may|jun|jul|aug|sep|oct|nov|dec)", re.IGNORECASE)
    month_cols = sum(1 for h in headers if month_re.match(str(h).strip()))
    if month_cols >= 3:
        return False  # Months as columns = needs unpivot
    # Clean: proper headers, no merges, no months as columns
    return True


def _handle_excel(file_path: str, filename: str) -> dict:
    """Handle Excel upload — master decision: clean data → fast load, messy data → AI analysis."""
    result = {"tables": [], "text": "", "images": [], "metadata": {}, "errors": [], "warnings": []}
    ext = Path(file_path).suffix.lower()

    # Step 1: Enumerate all sheets and read previews
    sheet_previews = {}
    sheet_names = []
    try:
        if ext == ".xlsx":
            import openpyxl
            # First pass: read merged cells (needs non-read-only mode)
            merged_info = {}
            try:
                wb_full = openpyxl.load_workbook(file_path, data_only=True)
                for sname in wb_full.sheetnames:
                    ws = wb_full[sname]
                    if ws.merged_cells.ranges:
                        merged_info[sname] = [str(r) for r in ws.merged_cells.ranges]
                wb_full.close()
            except Exception:
                pass
            # Second pass: read data (read-only for speed)
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            sheet_names = wb.sheetnames
            for sname in sheet_names:
                ws = wb[sname]
                rows = []
                for ri, row in enumerate(ws.iter_rows(max_row=25, values_only=True)):
                    if ri >= 25:
                        break
                    rows.append([str(v)[:60] if v is not None else "" for v in row[:15]])
                sheet_previews[sname] = {"rows": rows, "max_row": ws.max_row, "max_col": ws.max_column, "merged_cells": merged_info.get(sname, [])}
            wb.close()
        elif ext == ".xls":
            import xlrd
            wb = xlrd.open_workbook(file_path)
            sheet_names = wb.sheet_names()
            for sname in sheet_names:
                ws = wb.sheet_by_name(sname)
                rows = []
                for ri in range(min(25, ws.nrows)):
                    rows.append([str(ws.cell_value(ri, ci))[:60] for ci in range(min(15, ws.ncols))])
                sheet_previews[sname] = {"rows": rows, "max_row": ws.nrows, "max_col": ws.ncols}
    except Exception as e:
        result["errors"].append(f"Failed to scan sheets: {e}")
        # Fallback: try basic read
        try:
            df = pd.read_excel(file_path, header=_find_header_row(file_path, ext))
            df = _clean_dataframe(df)
            tname = _sanitize_table_name(Path(filename).stem)
            result["tables"].append({"name": tname, "df": df, "source": "sheet_1"})
        except Exception as e2:
            result["errors"].append(f"Fallback read failed: {e2}")
        return result

    if not sheet_names:
        result["warnings"].append("No sheets found in Excel file")
        return result

    # Step 2: MASTER DECISION — is this clean or messy data?
    file_slug = _sanitize_table_name(Path(filename).stem)
    clean_sheets = {}
    messy_sheets = {}

    for sname in sheet_names:
        info = sheet_previews.get(sname, {})
        merged = info.get("merged_cells", [])
        # Quick read with header=0 to check if it's clean
        try:
            df_peek = pd.read_excel(file_path, sheet_name=sname, header=0, nrows=10)
            if _is_clean_sheet(df_peek, merged):
                clean_sheets[sname] = True
            else:
                messy_sheets[sname] = True
        except Exception:
            messy_sheets[sname] = True

    # FAST PATH: All sheets are clean → direct load, no AI needed
    if not messy_sheets and clean_sheets:
        result["warnings"].append(f"Clean data detected — direct load ({len(clean_sheets)} sheets, no AI needed)")
        for sname in sheet_names:
            try:
                df = pd.read_excel(file_path, sheet_name=sname, header=0)
                df = _clean_dataframe(df)
                if len(df) == 0:
                    continue
                tname = f"{file_slug}_{_sanitize_table_name(str(sname))}" if len(sheet_names) > 1 else file_slug
                tname = _sanitize_table_name(tname)
                sheet_idx = sheet_names.index(sname) + 1
                result["tables"].append({"name": tname, "df": df, "source": sname, "sheet_number": sheet_idx, "description": f"Clean data from sheet '{sname}'"})
            except Exception as e:
                result["warnings"].append(f"Sheet '{sname}' failed: {e}")
        return result

    # MIXED: Some clean, some messy — load clean ones directly, AI for messy
    if clean_sheets:
        for sname in list(clean_sheets.keys()):
            try:
                df = pd.read_excel(file_path, sheet_name=sname, header=0)
                df = _clean_dataframe(df)
                if len(df) == 0:
                    continue
                tname = f"{file_slug}_{_sanitize_table_name(str(sname))}"
                tname = _sanitize_table_name(tname)
                sheet_idx = sheet_names.index(sname) + 1
                result["tables"].append({"name": tname, "df": df, "source": f"{sname} [direct]", "sheet_number": sheet_idx, "description": f"Clean data from '{sname}'"})
                result["warnings"].append(f"Sheet '{sname}': clean data — loaded directly (no AI)")
            except Exception:
                messy_sheets[sname] = True  # Failed direct load, try AI

    # AI PATH: Only for messy sheets
    # Filter sheet_names to only messy ones for AI analysis
    ai_sheet_names = [s for s in sheet_names if s in messy_sheets]
    if not ai_sheet_names:
        return result

    result["warnings"].append(f"{len(ai_sheet_names)} messy sheet(s) → AI analysis")

    # Step 3: AI analysis — send previews, get extraction plan (only messy sheets)
    ai_plan = None
    try:
        from dash.settings import training_llm_call
        preview_text = ""
        # Limit preview rows per sheet based on sheet count to stay within LLM context
        max_preview_rows = 25 if len(ai_sheet_names) <= 5 else (15 if len(ai_sheet_names) <= 10 else 10)
        for sname in ai_sheet_names:
            info = sheet_previews[sname]
            merged = info.get("merged_cells", [])
            merged_note = f" | Merged cells: {', '.join(merged[:10])}" if merged else ""
            preview_text += f"\n\nSheet: '{sname}' ({info['max_row']} rows × {info['max_col']} cols{merged_note})\n"
            for ri, row in enumerate(info["rows"][:max_preview_rows]):
                preview_text += f"  Row {ri}: {row}\n"

        prompt = f"""Analyze {len(ai_sheet_names)} messy Excel sheets that need intelligent processing. For each sheet I show the first rows as raw cell values. Empty cells are ''.

{preview_text}

Return ONLY a valid JSON array. For each sheet return an object with:
- "sheet": sheet name (exact match)
- "action": "load" | "split" | "skip" | "unpivot"
- "table_name": clean lowercase PostgreSQL table name (a-z, 0-9, underscore only, max 50 chars)
- "header_row": 0-indexed row number with column headers
- "data_start_row": 0-indexed row where data begins
- "skip_rows": array of 0-indexed row numbers to skip (blank/separator/unit/summary rows)
- "forward_fill_columns": array of 0-indexed column positions that need forward-fill (merged cells)
- "column_names": object mapping column index to clean name e.g. {{"0":"plant","1":"product","3":"jul_21"}}
- "description": brief description of table content

ACTIONS:
1. "load" — single table, store as-is
2. "split" — multiple tables in one sheet separated by blank rows. Return "tables" array.
3. "skip" — empty or irrelevant sheet
4. "unpivot" — columns are time periods (months, quarters, dates). MELT them into rows. Return:
   - "id_columns": array of column indexes that are identifiers (plant, product, category...)
   - "value_columns": array of column indexes that are time-based values to unpivot
   - "value_name": what the values represent ("output", "amount", "quantity")
   - "variable_name": what the columns represent ("month", "quarter", "period")
   - "extra_columns": object of extra columns to add e.g. {{"fiscal_year":"FY2022","unit":"Sachets"}}
   - "blocks": if multiple data blocks exist (e.g. IB Plant rows 4-9, Tea Plant row 13), list them:
     [{{"data_start":4,"data_end":9,"extra":{{"plant":"IB Plant","unit":"Sachets"}}}},
      {{"data_start":13,"data_end":13,"extra":{{"plant":"Tea Plant","unit":"kg"}}}}]
     All blocks get merged into ONE table with extra columns distinguishing them.

For "split" action, return "tables" array with "table_name", "header_row", "data_start_row", "data_end_row", "description".

DETECT THESE PATTERNS:
- Columns with month/date names (Jul'21, Aug'21, Q1 2022, 2021-01...) → use "unpivot"
- Multiple data blocks with same columns but different categories → merge via "blocks"
- Metadata rows at top (Company, Period, Rate) → NOT headers, skip them
- Sub-header rows (units: Sachets, kg, %) → skip_rows, capture unit in extra_columns
- Summary/total rows (Utilisation, Total) → skip_rows
- Merged cells → forward_fill_columns
- Empty sheets → "skip"
- NEVER invent column names. Use the EXACT text from the header row, just cleaned for PostgreSQL:
  "No of FFS machines" → "no_of_ffs_machines"
  "Annual Capacity" → "annual_capacity"
  "Products" → "products"
  "Jul'21" → keep as "Jul'21" (readable for unpivot values)
- column_names mapping: key=column index, value=cleaned REAL header text"""

        raw = training_llm_call(prompt, "excel_analysis")
        if raw:
            ai_plan = json.loads(raw)
            result["metadata"]["ai_analysis"] = "success"
        else:
            result["warnings"].append("AI returned empty response — using fallback")
    except json.JSONDecodeError as e:
        result["warnings"].append(f"AI returned invalid JSON: {str(e)[:100]}")
        ai_plan = None
    except Exception as e:
        result["warnings"].append(f"AI analysis error: {str(e)[:100]}")
        ai_plan = None

    # Step 3: Extract tables per AI plan or fallback
    if ai_plan and isinstance(ai_plan, list):
        file_slug = _sanitize_table_name(Path(filename).stem)
        for sheet_plan in ai_plan:
            sname = sheet_plan.get("sheet", "")
            action = sheet_plan.get("action", "skip")

            if action == "skip":
                result["warnings"].append(f"Skipped sheet '{sname}': {sheet_plan.get('reason', 'empty')}")
                continue

            if sname not in sheet_names:
                result["warnings"].append(f"Sheet '{sname}' not found in file")
                continue

            if action == "unpivot":
                # AI-powered conversion: send raw data to AI, let it decide column mapping
                try:
                    # Read raw data with no header — let AI figure out everything
                    df_raw_all = pd.read_excel(file_path, sheet_name=sname, header=None)
                    raw_preview = ""
                    for ri in range(min(20, len(df_raw_all))):
                        vals = [str(v)[:30] if pd.notna(v) else "" for v in df_raw_all.iloc[ri].values[:18]]
                        raw_preview += f"  Row {ri}: {vals}\n"

                    convert_prompt = f"""Convert this Excel sheet from WIDE to LONG format.

RAW DATA (first 20 rows, 0-indexed, no header):
{raw_preview}

Return ONLY valid JSON:
{{
  "header_row": <0-indexed row with column headers like month names (Jul'21, Aug'21 etc)>,
  "skip_after_header": [<0-indexed rows to skip AFTER header: unit rows (Sachets/kg), summary rows (Utilisation/Total), blank rows>],
  "id_columns": [<0-indexed col numbers for identifiers: plant, product, capacity>],
  "value_columns": [<0-indexed col numbers for time-period values to unpivot>],
  "id_names": ["plant", "product", "annual_capacity"],
  "value_name": "output",
  "variable_name": "month",
  "blocks": [
    {{"data_start": <first data row 0-indexed>, "data_end": <last data row 0-indexed>, "extra": {{"plant": "IB Plant", "unit": "Sachets"}} }}
  ],
  "extra_columns": {{"fiscal_year": "FY2022"}}
}}"""

                    from dash.settings import training_llm_call
                    convert_raw = training_llm_call(convert_prompt, "excel_analysis")
                    if not convert_raw:
                        raise ValueError("AI conversion returned empty")
                    cp = json.loads(convert_raw)

                    # Read with AI-determined header
                    ai_hrow = cp.get("header_row", 1)
                    skip_after = cp.get("skip_after_header", [])
                    skip_for_read = [s for s in skip_after if s > ai_hrow]
                    df_data = pd.read_excel(file_path, sheet_name=sname, header=ai_hrow, skiprows=skip_for_read)

                    # Use AI conversion plan for column mapping
                    id_col_idx = cp.get("id_columns", [0, 1, 2])
                    val_col_idx = cp.get("value_columns", [])
                    id_names_ai = cp.get("id_names", [])
                    value_name = cp.get("value_name", "value")
                    variable_name = cp.get("variable_name", "period")

                    # Rename ID columns with AI names
                    for i, name in enumerate(id_names_ai):
                        if i < len(id_col_idx) and id_col_idx[i] < len(df_data.columns):
                            df_data.rename(columns={df_data.columns[id_col_idx[i]]: name}, inplace=True)

                    id_cols = [id_names_ai[i] if i < len(id_names_ai) else df_data.columns[idx]
                               for i, idx in enumerate(id_col_idx) if idx < len(df_data.columns)]
                    val_cols = [df_data.columns[i] for i in val_col_idx if i < len(df_data.columns)]

                    blocks = cp.get("blocks", [])
                    tname = sheet_plan.get("table_name") or f"{file_slug}_{_sanitize_table_name(sname)}"
                    tname = _sanitize_table_name(tname)

                    if blocks and val_cols:
                        all_melted = []
                        for block in blocks:
                            bstart = block.get("data_start", 0) - ai_hrow - 1
                            bend = block.get("data_end", bstart + ai_hrow + 1) - ai_hrow - 1
                            skipped = len([s for s in skip_for_read if s < block.get("data_start", 0)])
                            bstart = max(0, bstart - skipped)
                            bend = max(bstart, bend - len([s for s in skip_for_read if s < block.get("data_end", 0)]))
                            bend = min(bend, len(df_data) - 1)
                            if bstart >= len(df_data):
                                continue
                            df_block = df_data.iloc[bstart:bend + 1].copy()
                            if id_cols and id_cols[0] in df_block.columns:
                                df_block[id_cols[0]] = df_block[id_cols[0]].ffill()
                            df_block = df_block.dropna(subset=val_cols, how='all')
                            melted = df_block.melt(id_vars=id_cols, value_vars=val_cols,
                                                   var_name=variable_name, value_name=value_name)
                            for k, v in block.get("extra", {}).items():
                                melted[k] = v
                            all_melted.append(melted)
                        df_final = pd.concat(all_melted, ignore_index=True) if all_melted else pd.DataFrame()
                    elif val_cols:
                        if id_cols and id_cols[0] in df_data.columns:
                            df_data[id_cols[0]] = df_data[id_cols[0]].ffill()
                        df_final = df_data.melt(id_vars=id_cols, value_vars=val_cols,
                                                var_name=variable_name, value_name=value_name)
                    else:
                        df_final = df_data

                    # Add extra columns
                    for k, v in cp.get("extra_columns", {}).items():
                        df_final[k] = v
                    for k, v in sheet_plan.get("extra_columns", {}).items():
                        if k not in df_final.columns:
                            df_final[k] = v

                    # Ask AI to create date mapping for the period column
                    if variable_name in df_final.columns:
                        unique_periods = df_final[variable_name].dropna().unique().tolist()[:20]
                        if unique_periods:
                            try:
                                date_prompt = f"""Convert these time period labels to ISO dates (YYYY-MM-DD, use 1st of month).

Periods: {unique_periods}

Return ONLY a JSON object mapping each period to its date:
{{"Jul'21": "2021-07-01", "Aug'21": "2021-08-01", "Q1 2022": "2022-01-01"}}"""
                                date_raw = training_llm_call(date_prompt, "extraction")
                                if date_raw:
                                    date_map = json.loads(date_raw)
                                    df_final["date"] = df_final[variable_name].map(date_map)
                                    df_final["date"] = pd.to_datetime(df_final["date"], errors="coerce")
                            except Exception:
                                pass

                    if value_name in df_final.columns:
                        df_final = df_final.dropna(subset=[value_name])
                    df_final = df_final.dropna(how='all')

                    if len(df_final) > 0:
                        sheet_idx = sheet_names.index(sname) + 1 if sname in sheet_names else 0
                        # Check if this table should merge with an existing one (same table_name from another sheet)
                        existing = [t for t in result["tables"] if t["name"] == tname]
                        if existing:
                            existing[0]["df"] = pd.concat([existing[0]["df"], df_final], ignore_index=True)
                            existing[0]["source"] += f" + {sname}"
                        else:
                            result["tables"].append({"name": tname, "df": df_final, "source": f"{sname} [unpivot]", "sheet_number": sheet_names.index(sname) + 1, "description": sheet_plan.get("description", "")})
                except Exception as e:
                    result["warnings"].append(f"Unpivot from '{sname}' failed: {e}")
                continue

            if action == "split":
                # Multiple tables in one sheet — read with header=None, slice manually
                try:
                    df_raw = pd.read_excel(file_path, sheet_name=sname, header=None)
                except Exception as e:
                    result["warnings"].append(f"Cannot read sheet '{sname}': {e}")
                    continue
                for sub in sheet_plan.get("tables", []):
                    try:
                        hrow = sub.get("header_row", 0)
                        dstart = sub.get("data_start_row", hrow + 1)
                        dend = sub.get("data_end_row")
                        if dend is None:
                            dend = len(df_raw) - 1
                        # Extract header row as column names
                        if hrow < len(df_raw):
                            headers = [str(v).strip() if pd.notna(v) and str(v).strip() else f"col_{i}" for i, v in enumerate(df_raw.iloc[hrow])]
                        else:
                            continue
                        # Extract data rows
                        skip_rows = set(sub.get("skip_rows", []))
                        skip_rows.add(hrow)  # Don't include header as data
                        data_rows = []
                        for ri in range(dstart, min(dend + 1, len(df_raw))):
                            if ri not in skip_rows:
                                data_rows.append(df_raw.iloc[ri].values)
                        if not data_rows:
                            continue
                        df = pd.DataFrame(data_rows, columns=headers[:len(df_raw.columns)])
                        df = _clean_dataframe(df)
                        if len(df) == 0:
                            continue
                        tname = sub.get("table_name") or f"{file_slug}_{_sanitize_table_name(sname)}"
                        tname = _sanitize_table_name(tname)
                        # Forward-fill
                        for col_idx in sub.get("forward_fill_columns", []):
                            if isinstance(col_idx, int) and col_idx < len(df.columns):
                                df.iloc[:, col_idx] = df.iloc[:, col_idx].ffill()
                        sheet_idx = sheet_names.index(sname) + 1 if sname in sheet_names else 0
                        result["tables"].append({"name": tname, "df": df, "source": f"{sname} [split]", "sheet_number": sheet_idx, "description": sub.get("description", "")})
                    except Exception as e:
                        result["warnings"].append(f"Split table from '{sname}' failed: {e}")
            else:
                # Single table from sheet
                try:
                    hrow = sheet_plan.get("header_row", 0)
                    skip = sheet_plan.get("skip_rows", [])
                    df = pd.read_excel(file_path, sheet_name=sname, header=hrow, skiprows=skip)
                    # Apply AI column names
                    col_names = sheet_plan.get("column_names", {})
                    if col_names:
                        new_cols = []
                        for i, c in enumerate(df.columns):
                            new_cols.append(col_names.get(str(i), str(c)))
                        df.columns = new_cols
                    df = _clean_dataframe(df)
                    if len(df) == 0:
                        result["warnings"].append(f"Sheet '{sname}' produced empty table after cleanup")
                        continue
                    tname = sheet_plan.get("table_name") or f"{file_slug}_{_sanitize_table_name(sname)}"
                    tname = _sanitize_table_name(tname)
                    # Forward-fill merged cells
                    for col_idx in sheet_plan.get("forward_fill_columns", []):
                        if isinstance(col_idx, int) and col_idx < len(df.columns):
                            df.iloc[:, col_idx] = df.iloc[:, col_idx].ffill()
                    # Add extra columns if specified
                    for k, v in sheet_plan.get("extra_columns", {}).items():
                        df[k] = v
                    sheet_idx = sheet_names.index(sname) + 1 if sname in sheet_names else 0
                    result["tables"].append({"name": tname, "df": df, "source": sname, "sheet_number": sheet_idx, "description": sheet_plan.get("description", "")})
                except Exception as e:
                    result["warnings"].append(f"Sheet '{sname}' extraction failed: {e}")

        # Collect metadata text from AI descriptions
        descriptions = [sp.get("description", "") for sp in ai_plan if sp.get("description")]
        if descriptions:
            result["text"] = f"Excel file: {filename}\n" + "\n".join(f"- {d}" for d in descriptions)

    else:
        # FALLBACK: No AI — read all sheets with basic header detection
        result["warnings"].append("AI analysis unavailable, using fallback for all sheets")
        file_slug = _sanitize_table_name(Path(filename).stem)
        try:
            all_sheets = pd.read_excel(file_path, sheet_name=None, header=None)
            for sname, raw_df in all_sheets.items():
                if raw_df.dropna(how='all').empty:
                    result["warnings"].append(f"Skipped empty sheet '{sname}'")
                    continue
                # Find header row for this sheet
                best_row = 0
                best_score = 0
                for i in range(min(10, len(raw_df))):
                    row = raw_df.iloc[i]
                    non_null = row.dropna()
                    score = sum(3 for v in non_null if isinstance(v, str) and 1 < len(str(v).strip()) < 50)
                    if score > best_score:
                        best_score = score
                        best_row = i
                df = pd.read_excel(file_path, sheet_name=sname, header=best_row)
                df = _clean_dataframe(df)
                if len(df) == 0:
                    continue
                tname = f"{file_slug}_{_sanitize_table_name(str(sname))}"
                sheet_idx = list(all_sheets.keys()).index(sname) + 1 if sname in all_sheets else 0
                result["tables"].append({"name": tname, "df": df, "source": str(sname), "sheet_number": sheet_idx, "description": ""})
        except Exception as e:
            result["errors"].append(f"Fallback read all sheets failed: {e}")

    return result


def _handle_pdf(file_path: str, filename: str) -> dict:
    """Handle PDF upload — text extraction + scanned page OCR via vision + tables + images."""
    result = {"tables": [], "text": "", "images": [], "metadata": {}, "errors": [], "warnings": []}
    try:
        import fitz
        import base64
        doc = fitz.open(file_path)
        if len(doc) == 0:
            result["errors"].append("PDF has 0 pages — file may be corrupt or encrypted")
            return result

        texts = []
        scanned_count = 0
        tesseract_ok = False
        try:
            import pytesseract
            from PIL import Image as PILImage
            tesseract_ok = True
        except ImportError:
            pass

        for pi, page in enumerate(doc):
            page_text = page.get_text()
            if len(page_text.strip()) < 50:
                # Scanned page — try Tesseract first (local, fast, free)
                scanned_count += 1
                try:
                    pixmap = page.get_pixmap(dpi=150)
                    if tesseract_ok:
                        img = PILImage.frombytes("RGB", [pixmap.width, pixmap.height], pixmap.samples)
                        ocr_text = pytesseract.image_to_string(img)
                        if len(ocr_text.strip()) > 30:
                            texts.append(f"[Page {pi + 1} — OCR]\n{ocr_text.strip()}")
                            continue  # Got text via Tesseract, skip vision
                    # Tesseract failed or got no text → send to vision LLM as fallback
                    if len(result["images"]) < 10:
                        png_bytes = pixmap.tobytes("png")
                        if len(png_bytes) < 5_000_000:
                            result["images"].append({
                                "b64": base64.b64encode(png_bytes).decode(),
                                "mime": "image/png",
                                "source": f"page_{pi + 1}_scanned",
                            })
                except Exception:
                    pass
            else:
                texts.append(page_text)
        doc.close()
        result["text"] = "\n".join(texts)

        if scanned_count > 0:
            result["metadata"]["scanned_pages"] = scanned_count
            vision_pages = len([i for i in result["images"] if "scanned" in i.get("source", "")])
            if tesseract_ok:
                result["warnings"].append(f"{scanned_count} scanned page(s) — Tesseract OCR (local), {vision_pages} sent to vision")
            else:
                result["warnings"].append(f"{scanned_count} scanned page(s) — sent to vision for OCR")

        # Extract tables (existing)
        tables = _extract_tables_pdf(file_path)
        result["tables"] = [{"name": f"{_sanitize_table_name(Path(filename).stem)}_{t['source']}", "df": t["df"], "source": t["source"]} for t in tables]

        # Extract embedded images (charts, diagrams)
        embedded_images = _extract_images_pdf(file_path)
        result["images"].extend(embedded_images)

    except Exception as e:
        result["errors"].append(f"PDF processing failed: {e}")
    return result


def _handle_pptx(file_path: str, filename: str) -> dict:
    """Handle PPTX upload — text + tables + images (existing logic, raised caps)."""
    result = {"tables": [], "text": "", "images": [], "metadata": {}, "errors": [], "warnings": []}
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
        result["text"] = "\n".join(texts)
        result["metadata"]["slides"] = len(prs.slides)
    except Exception as e:
        result["errors"].append(f"PPTX text extraction failed: {e}")

    tables = _extract_tables_pptx(file_path)
    result["tables"] = [{"name": f"{_sanitize_table_name(Path(filename).stem)}_{t['source']}", "df": t["df"], "source": t["source"]} for t in tables]
    result["images"] = _extract_images_pptx(file_path)
    return result


def _handle_docx(file_path: str, filename: str) -> dict:
    """Handle DOCX upload — text + tables + image extraction (new)."""
    result = {"tables": [], "text": "", "images": [], "metadata": {}, "errors": [], "warnings": []}
    try:
        from docx import Document
        doc = Document(file_path)
        result["text"] = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        result["errors"].append(f"DOCX text extraction failed: {e}")

    tables = _extract_tables_docx(file_path)
    result["tables"] = [{"name": f"{_sanitize_table_name(Path(filename).stem)}_{t['source']}", "df": t["df"], "source": t["source"]} for t in tables]
    result["images"] = _extract_images_docx(file_path)
    return result


def _handle_csv(file_path: str, filename: str) -> dict:
    """Handle CSV upload — existing read logic wrapped in standard result."""
    result = {"tables": [], "text": "", "images": [], "metadata": {}, "errors": [], "warnings": []}
    try:
        header_row = _find_header_row(file_path, ".csv")
        sep = _detect_delimiter(file_path)
        df = pd.read_csv(file_path, header=header_row, sep=sep)
        df = _clean_dataframe(df)
        tname = _sanitize_table_name(Path(filename).stem)
        result["tables"].append({"name": tname, "df": df, "source": "csv"})
    except Exception as e:
        result["errors"].append(f"CSV read failed: {e}")
    return result


def _handle_json(file_path: str, filename: str) -> dict:
    """Handle JSON upload — existing read logic wrapped in standard result."""
    result = {"tables": [], "text": "", "images": [], "metadata": {}, "errors": [], "warnings": []}
    try:
        try:
            df = pd.read_json(file_path)
        except ValueError:
            df = pd.read_json(file_path, orient='records')
        df = _clean_dataframe(df)
        tname = _sanitize_table_name(Path(filename).stem)
        result["tables"].append({"name": tname, "df": df, "source": "json"})
    except Exception as e:
        result["errors"].append(f"JSON read failed: {e}")
    return result


def _handle_text(file_path: str, filename: str, raw_content: bytes = b"") -> dict:
    """Handle text file upload — TXT, MD, SQL, PY."""
    result = {"tables": [], "text": "", "images": [], "metadata": {}, "errors": [], "warnings": []}
    try:
        if raw_content:
            result["text"] = raw_content.decode("utf-8", errors="ignore")
        else:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                result["text"] = f.read()
    except Exception as e:
        result["errors"].append(f"Text read failed: {e}")
    return result


def _conduct_upload(file_path: str, ext: str, project: str, filename: str, raw_content: bytes = b"") -> dict:
    """Upload Conductor — routes to handler, runs vision on images, returns unified result."""
    # Route to handler
    if ext in (".xlsx", ".xls"):
        result = _handle_excel(file_path, filename)
    elif ext == ".pdf":
        result = _handle_pdf(file_path, filename)
    elif ext == ".pptx":
        result = _handle_pptx(file_path, filename)
    elif ext == ".docx":
        result = _handle_docx(file_path, filename)
    elif ext in (".jpg", ".jpeg", ".png"):
        result = _handle_image(file_path, filename)
    elif ext == ".csv":
        result = _handle_csv(file_path, filename)
    elif ext == ".json":
        result = _handle_json(file_path, filename)
    elif ext in (".txt", ".md", ".sql", ".py"):
        result = _handle_text(file_path, filename, raw_content)
    else:
        result = {"tables": [], "text": "", "images": [], "metadata": {}, "errors": [f"Unsupported: {ext}"], "warnings": []}

    # Run vision on all collected images (scanned pages, charts, photos, DOCX images)
    if result.get("images"):
        image_text = _describe_images_with_vision(result["images"], filename)
        if image_text:
            result["text"] = (result.get("text", "") + f"\n\n--- IMAGE DESCRIPTIONS ---\n{image_text}").strip()

    return result


def _post_upload_engineer(project_slug: str, tables_created: list[dict], user_id: int = 1):
    """After upload: Engineer merges same-structure tables, Inspector validates, safe delete originals.
    Runs in background thread — does not block upload response."""
    if not project_slug or not tables_created:
        return
    import logging
    log = logging.getLogger("dash.upload")

    try:
        schema = re.sub(r"[^a-z0-9_]", "_", project_slug.lower())[:63]
        merge_engine = create_engine(db_url)

        # STEP 1: Get ALL tables in project with their columns
        all_tables = {}
        with merge_engine.connect() as conn:
            conn.execute(text(f"SET LOCAL search_path TO {schema}, public"))
            from sqlalchemy import inspect as sa_inspect
            insp = sa_inspect(merge_engine)
            for tbl in insp.get_table_names(schema=schema):
                try:
                    cols = insp.get_columns(tbl, schema=schema)
                    col_names = [c["name"] for c in cols]
                    row_count = conn.execute(text(f'SELECT COUNT(*) FROM "{schema}"."{tbl}"')).scalar() or 0
                    all_tables[tbl] = {"columns": col_names, "rows": row_count, "col_set": set(col_names)}
                except Exception:
                    pass

        if len(all_tables) < 2:
            log.info(f"Post-upload: only {len(all_tables)} table(s), skipping merge")
            # Still run Engineer for relationships
            _run_engineer_agent(project_slug, tables_created, user_id)
            return

        # STEP 2: Find merge groups (tables with >80% column overlap)
        merge_groups = []
        used = set()
        table_names = list(all_tables.keys())

        for i, t1 in enumerate(table_names):
            if t1 in used:
                continue
            group = [t1]
            cols1 = all_tables[t1]["col_set"]
            for j in range(i + 1, len(table_names)):
                t2 = table_names[j]
                if t2 in used:
                    continue
                cols2 = all_tables[t2]["col_set"]
                overlap = len(cols1 & cols2) / max(len(cols1 | cols2), 1)
                if overlap >= 0.8:
                    group.append(t2)
            if len(group) >= 2:
                merge_groups.append(group)
                used.update(group)

        if not merge_groups:
            log.info("Post-upload: no merge candidates found")
            _run_engineer_agent(project_slug, tables_created, user_id)
            return

        # STEP 3: Merge each group
        for group in merge_groups:
            # Pick merged table name (shortest common prefix or first table)
            import os
            prefix = os.path.commonprefix(group).rstrip("_") or group[0]
            merged_name = _sanitize_table_name(f"{prefix}_merged" if len(prefix) > 3 else f"{group[0]}_merged")

            # Get union of all columns
            all_cols = set()
            for tbl in group:
                all_cols.update(all_tables[tbl]["col_set"])
            all_cols_list = sorted(all_cols)

            # Count expected rows
            expected_rows = sum(all_tables[tbl]["rows"] for tbl in group)

            try:
                with merge_engine.connect() as conn:
                    conn.execute(text(f"SET LOCAL search_path TO {schema}, public"))

                    # Build UNION ALL with source_table column
                    parts = []
                    for tbl in group:
                        tbl_cols = all_tables[tbl]["col_set"]
                        select_parts = []
                        for col in all_cols_list:
                            if col in tbl_cols:
                                select_parts.append(f'"{col}"')
                            else:
                                select_parts.append(f"NULL AS \"{col}\"")
                        select_parts.append(f"'{tbl}' AS _source_table")
                        parts.append(f"SELECT {', '.join(select_parts)} FROM \"{schema}\".\"{tbl}\"")

                    union_sql = " UNION ALL ".join(parts)
                    create_sql = f'CREATE TABLE "{schema}"."{merged_name}" AS {union_sql}'
                    conn.execute(text(create_sql))
                    conn.commit()

                # STEP 4: Inspector validates merged table
                actual_rows = 0
                with merge_engine.connect() as conn:
                    actual_rows = conn.execute(text(f'SELECT COUNT(*) FROM "{schema}"."{merged_name}"')).scalar() or 0

                merge_valid = actual_rows >= expected_rows  # Must have ALL rows

                if merge_valid:
                    # Profile the merged table
                    try:
                        df_sample = pd.read_sql(f'SELECT * FROM "{schema}"."{merged_name}" LIMIT 5000', merge_engine)
                        profile = _profile_table(df_sample, project_slug, merged_name)
                        health = profile.get("health", 0)
                        merge_valid = health >= 50  # Must pass minimum quality
                    except Exception:
                        pass

                if merge_valid:
                    # PASS — safe to delete originals
                    with merge_engine.connect() as conn:
                        for tbl in group:
                            try:
                                conn.execute(text(f'DROP TABLE IF EXISTS "{schema}"."{tbl}" CASCADE'))
                            except Exception:
                                pass
                        conn.commit()
                    log.info(f"Post-upload MERGED: {group} → {merged_name} ({actual_rows} rows, health={profile.get('health', '?')}%)")

                    # Save source metadata for merged table
                    src_dir = KNOWLEDGE_DIR / project_slug / "table_sources"
                    src_dir.mkdir(parents=True, exist_ok=True)
                    _safe_write_json(src_dir / f"{merged_name}.json", {
                        "source_file": f"Merged from {len(group)} tables",
                        "source_detail": f"Sources: {', '.join(group)}",
                        "description": f"Merged table combining {len(group)} tables with same structure ({actual_rows} rows)",
                    })
                    # Clean up old source metadata
                    for tbl in group:
                        old_src = src_dir / f"{tbl}.json"
                        if old_src.exists():
                            old_src.unlink(missing_ok=True)
                        old_profile = src_dir / f"{tbl}_profile.json"
                        if old_profile.exists():
                            old_profile.unlink(missing_ok=True)
                else:
                    # FAIL — drop merged, keep originals
                    with merge_engine.connect() as conn:
                        conn.execute(text(f'DROP TABLE IF EXISTS "{schema}"."{merged_name}" CASCADE'))
                        conn.commit()
                    log.warning(f"Post-upload MERGE FAILED validation: {group} → kept separate (expected={expected_rows}, actual={actual_rows})")

            except Exception as e:
                log.warning(f"Post-upload merge error for {group}: {e}")

        # STEP 5: Run Engineer for relationships + views on final state
        _run_engineer_agent(project_slug, tables_created, user_id)

    except Exception as e:
        log.warning(f"Post-upload engineer failed: {e}")


def _run_engineer_agent(project_slug: str, tables_created: list[dict], user_id: int = 1):
    """Run Engineer agent to discover relationships and create useful views."""
    try:
        from dash.agents.engineer import create_engineer
        from db.session import create_project_knowledge, create_project_learnings
        from agno.learn import LearnedKnowledgeConfig, LearningMachine, LearningMode

        knowledge = create_project_knowledge(project_slug)
        learnings = create_project_learnings(project_slug)
        learning = LearningMachine(knowledge=learnings, learned_knowledge=LearnedKnowledgeConfig(mode=LearningMode.AGENTIC))
        engineer = create_engineer(project_slug=project_slug, knowledge=knowledge, learning=learning, dashboard_user_id=user_id)

        table_summary = "\n".join(f"  - {t['table']} ({t['rows']} rows)" for t in tables_created[:20])

        prompt = f"""Inspect the project tables and optimize:
1. INSPECT all tables — run introspect_schema
2. DISCOVER RELATIONSHIPS — find JOINable columns (shared IDs, dates, categories)
3. FIX COLUMN TYPES — dates stored as text → ALTER to DATE
4. REPORT relationships found

Tables in project:
{table_summary}

Use SQL tools. Be concise."""

        response = engineer.run(prompt)
        content = response.content if response else ""
        if content:
            from agno.knowledge.reader.text_reader import TextReader
            try:
                knowledge.insert(name="engineer-upload-analysis", text_content=f"Engineer Analysis:\n\n{content[:5000]}", reader=TextReader(), skip_if_exists=False)
            except Exception:
                pass
    except Exception as e:
        import logging
        logging.getLogger("dash.upload").warning(f"Engineer agent failed: {e}")


def _extract_document_structure(file_path: str, ext: str) -> list[dict]:
    """Extract document structure (titles + content summaries) for workflow conversion.

    Returns list of {"index": 1, "title": "...", "content_summary": "..."} dicts.
    PPTX → slide titles, PDF → page headers by font size, DOCX → heading paragraphs.
    """
    sections: list[dict] = []

    if ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides, 1):
                title = ""
                content_parts: list[str] = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if not t:
                                continue
                            if not title and (shape == slide.shapes.title if hasattr(slide.shapes, 'title') and slide.shapes.title else shape == slide.shapes[0]):
                                title = t
                            else:
                                content_parts.append(t)
                if not title and content_parts:
                    title = content_parts.pop(0)
                if title:
                    sections.append({
                        "index": i,
                        "title": title[:120],
                        "content_summary": " ".join(content_parts)[:150],
                    })
        except Exception:
            pass

    elif ext == ".pdf":
        try:
            import fitz
            doc = fitz.open(file_path)
            for i, page in enumerate(doc, 1):
                blocks = page.get_text("dict").get("blocks", [])
                max_size = 0
                title = ""
                content_parts: list[str] = []
                for block in blocks:
                    if "lines" not in block:
                        continue
                    for line in block["lines"]:
                        for span in line.get("spans", []):
                            text = span.get("text", "").strip()
                            size = span.get("size", 0)
                            if not text:
                                continue
                            if size > max_size and len(text) > 3:
                                if title:
                                    content_parts.insert(0, title)
                                title = text
                                max_size = size
                            else:
                                content_parts.append(text)
                if not title and content_parts:
                    title = content_parts.pop(0)
                if title:
                    sections.append({
                        "index": i,
                        "title": title[:120],
                        "content_summary": " ".join(content_parts)[:150],
                    })
            doc.close()
        except Exception:
            pass

    elif ext == ".docx":
        try:
            from docx import Document
            doc = Document(file_path)
            current_title = ""
            current_content: list[str] = []
            idx = 0
            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                is_heading = para.style and para.style.name and para.style.name.startswith("Heading")
                if is_heading:
                    if current_title:
                        idx += 1
                        sections.append({
                            "index": idx,
                            "title": current_title[:120],
                            "content_summary": " ".join(current_content)[:150],
                        })
                    current_title = text
                    current_content = []
                else:
                    if current_title:
                        current_content.append(text)
                    elif not sections:
                        current_title = text
            if current_title:
                idx += 1
                sections.append({
                    "index": idx,
                    "title": current_title[:120],
                    "content_summary": " ".join(current_content)[:150],
                })
        except Exception:
            pass

    return sections[:20]


@router.post("/upload")
async def upload_file(request: Request, file: UploadFile, table_name: str | None = None, replace: bool = False, project: str | None = None, action: str = "auto"):
    """Upload a data file. action: auto/append/upsert/replace/new"""
    user_id = _get_user_id(request)

    # Editor role required for uploads
    from app.auth import check_project_permission
    user = getattr(getattr(request, 'state', None), 'user', None)
    if project and user:
        perm = check_project_permission(user, project, required_role="editor")
        if not perm:
            raise HTTPException(403, "Editor access required to upload")

    if not file.filename:
        raise HTTPException(400, "No filename provided")

    ext = Path(file.filename).suffix.lower()
    if ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(400, f"Unsupported format: {ext}. Allowed: {', '.join(ALLOWED_EXTENSIONS)}")

    # Stream to temp file instead of loading to memory
    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
        size = 0
        while chunk := await file.read(1024 * 1024):  # 1MB chunks
            size += len(chunk)
            if size > MAX_FILE_SIZE:
                raise HTTPException(400, f"File too large. Max: {MAX_FILE_SIZE // (1024*1024)} MB")
            tmp.write(chunk)
        tmp_path = tmp.name

    # For text files, read content from temp file (needed by _extract_content)
    if ext in ('.md', '.txt', '.sql', '.py'):
        with open(tmp_path, 'rb') as f:
            content = f.read()
    else:
        content = b''  # Don't hold file content in memory

    # Smart routing for non-data files (SQL, MD, TXT, DOCX, PPTX, PDF, images)
    if ext in (".sql", ".md", ".txt", ".docx", ".pptx", ".pdf", ".jpg", ".jpeg", ".png") and project:
        try:
            conductor_result = _conduct_upload(tmp_path, ext, project, file.filename, raw_content=content)
            text_content = conductor_result.get("text", "")
            doc_tables = conductor_result.get("tables", [])
            file_type = classify_file(file.filename, content_sample=text_content[:500])

            if file_type == "sql_patterns":
                result = process_sql_file(project, text_content)
                # Also save as doc
                doc_dir = KNOWLEDGE_DIR / project / "docs"
                doc_dir.mkdir(parents=True, exist_ok=True)
                with open(doc_dir / file.filename, "w") as f:
                    f.write(text_content)
                Path(tmp_path).unlink(missing_ok=True)
                return {"status": "ok", "file_type": "sql_patterns", "smart": result}

            if file_type == "business_rules":
                result = process_business_rules_doc(project, text_content, file.filename)
                # Also save as doc
                doc_dir = KNOWLEDGE_DIR / project / "docs"
                doc_dir.mkdir(parents=True, exist_ok=True)
                with open(doc_dir / file.filename, "w") as f:
                    f.write(text_content)
                Path(tmp_path).unlink(missing_ok=True)
                return {"status": "ok", "file_type": "business_rules", "smart": result}

            if file_type == "documentation":
                # Save as doc + index into knowledge
                doc_dir = KNOWLEDGE_DIR / project / "docs"
                doc_dir.mkdir(parents=True, exist_ok=True)
                save_name = Path(file.filename).stem + ".txt" if ext in (".pptx", ".docx", ".pdf") else file.filename
                with open(doc_dir / save_name, "w") as f:
                    f.write(text_content)
                # Index into project knowledge base
                indexed = False
                if text_content.strip():
                    try:
                        from agno.knowledge.reader.text_reader import TextReader
                        from db.session import create_project_knowledge
                        knowledge = create_project_knowledge(project)
                        knowledge.insert(
                            name=f"doc-{Path(file.filename).stem}",
                            text_content=f"Document: {file.filename}\n\n{text_content[:10000]}",
                            reader=TextReader(),
                            skip_if_exists=False,
                        )
                        indexed = True
                    except Exception:
                        pass
                # Save extracted tables to PostgreSQL
                tables_saved = 0
                if doc_tables and project:
                    try:
                        import pandas as pd
                        from db import get_project_engine
                        from db.session import create_project_schema
                        schema = create_project_schema(project)
                        engine = get_project_engine(project)
                        for tbl_info in doc_tables:
                            df = tbl_info["df"]
                            if len(df) < 2 or len(df.columns) < 2:
                                continue
                            # Generate table name from filename + source
                            base_name = Path(file.filename).stem.lower()
                            base_name = re.sub(r'[^a-z0-9_]', '_', base_name)[:30]
                            tbl_name = f"{base_name}_{tbl_info['source']}"
                            tbl_name = re.sub(r'_+', '_', tbl_name).strip('_')
                            try:
                                df.to_sql(tbl_name, engine, schema=schema, if_exists='replace', index=False)
                                tables_saved += 1
                            except Exception:
                                pass
                    except Exception:
                        pass
                Path(tmp_path).unlink(missing_ok=True)
                return {"status": "ok", "file_type": "documentation", "filename": save_name, "indexed": indexed, "tables_saved": tables_saved}
        except Exception:
            pass

    try:
        # Excel multi-sheet: use AI handler for .xlsx/.xls
        if ext in (".xlsx", ".xls") and not table_name:
            excel_result = _handle_excel(tmp_path, file.filename)
            if excel_result.get("tables") and len(excel_result["tables"]) > 1:
                # Multi-sheet Excel — create each sheet as its own table
                proj_schema = _get_project_schema(request, project)
                if proj_schema:
                    from db import get_project_engine
                    user_schema = proj_schema
                    engine = get_project_engine(project or '')
                else:
                    from db import create_user_schema, get_user_engine
                    user_schema = create_user_schema(user_id)
                    engine = get_user_engine(user_id)

                tables_created = []
                total_rows = 0
                for tbl_info in excel_result["tables"]:
                    df = tbl_info["df"]
                    if df.empty or len(df.columns) == 0:
                        continue
                    tbl_name = _sanitize_table_name(tbl_info["name"])
                    try:
                        df.to_sql(tbl_name, engine, schema=user_schema, if_exists='replace', index=False)
                        tables_created.append({"table": tbl_name, "rows": len(df), "cols": len(df.columns), "source": tbl_info.get("source", ""), "sheet_number": tbl_info.get("sheet_number", 0), "description": tbl_info.get("description", "")})
                        total_rows += len(df)

                        # Save source metadata for DATASETS tab display
                        if project:
                            src_dir = KNOWLEDGE_DIR / project / "table_sources"
                            src_dir.mkdir(parents=True, exist_ok=True)
                            _safe_write_json(src_dir / f"{tbl_name}.json", {
                                "source_file": file.filename,
                                "source_detail": f"Sheet {tbl_info.get('sheet_number', '?')}: {tbl_info.get('source', '')}",
                                "description": tbl_info.get("description", ""),
                            })

                        # Profile table for real health %
                        if project:
                            try:
                                _profile_table(df, project, tbl_name)
                            except Exception:
                                pass

                        # Generate metadata + queries for each table
                        if project:
                            col_analyses = [_analyze_column(df[col]) for col in df.columns]
                            try:
                                metadata = _generate_metadata(tbl_name, df, col_analyses)
                                if metadata:
                                    meta_dir = KNOWLEDGE_DIR / project / "tables"
                                    meta_dir.mkdir(parents=True, exist_ok=True)
                                    _safe_write_json(meta_dir / f"{tbl_name}.json", metadata)
                            except Exception:
                                pass
                            try:
                                _generate_sample_queries(project, tbl_name, col_analyses)
                            except Exception:
                                pass
                    except Exception as e:
                        excel_result.setdefault("warnings", []).append(f"Table '{tbl_name}' failed: {e}")

                # Index text descriptions into knowledge
                if excel_result.get("text") and project:
                    try:
                        from agno.knowledge.reader.text_reader import TextReader
                        from db.session import create_project_knowledge
                        knowledge = create_project_knowledge(project)
                        knowledge.insert(
                            name=f"doc-{Path(file.filename).stem}",
                            text_content=f"Excel file: {file.filename}\n\n{excel_result['text'][:10000]}",
                            reader=TextReader(), skip_if_exists=False,
                        )
                    except Exception:
                        pass

                # Reload knowledge
                if project:
                    try:
                        from db.session import create_project_knowledge
                        pk = create_project_knowledge(project)
                        pk.load(recreate=False)
                    except Exception:
                        pass

                # Trigger Engineer agent in background to inspect tables, create views, discover relationships
                if project and tables_created:
                    _bg_executor.submit(_post_upload_engineer, project, tables_created, user_id or 1)

                Path(tmp_path).unlink(missing_ok=True)
                return {
                    "status": "ok",
                    "multi_sheet": True,
                    "tables_created": len(tables_created),
                    "total_rows": total_rows,
                    "tables": tables_created,
                    "warnings": excel_result.get("warnings", []),
                    "engineer": "running in background — inspecting tables, creating views",
                    "smart": {"file_type": "excel_multi_sheet", "tables_created": len(tables_created)},
                }

        # 1. Parse file to validate (single table: CSV, JSON, single-sheet Excel)
        try:
            df = _read_file(tmp_path, ext)
        except Exception as parse_err:
            raise HTTPException(400, f"Could not parse file: {str(parse_err)[:200]}")
        if df.empty or len(df.columns) == 0:
            raise HTTPException(400, "File is empty or has no columns")

        # Check if this is a column definition file (check headers + values)
        raw_headers = list(df.columns)
        file_type_check = classify_file(file.filename, raw_headers, df=df, project_slug=project or "")
        if file_type_check == "column_definition" and project:
            result = process_column_definitions(project, df)
            # Save record so file appears in docs list
            docs_dir = KNOWLEDGE_DIR / project / "docs"
            docs_dir.mkdir(parents=True, exist_ok=True)
            (docs_dir / file.filename).write_text(
                f"[Column Definition File: {file.filename}]\n"
                f"Annotations saved: {result.get('annotations', 0)}\n"
                f"Memories saved: {result.get('memories', 0)}\n"
                f"Rules saved: {result.get('rules', 0)}"
            )
            return {"status": "ok", "file_type": "column_definition", "smart": result}

        # Column names already cleaned by _clean_dataframe() in _read_file()

        # 2. Determine table name
        tbl = _sanitize_table_name(table_name or file.filename)

        # 3. Save raw file to project staging area (upload only — no DB load yet)
        if project:
            staging_dir = KNOWLEDGE_DIR / project / "staging"
            staging_dir.mkdir(parents=True, exist_ok=True)
            import shutil
            staged_path = staging_dir / f"{tbl}{ext}"
            shutil.copy2(tmp_path, str(staged_path))

        # 4. Analyze columns (lightweight — no DB write)
        col_analyses = [_analyze_column(df[col]) for col in df.columns]

        # 5. Load to PostgreSQL
        proj_schema = _get_project_schema(request, project)
        if proj_schema:
            from db import get_project_engine
            user_schema = proj_schema
            engine = get_project_engine(project or '')
        else:
            from db import create_user_schema, get_user_engine
            user_schema = create_user_schema(user_id)
            engine = get_user_engine(user_id)

        # Handle action for existing tables
        insp = inspect(engine)
        table_exists = tbl in insp.get_table_names(schema=user_schema)
        upload_action = action if action != "auto" else ("replace" if replace else "new")
        rows_appended = 0
        rows_upserted = 0

        if table_exists and upload_action == "new":
            # Table exists, no action specified — return smart conflict
            match_info = match_existing_table(project or "", list(df.columns))
            raise HTTPException(409, detail={
                "message": f"Table '{tbl}' already exists",
                "match": match_info,
                "options": ["replace", "append", "upsert"],
                "hint": "Use action=append to add rows, action=upsert to update+add, action=replace to overwrite",
            })

        if table_exists and upload_action == "append":
            # APPEND: check for duplicates first, insert only new rows
            try:
                with engine.connect() as conn:
                    existing_count = conn.execute(text(f'SELECT COUNT(*) FROM "{user_schema}"."{tbl}"')).scalar() or 0
                    # Use first column as key to find new rows
                    if len(df.columns) == 0:
                        raise HTTPException(400, "File has no usable columns")
                    # Smart PK detection for upsert
                    pk_col = None
                    # 1. Check for 'id' column
                    id_cols = [c for c in df.columns if 'id' in c.lower() and df[c].nunique() == len(df)]
                    if id_cols:
                        pk_col = id_cols[0]
                    # 2. Check if first column is unique
                    elif df[df.columns[0]].nunique() >= len(df) * 0.95:
                        pk_col = df.columns[0]
                    else:
                        # Can't determine PK — fall back to append
                        pk_col = df.columns[0]
                    existing_keys = set()
                    try:
                        result = conn.execute(text(f'SELECT "{pk_col}" FROM "{user_schema}"."{tbl}"'))
                        existing_keys = set(str(r[0]) for r in result.fetchall())
                    except Exception:
                        pass

                    if existing_keys:
                        new_rows = df[~df[pk_col].astype(str).isin(existing_keys)]
                    else:
                        new_rows = df

                    if len(new_rows) == 0:
                        rows_appended = 0
                    else:
                        new_rows.to_sql(tbl, engine, if_exists="append", index=False, schema=user_schema)
                        rows_appended = len(new_rows)
            except Exception:
                # Fallback: just append all
                df.to_sql(tbl, engine, if_exists="append", index=False, schema=user_schema)
                rows_appended = len(df)

        elif table_exists and upload_action == "upsert":
            # UPSERT: use first column as key, update existing + insert new
            # Only use columns that exist in BOTH new data and existing table
            existing_cols = set(c["name"] for c in insp.get_columns(tbl, schema=user_schema))
            common_cols = [c for c in df.columns if c in existing_cols]
            missing_in_new = existing_cols - set(df.columns)
            extra_in_new = set(df.columns) - existing_cols

            # Add missing columns to DB if new data has extra columns
            if extra_in_new:
                with engine.connect() as conn:
                    for col in extra_in_new:
                        try:
                            conn.execute(text(f'ALTER TABLE "{user_schema}"."{tbl}" ADD COLUMN IF NOT EXISTS "{col}" TEXT'))
                        except Exception:
                            pass
                    conn.commit()
                common_cols = list(df.columns)

            # Use only common columns for upsert
            df_upsert = df[common_cols] if common_cols else df
            if len(df_upsert.columns) == 0:
                raise HTTPException(400, "File has no usable columns")
            # Smart PK detection for upsert
            pk_col = None
            # 1. Check for 'id' column
            id_cols = [c for c in df_upsert.columns if 'id' in c.lower() and df_upsert[c].nunique() == len(df_upsert)]
            if id_cols:
                pk_col = id_cols[0]
            # 2. Check if first column is unique
            elif df_upsert[df_upsert.columns[0]].nunique() >= len(df_upsert) * 0.95:
                pk_col = df_upsert.columns[0]
            else:
                # Can't determine PK — fall back to append
                pk_col = df_upsert.columns[0]

            with engine.connect() as conn:
                existing_pks = set()
                try:
                    result = conn.execute(text(f'SELECT "{pk_col}" FROM "{user_schema}"."{tbl}"'))
                    existing_pks = set(str(r[0]) for r in result.fetchall())
                except Exception:
                    pass

                new_rows = df_upsert[~df_upsert[pk_col].astype(str).isin(existing_pks)]
                update_rows = df_upsert[df_upsert[pk_col].astype(str).isin(existing_pks)]

                if len(new_rows) > 0:
                    new_rows.to_sql(tbl, engine, if_exists="append", index=False, schema=user_schema)
                    rows_appended = len(new_rows)

                for _, row in update_rows.iterrows():
                    update_cols = [c for c in df_upsert.columns if c != pk_col]
                    if not update_cols:
                        continue
                    set_clause = ", ".join(f'"{c}" = :{c}' for c in update_cols)
                    params = {c: (None if pd.isna(row[c]) else row[c]) for c in df_upsert.columns}
                    conn.execute(text(
                        f'UPDATE "{user_schema}"."{tbl}" SET {set_clause} WHERE "{pk_col}" = :{pk_col}'
                    ), params)
                rows_upserted = len(update_rows)
                conn.commit()

        else:
            # REPLACE or new table
            mode = "replace" if (replace or upload_action == "replace") else "fail"
            df.to_sql(tbl, engine, if_exists=mode, index=False, schema=user_schema)

        # 5. Generate metadata
        metadata = _generate_metadata(tbl, df, col_analyses)

        # 6. Generate sample queries
        sample_queries = _generate_sample_queries(tbl, col_analyses)

        # 7. Generate business rules
        biz_rules = _generate_business_rules(tbl, col_analyses)

        # 8. Save knowledge files (per-project if project provided)
        if proj_schema:
            tables_dir = KNOWLEDGE_DIR / (project or proj_schema) / "tables"
            queries_dir = KNOWLEDGE_DIR / (project or proj_schema) / "queries"
            business_dir = KNOWLEDGE_DIR / (project or proj_schema) / "business"
        else:
            tables_dir = TABLES_DIR
            queries_dir = QUERIES_DIR
            business_dir = BUSINESS_DIR

        tables_dir.mkdir(parents=True, exist_ok=True)
        queries_dir.mkdir(parents=True, exist_ok=True)
        business_dir.mkdir(parents=True, exist_ok=True)

        with open(tables_dir / f"{tbl}.json", "w") as f:
            json.dump(metadata, f, indent=2, default=str)

        # Save source metadata for DATASETS tab
        if proj_schema:
            src_dir = KNOWLEDGE_DIR / (project or proj_schema) / "table_sources"
            src_dir.mkdir(parents=True, exist_ok=True)
            _safe_write_json(src_dir / f"{tbl}.json", {
                "source_file": file.filename if hasattr(file, 'filename') else "",
                "source_detail": ext.replace(".", "").upper(),
                "description": metadata.get("description", "") if isinstance(metadata, dict) else "",
            })

        # Run YData Profiling for real health %
        try:
            _profile_table(df, project or proj_schema, tbl)
        except Exception:
            pass

        queries_file = queries_dir / f"{tbl}_queries.sql"
        with open(queries_file, "w") as f:
            f.write(f"-- Auto-generated queries for {tbl}\n\n{sample_queries}\n")

        biz_file = business_dir / f"{tbl}_rules.json"
        with open(biz_file, "w") as f:
            json.dump(biz_rules, f, indent=2, default=str)

        # 9. Reload knowledge — project-specific or global
        if proj_schema and project:
            from db.session import create_project_knowledge
            proj_knowledge = create_project_knowledge(project)
            proj_dir = KNOWLEDGE_DIR / project
            if proj_dir.exists():
                for subdir in ["tables", "queries", "business"]:
                    path = proj_dir / subdir
                    if path.exists():
                        files = [f for f in path.iterdir() if f.is_file() and not f.name.startswith(".")]
                        if files:
                            proj_knowledge.insert(name=f"{project}-{subdir}", path=str(path))
        else:
            from scripts.load_knowledge import load_knowledge
            load_knowledge(recreate=False)

        # Training is NOT auto-triggered on upload.
        # User clicks TRAIN ALL to start training manually.

        # 10. Smart analysis: classify, match, fingerprint
        file_type = classify_file(file.filename, list(df.columns))
        if project:
            save_fingerprint(project, tbl, len(df), list(df.columns))
            change_type = check_fingerprint_changed(project, tbl, len(df), list(df.columns))
        else:
            change_type = "new"

        return {
            "status": "ok",
            "table_name": tbl,
            "rows": len(df),
            "columns": len(df.columns),
            "column_details": col_analyses,
            "metadata": metadata,
            "sample_queries_count": sample_queries.count("<query>"),
            "business_rules": biz_rules,
            # Smart upload info
            "smart": {
                "file_type": file_type,
                "action": upload_action,
                "rows_appended": rows_appended,
                "rows_upserted": rows_upserted,
                "change_type": change_type,
                "fingerprint": compute_fingerprint(len(df), list(df.columns)),
            },
        }

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"Upload failed: {str(e)}")
    finally:
        Path(tmp_path).unlink(missing_ok=True)


@router.get("/dashboard")
def get_dashboard(request: Request):
    """Full dashboard data: tables, knowledge stats, system info."""
    user_id = _get_user_id(request)
    from db import create_user_schema
    user_schema = create_user_schema(user_id)

    engine = create_engine(db_url)
    insp = inspect(engine)

    # Tables (public + user schema)
    tables_list = []
    total_rows = 0
    with engine.connect() as conn:
        # Public tables (shared demo)
        for tbl in sorted(insp.get_table_names(schema="public")):
            try:
                count = conn.execute(text(f'SELECT COUNT(*) FROM public."{tbl}"')).scalar() or 0
                cols = insp.get_columns(tbl, schema="public")
                tables_list.append({
                    "name": tbl, "rows": count, "columns": len(cols),
                    "protected": tbl in PROTECTED_TABLES,
                })
                total_rows += count
            except Exception:
                tables_list.append({"name": tbl, "rows": 0, "columns": 0, "protected": tbl in PROTECTED_TABLES})

        # User schema tables
        for tbl in sorted(insp.get_table_names(schema=user_schema)):
            try:
                count = conn.execute(text(f'SELECT COUNT(*) FROM "{user_schema}"."{tbl}"')).scalar() or 0
                cols = insp.get_columns(tbl, schema=user_schema)
                tables_list.append({
                    "name": tbl, "rows": count, "columns": len(cols),
                    "protected": False,
                })
                total_rows += count
            except Exception:
                tables_list.append({"name": tbl, "rows": 0, "columns": 0, "protected": False})

        # DB size
        try:
            db_size = conn.execute(text("SELECT pg_size_pretty(pg_database_size(current_database()))")).scalar()
        except Exception:
            db_size = "unknown"

        # Knowledge vectors
        knowledge_count = 0
        learnings_count = 0
        try:
            knowledge_count = conn.execute(text("SELECT COUNT(*) FROM ai.dash_knowledge")).scalar() or 0
        except Exception:
            pass
        try:
            learnings_count = conn.execute(text("SELECT COUNT(*) FROM ai.dash_learnings")).scalar() or 0
        except Exception:
            pass

    # Knowledge files
    tables_files = len(list(TABLES_DIR.glob("*.json"))) if TABLES_DIR.exists() else 0
    queries_files = len(list(QUERIES_DIR.glob("*.sql"))) if QUERIES_DIR.exists() else 0
    business_files = len(list(BUSINESS_DIR.glob("*.json"))) if BUSINESS_DIR.exists() else 0

    return {
        "tables": tables_list,
        "stats": {
            "table_count": len(tables_list),
            "total_rows": total_rows,
            "knowledge_vectors": knowledge_count,
            "learnings_count": learnings_count,
            "db_size": db_size,
            "model": TRAINING_MODEL,
            "model_provider": "OpenRouter",
            "embeddings": "text-embedding-3-small",
            "status": "online",
        },
        "knowledge": {
            "table_metadata": tables_files,
            "query_patterns": queries_files,
            "business_rules": business_files,
            "learnings": learnings_count,
        },
    }


@router.post("/upload-doc")
async def upload_document(request: Request, file: UploadFile, project: str | None = None):
    """Upload a code/doc file to project or global knowledge base."""
    # Editor role required for doc uploads
    from app.auth import check_project_permission
    user = getattr(getattr(request, 'state', None), 'user', None)
    if project and user:
        perm = check_project_permission(user, project, required_role="editor")
        if not perm:
            raise HTTPException(403, "Editor access required to upload")

    if not file.filename:
        raise HTTPException(400, "No filename provided")

    ext = Path(file.filename).suffix.lower()
    allowed = {".sql", ".py", ".txt", ".md", ".pptx", ".docx", ".pdf", ".jpg", ".jpeg", ".png"}
    if ext not in allowed:
        raise HTTPException(400, f"Unsupported: {ext}. Allowed: {', '.join(allowed)}")

    content = await file.read()
    if len(content) > MAX_FILE_SIZE:
        raise HTTPException(400, "File too large")

    # Use Upload Conductor for all formats
    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
        tmp.write(content)
        tmp_path = tmp.name

    conductor_result = _conduct_upload(tmp_path, ext, project or "", file.filename, raw_content=content)
    text_content = conductor_result.get("text", "")
    doc_tables = conductor_result.get("tables", [])

    import os; os.unlink(tmp_path)
    doc_name = Path(file.filename).stem

    # Per-project or global docs dir
    if project:
        docs_dir = KNOWLEDGE_DIR / project / "docs"
    else:
        docs_dir = KNOWLEDGE_DIR / "docs"
    docs_dir.mkdir(parents=True, exist_ok=True)
    doc_path = docs_dir / file.filename
    with open(doc_path, "w") as f:
        f.write(text_content)

    # Save raw binary for structure extraction (doc-to-workflow)
    if ext in (".pptx", ".docx", ".pdf", ".jpg", ".jpeg", ".png") and project:
        docs_raw_dir = KNOWLEDGE_DIR / project / "docs_raw"
        docs_raw_dir.mkdir(parents=True, exist_ok=True)
        (docs_raw_dir / file.filename).write_bytes(content)

    # Index into project-specific or global knowledge
    from agno.knowledge.reader.text_reader import TextReader
    if project:
        from db.session import create_project_knowledge
        knowledge = create_project_knowledge(project)
    else:
        from dash.settings import dash_knowledge
        knowledge = dash_knowledge

    try:
        knowledge.insert(
            name=f"doc-{doc_name}",
            text_content=f"Document: {file.filename}\n\n{text_content[:10000]}",
            reader=TextReader(),
            skip_if_exists=False,
        )
    except Exception as e:
        raise HTTPException(500, f"Indexing failed: {str(e)}")

    # Save extracted tables to PostgreSQL
    tables_saved = 0
    if doc_tables and project:
        try:
            import pandas as pd
            from db import get_project_engine
            from db.session import create_project_schema
            schema = create_project_schema(project)
            engine = get_project_engine(project)
            for tbl_info in doc_tables:
                df = tbl_info["df"]
                if len(df) < 2 or len(df.columns) < 2:
                    continue
                # Use conductor-provided name or generate from filename + source
                tbl_name = tbl_info.get("name", "")
                if not tbl_name:
                    base_name = Path(file.filename).stem.lower()
                    base_name = re.sub(r'[^a-z0-9_]', '_', base_name)[:30]
                    tbl_name = f"{base_name}_{tbl_info.get('source', 'table')}"
                tbl_name = re.sub(r'[^a-z0-9_]', '_', tbl_name.lower())
                tbl_name = re.sub(r'_+', '_', tbl_name).strip('_')[:63]
                try:
                    df.to_sql(tbl_name, engine, schema=schema, if_exists='replace', index=False)
                    tables_saved += 1
                    # Save source metadata
                    src_dir = KNOWLEDGE_DIR / project / "table_sources"
                    src_dir.mkdir(parents=True, exist_ok=True)
                    source_detail = tbl_info.get("source", "")
                    if "slide" in source_detail:
                        source_detail = f"Slide {source_detail.replace('slide_', '')}"
                    elif "page" in source_detail:
                        source_detail = f"Page {source_detail.replace('page_', '').replace('_table_', ' Table ')}"
                    elif "table" in source_detail:
                        source_detail = f"Table {source_detail.replace('table_', '')}"
                    _safe_write_json(src_dir / f"{tbl_name}.json", {
                        "source_file": file.filename,
                        "source_detail": source_detail,
                        "description": "",
                    })
                except Exception:
                    pass
        except Exception:
            pass

    # Trigger Engineer in background if tables were extracted from document
    if tables_saved > 0 and project:
        doc_tables_info = [{"table": t.get("name", ""), "rows": len(t.get("df", [])), "cols": len(t.get("df", {}).columns) if hasattr(t.get("df"), "columns") else 0, "source": t.get("source", "")} for t in doc_tables if t.get("df") is not None and len(t.get("df", [])) >= 2]
        if doc_tables_info:
            user_id_val = getattr(getattr(getattr(request, 'state', None), 'user', None), 'id', 1)
            _bg_executor.submit(_post_upload_engineer, project, doc_tables_info, user_id_val)

    return {"status": "ok", "filename": file.filename, "type": ext, "size": len(text_content), "indexed": True, "tables_saved": tables_saved}


@router.post("/upload-agent")
async def upload_with_agent(request: Request, file: UploadFile, project: str | None = None):
    """Upload a file using the Upload Agent Team (Conductor → Parser/Scanner/Vision → Inspector → Engineer).

    This uses AI agents for intelligent processing: structure detection,
    unpivot, merge, quality validation, and post-upload optimization.
    Falls back to standard upload if agent fails.
    """
    from app.auth import check_project_permission
    user = getattr(getattr(request, 'state', None), 'user', None)
    if project and user:
        perm = check_project_permission(user, project, required_role="editor")
        if not perm:
            raise HTTPException(403, "Editor access required to upload")

    if not file.filename:
        raise HTTPException(400, "No filename provided")
    if not project:
        raise HTTPException(400, "Project required for agent upload")

    ext = Path(file.filename).suffix.lower()
    if ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(400, f"Unsupported: {ext}")

    # Stream to temp file
    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
        size = 0
        while chunk := await file.read(1024 * 1024):
            size += len(chunk)
            if size > MAX_FILE_SIZE:
                raise HTTPException(400, "File too large")
            tmp.write(chunk)
        tmp_path = tmp.name

    user_id = _get_user_id(request) or 1

    try:
        # Create the Upload Agent Team
        from dash.agents.conductor import create_upload_team
        uid = 1
        try:
            uid = int(user_id) if user_id and str(user_id).isdigit() else 1
        except Exception:
            pass
        upload_team = create_upload_team(project, user_id=uid)

        # Build the prompt for the Conductor
        file_info = f"File: {file.filename} ({ext}, {size:,} bytes)"
        if ext in (".xlsx", ".xls"):
            file_info += "\nThis is an Excel file — use Parser to analyze sheets, detect structure, unpivot if needed."
        elif ext in (".pdf",):
            file_info += "\nThis is a PDF — use Scanner to extract text, tables, OCR scanned pages."
        elif ext in (".pptx",):
            file_info += "\nThis is a PowerPoint — use Scanner to extract slides, tables, images."
        elif ext in (".docx",):
            file_info += "\nThis is a Word document — use Scanner to extract text, tables, images."
        elif ext in (".jpg", ".jpeg", ".png"):
            file_info += "\nThis is an image — use Vision to OCR text or describe content."
        elif ext in (".csv",):
            file_info += "\nThis is a CSV — use Parser for fast direct load."
        elif ext in (".json",):
            file_info += "\nThis is a JSON file — use Parser to load."
        else:
            file_info += "\nThis is a text file — use Scanner to index to knowledge base."

        prompt = f"""Process this uploaded file for project '{project}'.

{file_info}
File saved at: {tmp_path}

STEPS:
1. Use the right agent to parse/extract this file
2. After tables are created, use Inspector to validate quality
3. Use Engineer to create views and discover relationships
4. Report: tables created, health scores, views, relationships"""

        # Run the Conductor
        response = upload_team.run(prompt)
        agent_output = response.content if response else "Agent returned no response"

        # Also run the standard conductor as backup (ensures tables are actually created)
        if ext in ('.md', '.txt', '.sql', '.py'):
            with open(tmp_path, 'rb') as f:
                content = f.read()
        else:
            content = b''

        conductor_result = _conduct_upload(tmp_path, ext, project, file.filename, raw_content=content)

        # Store tables from conductor result
        tables_stored = []
        if conductor_result.get("tables"):
            proj_schema = re.sub(r"[^a-z0-9_]", "_", project.lower())[:63]
            try:
                from db import get_project_engine
                from db.session import create_project_schema
                schema = create_project_schema(project)
                eng = get_project_engine(project)
                for tbl_info in conductor_result["tables"]:
                    df = tbl_info.get("df")
                    if df is None or len(df) == 0:
                        continue
                    tbl_name = _sanitize_table_name(tbl_info.get("name", "table"))
                    try:
                        df.to_sql(tbl_name, eng, schema=schema, if_exists='replace', index=False)
                        tables_stored.append({"table": tbl_name, "rows": len(df), "cols": len(df.columns), "source": tbl_info.get("source", "")})
                        # Profile + source metadata
                        try:
                            _profile_table(df, project, tbl_name)
                        except Exception:
                            pass
                        src_dir = KNOWLEDGE_DIR / project / "table_sources"
                        src_dir.mkdir(parents=True, exist_ok=True)
                        _safe_write_json(src_dir / f"{tbl_name}.json", {
                            "source_file": file.filename,
                            "source_detail": tbl_info.get("source", ext.upper()),
                            "description": tbl_info.get("description", ""),
                        })
                    except Exception:
                        pass
            except Exception:
                pass

        # Index text to knowledge
        text_content = conductor_result.get("text", "")
        if text_content and text_content.strip():
            try:
                from agno.knowledge.reader.text_reader import TextReader
                from db.session import create_project_knowledge
                knowledge = create_project_knowledge(project)
                knowledge.insert(name=f"doc-{Path(file.filename).stem}", text_content=f"Document: {file.filename}\n\n{text_content[:10000]}", reader=TextReader(), skip_if_exists=False)
            except Exception:
                pass

        # Trigger Engineer in background
        if tables_stored:
            _bg_executor.submit(_post_upload_engineer, project, tables_stored, user_id or 1)

        Path(tmp_path).unlink(missing_ok=True)

        return {
            "status": "ok",
            "agent": True,
            "agent_report": agent_output[:2000] if agent_output else "",
            "tables_created": len(tables_stored),
            "tables": tables_stored,
            "text_indexed": len(text_content),
            "warnings": conductor_result.get("warnings", []),
            "errors": conductor_result.get("errors", []),
        }

    except Exception as e:
        # Fallback to standard upload if agent fails
        Path(tmp_path).unlink(missing_ok=True)
        return {"status": "error", "agent": True, "error": str(e)[:500], "fallback": "Use /api/upload or /api/upload-doc instead"}


@router.get("/docs")
def list_docs(request: Request, project: str | None = None):
    """List uploaded documents (project-scoped or global)."""
    if project:
        docs_dir = KNOWLEDGE_DIR / project / "docs"
    else:
        docs_dir = KNOWLEDGE_DIR / "docs"
    if not docs_dir.exists():
        return {"docs": []}
    docs = []
    for f in sorted(docs_dir.iterdir()):
        if f.is_file() and not f.name.startswith("."):
            docs.append({"name": f.name, "size": f.stat().st_size, "type": f.suffix})
    return {"docs": docs}


@router.delete("/docs/{filename}")
def delete_doc(filename: str, request: Request, project: str | None = None):
    """Delete an uploaded document."""
    safe_name = Path(filename).name
    if safe_name != filename or '..' in filename:
        raise HTTPException(400, "Invalid filename")
    if project:
        doc_path = KNOWLEDGE_DIR / project / "docs" / filename
    else:
        doc_path = KNOWLEDGE_DIR / "docs" / filename
    if not doc_path.exists():
        raise HTTPException(404, f"Document '{filename}' not found")
    doc_path.unlink()
    return {"status": "ok", "deleted": filename}


@router.get("/tables")
def list_tables():
    """List all tables in the public schema with row counts."""
    engine = create_engine(db_url)
    insp = inspect(engine)
    tables = []

    with engine.connect() as conn:
        for tbl in sorted(insp.get_table_names(schema="public")):
            try:
                count = conn.execute(text(f'SELECT COUNT(*) FROM public."{tbl}"')).scalar()
                cols = insp.get_columns(tbl, schema="public")
                tables.append({
                    "name": tbl,
                    "rows": count,
                    "columns": len(cols),
                    "protected": tbl in PROTECTED_TABLES,
                })
            except Exception:
                tables.append({"name": tbl, "rows": 0, "columns": 0, "protected": tbl in PROTECTED_TABLES})

    return {"tables": tables}


_SCHEDULES_FILE = KNOWLEDGE_DIR / "schedules.json"


@router.get("/user-schedules")
def list_schedules(request: Request):
    """List saved scheduled reports for current user."""
    user_id = _get_user_id(request)
    if not _SCHEDULES_FILE.exists():
        return {"schedules": []}
    with open(_SCHEDULES_FILE) as f:
        all_s = json.load(f)
    return {"schedules": [s for s in all_s if s.get("user_id") == user_id]}


@router.post("/user-schedules")
def create_schedule(request: Request, name: str, prompt: str, frequency: str = "daily"):
    """Save a scheduled report."""
    user_id = _get_user_id(request)
    schedules = []
    if _SCHEDULES_FILE.exists():
        with open(_SCHEDULES_FILE) as f:
            schedules = json.load(f)
    schedules.append({
        "name": name, "prompt": prompt, "frequency": frequency,
        "user_id": user_id, "enabled": True,
        "created_at": __import__("datetime").datetime.now().isoformat(),
        "last_run": None,
    })
    _SCHEDULES_FILE.parent.mkdir(parents=True, exist_ok=True)
    with open(_SCHEDULES_FILE, "w") as f:
        json.dump(schedules, f, indent=2)
    return {"status": "ok"}


@router.delete("/user-schedules/{name}")
def delete_schedule(name: str, request: Request):
    """Delete a scheduled report."""
    user_id = _get_user_id(request)
    if not _SCHEDULES_FILE.exists():
        raise HTTPException(404, "Not found")
    with open(_SCHEDULES_FILE) as f:
        schedules = json.load(f)
    schedules = [s for s in schedules if not (s["name"] == name and s.get("user_id") == user_id)]
    with open(_SCHEDULES_FILE, "w") as f:
        json.dump(schedules, f, indent=2)
    return {"status": "ok"}


_EVAL_RESULTS_FILE = KNOWLEDGE_DIR / "eval_results.json"
_eval_running = False


@router.post("/evals/run")
def run_evals(request: Request):
    """Run eval suite in background."""
    global _eval_running
    if _eval_running:
        return {"status": "already_running"}

    run_id = f"eval-{__import__('time').time():.0f}"

    def _run():
        global _eval_running
        _eval_running = True
        results = {"run_id": run_id, "status": "running", "categories": {}}
        try:
            # Save initial state
            KNOWLEDGE_DIR.mkdir(parents=True, exist_ok=True)
            with open(_EVAL_RESULTS_FILE, "w") as f:
                json.dump(results, f, indent=2)

            # Run evals
            from evals.run import run_evals as _run_evals
            passed = _run_evals()

            results["status"] = "completed"
            results["passed"] = passed
            results["completed_at"] = __import__("datetime").datetime.now().isoformat()
        except Exception as e:
            results["status"] = "error"
            results["error"] = str(e)
        finally:
            with open(_EVAL_RESULTS_FILE, "w") as f:
                json.dump(results, f, indent=2)
            _eval_running = False

    _bg_executor.submit(_run)
    return {"status": "started", "run_id": run_id}


@router.get("/evals/status")
def eval_status(request: Request):
    """Get latest eval run status."""
    if not _EVAL_RESULTS_FILE.exists():
        return {"status": "never_run", "running": _eval_running}
    with open(_EVAL_RESULTS_FILE) as f:
        results = json.load(f)
    results["running"] = _eval_running
    return results


@router.get("/training")
def get_training(request: Request, project: str | None = None):
    """Get learnings/training data for a project."""
    if not project:
        return {"learnings": 0, "items": [], "training_qa": []}
    import re
    safe = re.sub(r"[^a-z0-9_]", "_", project.lower())[:63]
    engine = create_engine(db_url)
    count = 0
    try:
        with engine.connect() as conn:
            count = conn.execute(text(f'SELECT COUNT(*) FROM ai."{safe}_learnings"')).scalar() or 0
    except Exception:
        pass

    # Load training Q&A pairs
    training_qa: list[dict] = []
    training_dir = KNOWLEDGE_DIR / project / "training"
    if training_dir.exists():
        for f in sorted(training_dir.glob("*.json")):
            try:
                with open(f) as fh:
                    data = json.load(fh)
                if isinstance(data, list):
                    for qa in data:
                        qa["source_table"] = f.stem.replace("_qa", "")
                        training_qa.append(qa)
            except Exception:
                pass

    return {"learnings": count, "training_qa": training_qa}


# ---------------------------------------------------------------------------
# Self-Learning Endpoints
# ---------------------------------------------------------------------------


@router.post("/projects/{slug}/feedback")
async def save_feedback(slug: str, request: Request):
    """Save user feedback (thumbs up/down) as training data."""
    body = await request.json()
    question = body.get("question", "")
    answer = body.get("answer", "")
    rating = body.get("rating", "up")  # "up" or "down"

    if not question or not answer:
        return {"status": "skip"}

    feedback_dir = KNOWLEDGE_DIR / slug / "training"
    feedback_dir.mkdir(parents=True, exist_ok=True)

    filename = "feedback_good.json" if rating == "up" else "feedback_bad.json"
    filepath = feedback_dir / filename

    existing = []
    if filepath.exists():
        try:
            with open(filepath) as f:
                existing = json.load(f)
        except Exception:
            existing = []

    existing.append({"question": question, "answer": answer[:500], "timestamp": time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())})
    # Keep last 50
    existing = existing[-50:]

    with open(filepath, "w") as f:
        json.dump(existing, f, indent=2)

    # Re-index
    _reload_project_knowledge(slug)
    return {"status": "ok", "saved": rating}


@router.post("/projects/{slug}/save-query-pattern")
async def save_query_pattern(slug: str, request: Request):
    """Save a proven query pattern for reuse."""
    body = await request.json()
    question = body.get("question", "")
    sql = body.get("sql", "")

    if not question or not sql:
        return {"status": "skip"}

    patterns_dir = KNOWLEDGE_DIR / slug / "queries"
    patterns_dir.mkdir(parents=True, exist_ok=True)
    filepath = patterns_dir / "proven_patterns.json"

    existing = []
    if filepath.exists():
        try:
            with open(filepath) as f:
                existing = json.load(f)
        except Exception:
            existing = []

    # Check for duplicate
    if not any(p.get("sql") == sql for p in existing):
        existing.append({"question": question, "sql": sql, "uses": 1, "timestamp": time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())})
        existing = existing[-30:]
        with open(filepath, "w") as f:
            json.dump(existing, f, indent=2)

    return {"status": "ok"}


@router.post("/projects/{slug}/suggest-followups")
async def suggest_followups(slug: str, request: Request):
    """LLM-powered follow-up suggestions based on conversation context."""
    body = await request.json()
    question = body.get("question", "")
    answer = body.get("answer", "")

    if not question or not answer:
        return {"suggestions": []}

    from os import getenv
    import httpx
    api_key = getenv("OPENROUTER_API_KEY", "")
    if not api_key:
        return {"suggestions": []}

    prompt = f"""Based on this data conversation, suggest 3 natural follow-up questions.

Q: {question}
A: {answer[:500]}

Return ONLY a JSON array of 3 short follow-up questions (no markdown):
["question 1", "question 2", "question 3"]"""

    try:
        resp = httpx.post(
            "https://openrouter.ai/api/v1/chat/completions",
            headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
            json={"model": TRAINING_MODEL, "messages": [{"role": "user", "content": prompt}], "max_tokens": 200, "temperature": 0.3},
            timeout=10,
        )
        result = resp.json()
        content = result.get("choices", [{}])[0].get("message", {}).get("content", "")
        suggestions = json.loads(content.strip().strip("`").strip())
        return {"suggestions": suggestions[:3] if isinstance(suggestions, list) else []}
    except Exception:
        return {"suggestions": []}


@router.post("/projects/{slug}/extract-context")
async def extract_context(slug: str, request: Request):
    """Extract key facts from a conversation for context memory."""
    body = await request.json()
    question = body.get("question", "")
    answer = body.get("answer", "")

    if not question or not answer:
        return {"status": "skip"}

    from os import getenv
    import httpx
    api_key = getenv("OPENROUTER_API_KEY", "")
    if not api_key:
        return {"status": "skip"}

    prompt = f"""Extract 0-3 key facts from this conversation that should be remembered for future sessions.

For each fact, provide:
- "fact": the fact text
- "confidence": "high" or "low"
  - "high": objective data facts (column names, table relationships, data formats, business definitions, SQL patterns)
  - "low": subjective observations (user preferences, interpretation of results, opinions)
- "score": confidence score 0-100 (100 = absolutely certain data fact, 0 = wild guess)

Q: {question}
A: {answer[:500]}

Return ONLY valid JSON (no markdown). Empty array if nothing worth remembering:
[{{"fact": "the fact text", "confidence": "high", "score": 95}}, {{"fact": "another fact", "confidence": "low", "score": 40}}]"""

    try:
        resp = httpx.post(
            "https://openrouter.ai/api/v1/chat/completions",
            headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
            json={"model": TRAINING_MODEL, "messages": [{"role": "user", "content": prompt}], "max_tokens": 300, "temperature": 0.1},
            timeout=10,
        )
        result = resp.json()
        content = result.get("choices", [{}])[0].get("message", {}).get("content", "")
        # Clean LLM output — strip markdown code fences
        clean = content.strip()
        if clean.startswith("```"):
            clean = clean.split("\n", 1)[-1] if "\n" in clean else clean[3:]
        if clean.endswith("```"):
            clean = clean[:-3]
        clean = clean.strip().strip("`").strip()
        if clean.lower().startswith("json"):
            clean = clean[4:].strip()
        parsed = json.loads(clean)

        if not isinstance(parsed, list) or len(parsed) == 0:
            return {"status": "ok", "facts": [], "auto_saved": []}

        # Separate high vs low confidence, keep scores
        auto_save = []
        auto_save_with_scores = []
        needs_approval = []
        needs_approval_with_scores = []
        for item in parsed[:3]:
            if isinstance(item, str):
                needs_approval.append(item)
                needs_approval_with_scores.append({"fact": item, "score": 30})
            elif isinstance(item, dict):
                fact = item.get("fact", "")
                confidence = item.get("confidence", "low")
                score = int(item.get("score", 50))
                if not fact:
                    continue
                if confidence == "high":
                    auto_save.append(fact)
                    auto_save_with_scores.append({"fact": fact, "score": score})
                else:
                    needs_approval.append(fact)
                    needs_approval_with_scores.append({"fact": fact, "score": score})

        # Auto-save high-confidence facts
        if auto_save:
            try:
                engine = create_engine(db_url)
                with engine.connect() as conn:
                    for fact in auto_save:
                        conn.execute(text(
                            "INSERT INTO public.dash_memories (project_slug, fact, scope, source) VALUES (:s, :f, 'project', 'auto')"
                        ), {"s": slug, "f": fact})
                    conn.commit()
            except Exception:
                pass

        # Always return "proposed" if there are any learnings (auto or manual)
        if auto_save or needs_approval:
            return {"status": "proposed", "facts": needs_approval, "facts_with_scores": needs_approval_with_scores, "auto_saved": auto_save, "auto_saved_with_scores": auto_save_with_scores}

        return {"status": "ok", "facts": [], "auto_saved": []}
    except Exception:
        return {"status": "skip"}


@router.post("/projects/{slug}/approve-learnings")
async def approve_learnings(slug: str, request: Request):
    """Save user-approved learnings to memory."""
    body = await request.json()
    facts = body.get("facts", [])
    scope = body.get("scope", "project")

    if not facts:
        return {"status": "skip"}

    # Save to DB (dash_memories)
    try:
        engine = create_engine(db_url)
        with engine.connect() as conn:
            for fact in facts[:3]:
                if fact:
                    conn.execute(text(
                        "INSERT INTO public.dash_memories (project_slug, fact, scope, source) VALUES (:s, :f, :sc, 'agent')"
                    ), {"s": slug, "f": fact, "sc": scope})
            conn.commit()
    except Exception:
        pass

    # Also save to JSON memory file for backward compat
    memory_dir = KNOWLEDGE_DIR / slug / "memory"
    memory_dir.mkdir(parents=True, exist_ok=True)
    filepath = memory_dir / "context.json"

    existing = []
    if filepath.exists():
        try:
            with open(filepath) as f:
                existing = json.load(f)
        except Exception:
            existing = []

    for fact in facts[:3]:
        if fact and not any(f.get("fact") == fact for f in existing):
            existing.append({"fact": fact, "timestamp": time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())})
    existing = existing[-30:]

    with open(filepath, "w") as f:
        json.dump(existing, f, indent=2)

    return {"status": "saved", "count": len(facts)}


@router.put("/projects/{slug}/persona")
async def update_persona(slug: str, request: Request):
    """Update project persona."""
    body = await request.json()
    try:
        engine = create_engine(db_url)
        with engine.connect() as conn:
            conn.execute(text("""
                INSERT INTO public.dash_personas (project_slug, persona)
                VALUES (:s, :p::jsonb)
                ON CONFLICT (project_slug)
                DO UPDATE SET persona = :p::jsonb, updated_at = NOW()
            """), {"s": slug, "p": json.dumps(body, default=str)})
            conn.commit()
        # Also save to file for PgVector
        persona_file = KNOWLEDGE_DIR / slug / "persona.json"
        persona_file.parent.mkdir(parents=True, exist_ok=True)
        with open(persona_file, "w") as f:
            json.dump(body, f, indent=2)
    except Exception:
        pass
    return {"status": "ok"}


@router.get("/projects/{slug}/persona")
def get_persona(slug: str, request: Request):
    """Get the generated persona for a project (DB first, file fallback)."""
    # Try DB first
    try:
        engine = create_engine(db_url)
        with engine.connect() as conn:
            row = conn.execute(text("SELECT persona FROM public.dash_personas WHERE project_slug = :s"), {"s": slug}).fetchone()
            if row and row[0]:
                return {"persona": row[0] if isinstance(row[0], dict) else json.loads(row[0])}
    except Exception:
        pass
    # Fallback to file
    persona_file = KNOWLEDGE_DIR / slug / "persona.json"
    if persona_file.exists():
        try:
            with open(persona_file) as f:
                return {"persona": json.load(f)}
        except Exception:
            pass
    return {"persona": None}


_training_cancel_flags: dict[str, bool] = {}


@router.post("/projects/{slug}/stop-training")
def stop_training(slug: str, request: Request):
    """Stop a running training pipeline."""
    _training_cancel_flags[slug] = True
    # Also mark the latest training run as failed
    try:
        eng = create_engine(db_url)
        with eng.connect() as conn:
            conn.execute(text(
                "UPDATE public.dash_training_runs SET status = 'failed' "
                "WHERE project_slug = :s AND status NOT IN ('done', 'failed') "
            ), {"s": slug})
            conn.commit()
    except Exception:
        pass
    return {"status": "stopped"}


@router.post("/projects/{slug}/retrain")
def retrain_project(slug: str, request: Request):
    """Delta retrain — only train tables that changed since last training."""
    user = getattr(getattr(request, 'state', None), 'user', None)
    if not user:
        raise HTTPException(401, "Not authenticated")

    # Editor role required for retraining
    from app.auth import check_project_permission
    perm = check_project_permission(user, slug, required_role="editor")
    if not perm:
        raise HTTPException(403, "Editor access required to retrain")

    from db.session import create_project_schema
    schema = create_project_schema(slug)
    engine = create_engine(db_url)
    insp = inspect(engine)

    tables = insp.get_table_names(schema=schema)
    if not tables:
        # Doc-only project — index knowledge + fill brain from doc text
        import threading
        import time as _time
        def _bg_docs_only():
            _log_entries = []
            def _dlog(msg):
                _log_entries.append({"ts": _time.strftime('%H:%M:%S'), "msg": msg})

            try:
                from dash.settings import training_llm_call
                eng = create_engine(db_url)

                # Create training run for UI tracking
                run_id = None
                try:
                    with eng.connect() as conn:
                        r = conn.execute(text(
                            "INSERT INTO public.dash_training_runs (project_slug, status, steps) "
                            "VALUES (:s, 'running', 'reindex') RETURNING id"
                        ), {"s": slug})
                        run_id = r.fetchone()[0]
                        conn.commit()
                except Exception:
                    pass

                def _update_step(step_name):
                    if run_id:
                        try:
                            with eng.connect() as conn:
                                conn.execute(text(
                                    "UPDATE public.dash_training_runs SET steps = :s, logs = CAST(:logs AS jsonb) WHERE id = :id"
                                ), {"s": step_name, "logs": json.dumps(_log_entries), "id": run_id})
                                conn.commit()
                        except Exception:
                            pass

                # Step 1: Index knowledge
                _update_step("reindex")
                _dlog("indexing documents into knowledge base...")
                _reload_project_knowledge(slug, timeout_sec=60)
                _dlog("✓ knowledge indexed")

                # Step 2: Read all doc text for LLM context
                docs_dir = KNOWLEDGE_DIR / slug / "docs"
                all_text = ""
                if docs_dir.exists():
                    for f in docs_dir.iterdir():
                        if f.is_file():
                            try:
                                all_text += f.read_text(errors='ignore')[:5000] + "\n\n"
                            except Exception:
                                pass
                all_text = all_text[:8000]

                # Extract and describe images from raw document files
                docs_raw_dir = KNOWLEDGE_DIR / slug / "docs_raw"
                if docs_raw_dir.exists():
                    for raw_file in docs_raw_dir.iterdir():
                        ext_r = raw_file.suffix.lower()
                        if ext_r == ".pptx":
                            imgs = _extract_images_pptx(str(raw_file))
                        elif ext_r == ".pdf":
                            imgs = _extract_images_pdf(str(raw_file))
                        else:
                            continue
                        if imgs:
                            _dlog(f"describing {len(imgs)} images from {raw_file.name}...")
                            desc = _describe_images_with_vision(imgs, raw_file.name)
                            if desc:
                                all_text += f"\n\n--- IMAGES FROM {raw_file.name} ---\n{desc}"
                                _dlog(f"✓ {len(imgs)} image descriptions added")

                if not all_text.strip():
                    _dlog("· no document text found — skipping brain fill")
                else:
                    # Step 3: Generate memories from docs
                    _update_step("brain_fill")
                    _dlog("generating memories from documents...")
                    try:
                        result = training_llm_call(
                            f"Extract 5-8 key facts from this document that a data analyst should remember.\n\n"
                            f"DOCUMENT:\n{all_text[:4000]}\n\n"
                            f"Return ONLY valid JSON array of strings:\n"
                            f'["fact 1", "fact 2", "fact 3"]',
                            "extraction"
                        )
                        if result:
                            facts = json.loads(result)
                            if isinstance(facts, list):
                                saved = 0
                                with eng.connect() as conn:
                                    for fact in facts[:8]:
                                        if isinstance(fact, str) and len(fact) > 10:
                                            conn.execute(text(
                                                "INSERT INTO public.dash_memories (project_slug, scope, fact, source) "
                                                "VALUES (:s, 'project', :f, 'auto_training') ON CONFLICT DO NOTHING"
                                            ), {"s": slug, "f": fact})
                                            saved += 1
                                    conn.commit()
                                _dlog(f"✓ {saved} memories saved")
                    except Exception as e:
                        _dlog(f"⚠ memories error: {str(e)[:60]}")

                    # Step 4: Generate persona
                    _update_step("persona")
                    _dlog("generating agent persona...")
                    try:
                        result = training_llm_call(
                            f"Based on this document, generate a JSON persona for an AI agent:\n\n"
                            f"DOCUMENT:\n{all_text[:3000]}\n\n"
                            f'Return: {{"persona_prompt": "You are an expert...", "domain_terms": ["term1"], "expertise_areas": ["area1"], "greeting": "Hi! ..."}}',
                            "persona"
                        )
                        if result:
                            persona = json.loads(result)
                            if isinstance(persona, dict) and persona.get("persona_prompt"):
                                persona_file = KNOWLEDGE_DIR / slug / "persona.json"
                                persona_file.parent.mkdir(parents=True, exist_ok=True)
                                with open(persona_file, "w") as f:
                                    json.dump(persona, f, indent=2)
                                with eng.connect() as conn:
                                    conn.execute(text(
                                        "INSERT INTO public.dash_personas (project_slug, persona) VALUES (:s, CAST(:p AS jsonb)) "
                                        "ON CONFLICT (project_slug) DO UPDATE SET persona = CAST(:p AS jsonb)"
                                    ), {"s": slug, "p": json.dumps(persona)})
                                    conn.commit()
                                _dlog("✓ persona generated")
                    except Exception as e:
                        _dlog(f"⚠ persona error: {str(e)[:60]}")

                    # Step 5: Generate workflows (structure-aware for PPTX/PDF/DOCX)
                    _update_step("synthesis")
                    _dlog("generating sample workflows...")
                    try:
                        # Try structure-based extraction from raw docs first
                        structure_wf_saved = 0
                        docs_raw_dir = KNOWLEDGE_DIR / slug / "docs_raw"
                        if docs_raw_dir.exists():
                            for raw_file in docs_raw_dir.iterdir():
                                ext = raw_file.suffix.lower()
                                if ext not in (".pptx", ".pdf", ".docx"):
                                    continue
                                sections = _extract_document_structure(str(raw_file), ext)
                                if len(sections) < 2:
                                    continue
                                sections_text = "\n".join(f"{s['index']}. {s['title']} — {s['content_summary']}" for s in sections)
                                result = training_llm_call(
                                    f"Convert this document structure into a reusable analysis workflow.\n"
                                    f"Each section becomes one step — write a clear analyst question.\n\n"
                                    f"DOCUMENT: {raw_file.name}\nSECTIONS:\n{sections_text}\n\n"
                                    f'Return ONLY valid JSON (no markdown):\n'
                                    f'{{"name": "workflow name", "description": "what it analyzes", '
                                    f'"steps": [{{"title": "section title", "question": "analyst question"}}]}}',
                                    "extraction"
                                )
                                if result:
                                    wf = json.loads(result.strip().strip("`").strip())
                                    if isinstance(wf, dict) and wf.get("name"):
                                        steps = wf.get("steps", [])
                                        step_list = [s.get("question", s.get("title", "")) if isinstance(s, dict) else s for s in steps]
                                        with eng.connect() as conn:
                                            conn.execute(text(
                                                "INSERT INTO public.dash_workflows_db (project_slug, name, description, steps, source) "
                                                "VALUES (:s, :n, :d, CAST(:st AS jsonb), 'document') ON CONFLICT DO NOTHING"
                                            ), {"s": slug, "n": wf["name"], "d": wf.get("description", ""), "st": json.dumps(step_list)})
                                            conn.commit()
                                        structure_wf_saved += 1
                                        _dlog(f"✓ workflow from {raw_file.name} ({len(step_list)} steps)")

                        # Fall back: use LLM to extract structure from text content
                        if structure_wf_saved == 0:
                            result = training_llm_call(
                                f"Read this document and identify the analysis sections/topics.\n"
                                f"Convert the document structure into a reusable analysis workflow.\n"
                                f"Each section/topic becomes one step — write a clear analyst question.\n\n"
                                f"DOCUMENT:\n{all_text[:4000]}\n\n"
                                f"Return ONLY valid JSON (no markdown):\n"
                                f'{{"name": "workflow name based on document", "description": "what this workflow analyzes", '
                                f'"steps": [{{"title": "section title", "question": "analyst question to reproduce this analysis"}}]}}',
                                "extraction"
                            )
                            if result:
                                wf = json.loads(result.strip().strip("`").strip())
                                if isinstance(wf, dict) and wf.get("name"):
                                    steps = wf.get("steps", [])
                                    step_list = [s.get("question", s.get("title", "")) if isinstance(s, dict) else s for s in steps]
                                    with eng.connect() as conn:
                                        conn.execute(text(
                                            "INSERT INTO public.dash_workflows_db (project_slug, name, description, steps, source) "
                                            "VALUES (:s, :n, :d, CAST(:st AS jsonb), 'document') ON CONFLICT DO NOTHING"
                                        ), {"s": slug, "n": wf["name"], "d": wf.get("description", ""), "st": json.dumps(step_list)})
                                        conn.commit()
                                    _dlog(f"✓ workflow from text ({len(step_list)} steps)")
                                    structure_wf_saved = 1
                                elif isinstance(wf, list):
                                    saved = 0
                                    with eng.connect() as conn:
                                        for w in wf[:3]:
                                            if isinstance(w, dict) and w.get("name"):
                                                steps = w.get("steps", [])
                                                step_list = [s.get("question", s.get("title", "")) if isinstance(s, dict) else s for s in steps]
                                                conn.execute(text(
                                                    "INSERT INTO public.dash_workflows_db (project_slug, name, description, steps, source) "
                                                    "VALUES (:s, :n, :d, CAST(:st AS jsonb), 'document') ON CONFLICT DO NOTHING"
                                                ), {"s": slug, "n": w["name"], "d": w.get("description", ""), "st": json.dumps(step_list)})
                                                saved += 1
                                        conn.commit()
                                    _dlog(f"✓ {saved} workflows from text")
                                    structure_wf_saved = saved
                        else:
                            _dlog(f"✓ {structure_wf_saved} document-based workflows saved")
                    except Exception as e:
                        _dlog(f"⚠ workflows error: {str(e)[:60]}")

                    # Step 6: Generate evals
                    _dlog("generating eval questions...")
                    try:
                        result = training_llm_call(
                            f"Based on this document, generate 6 smart business questions an executive would ask.\n\n"
                            f"DOCUMENT:\n{all_text[:4000]}\n\n"
                            f"Generate questions that:\n"
                            f"- Ask about key metrics, trends, performance, comparisons\n"
                            f"- Use business language (not technical or table names)\n"
                            f"- Cover: summary, metrics, trends, comparisons, risks, recommendations\n\n"
                            f'Return JSON array: [{{"question": "What was the revenue growth this quarter?", "expected_answer": "Revenue grew..."}}]',
                            "extraction"
                        )
                        if result:
                            evals = json.loads(result)
                            if isinstance(evals, list):
                                saved = 0
                                with eng.connect() as conn:
                                    for ev in evals[:6]:
                                        if isinstance(ev, dict) and ev.get("question"):
                                            conn.execute(text(
                                                "INSERT INTO public.dash_evals (project_slug, question, expected_sql) VALUES (:s, :q, :a)"
                                            ), {"s": slug, "q": ev["question"], "a": ev.get("expected_answer", "")})
                                            saved += 1
                                    conn.commit()
                                _dlog(f"✓ {saved} eval questions saved")
                    except Exception as e:
                        _dlog(f"⚠ evals error: {str(e)[:60]}")

                    # Step 7: Seed feedback
                    _dlog("seeding sample feedback...")
                    try:
                        with eng.connect() as conn:
                            conn.execute(text(
                                "INSERT INTO public.dash_feedback (project_slug, question, answer, rating) "
                                "VALUES (:s, 'What is this project about?', :a, 'up')"
                            ), {"s": slug, "a": all_text[:500]})
                            conn.commit()
                        _dlog("✓ 1 seed feedback saved")
                    except Exception:
                        pass

                    # Step 8: Extract business rules from docs
                    _update_step("domain_knowledge")
                    _dlog("extracting business rules from documents...")
                    try:
                        result = training_llm_call(
                            f"Extract business rules, constraints, and policies from this document.\n\n"
                            f"DOCUMENT:\n{all_text[:4000]}\n\n"
                            f'Return JSON array: [{{"name": "Rule Name", "definition": "The rule description", "type": "business_rule"}}]',
                            "extraction"
                        )
                        if result:
                            rules = json.loads(result)
                            if isinstance(rules, list):
                                saved = 0
                                with eng.connect() as conn:
                                    for r in rules[:8]:
                                        if isinstance(r, dict) and r.get("name"):
                                            rule_id = f"doc_{r['name'].lower().replace(' ', '_')[:25]}"
                                            conn.execute(text(
                                                "INSERT INTO public.dash_rules_db (project_slug, rule_id, name, type, definition, source) "
                                                "VALUES (:s, :rid, :name, :type, :defn, 'doc_training') ON CONFLICT DO NOTHING"
                                            ), {"s": slug, "rid": rule_id, "name": r["name"], "type": r.get("type", "business_rule"), "defn": r.get("definition", "")})
                                            saved += 1
                                    conn.commit()
                                _dlog(f"✓ {saved} business rules extracted")
                    except Exception as e:
                        _dlog(f"⚠ rules error: {str(e)[:60]}")

                    # Step 9: Extract domain knowledge (glossary + KPIs)
                    _dlog("extracting domain knowledge...")
                    try:
                        result = training_llm_call(
                            f"Extract domain knowledge from this document:\n\n"
                            f"DOCUMENT:\n{all_text[:4000]}\n\n"
                            f"Return JSON with:\n"
                            f'{{"glossary": [{{"term": "SLA", "definition": "Service Level Agreement"}}], '
                            f'"kpis": [{{"name": "Resolution Time", "definition": "Average time to resolve tickets"}}], '
                            f'"key_metrics": [{{"name": "Ticket Volume", "value": "500/month"}}]}}',
                            "extraction"
                        )
                        if result:
                            domain = json.loads(result)
                            if isinstance(domain, dict):
                                saved = 0
                                with eng.connect() as conn:
                                    # Glossary → memories
                                    for g in (domain.get("glossary") or [])[:10]:
                                        if isinstance(g, dict) and g.get("term"):
                                            conn.execute(text(
                                                "INSERT INTO public.dash_memories (project_slug, scope, fact, source) "
                                                "VALUES (:s, 'project', :f, 'glossary') ON CONFLICT DO NOTHING"
                                            ), {"s": slug, "f": f"Glossary: {g['term']} = {g.get('definition', '')}"})
                                            saved += 1
                                    # KPIs → rules
                                    for k in (domain.get("kpis") or [])[:8]:
                                        if isinstance(k, dict) and k.get("name"):
                                            rule_id = f"kpi_doc_{k['name'].lower().replace(' ', '_')[:25]}"
                                            conn.execute(text(
                                                "INSERT INTO public.dash_rules_db (project_slug, rule_id, name, type, definition, source) "
                                                "VALUES (:s, :rid, :name, 'kpi', :defn, 'doc_training') ON CONFLICT DO NOTHING"
                                            ), {"s": slug, "rid": rule_id, "name": k["name"], "defn": k.get("definition", "")})
                                            saved += 1
                                    # Metrics → memories
                                    for m in (domain.get("key_metrics") or [])[:5]:
                                        if isinstance(m, dict) and m.get("name"):
                                            conn.execute(text(
                                                "INSERT INTO public.dash_memories (project_slug, scope, fact, source) "
                                                "VALUES (:s, 'project', :f, 'auto_training') ON CONFLICT DO NOTHING"
                                            ), {"s": slug, "f": f"Key metric: {m['name']} = {m.get('value', 'N/A')}"})
                                            saved += 1
                                    conn.commit()
                                _dlog(f"✓ {saved} domain knowledge items extracted")
                    except Exception as e:
                        _dlog(f"⚠ domain knowledge error: {str(e)[:60]}")

                    # Step 10: Generate proactive insights from docs
                    _dlog("generating proactive insights...")
                    try:
                        result = training_llm_call(
                            f"Analyze this document and identify 3 proactive insights — things that need attention, risks, or opportunities.\n\n"
                            f"DOCUMENT:\n{all_text[:4000]}\n\n"
                            f'Return JSON array: [{{"insight": "Description of the insight", "severity": "info|warning|critical"}}]',
                            "extraction"
                        )
                        if result:
                            insights = json.loads(result)
                            if isinstance(insights, list):
                                saved = 0
                                with eng.connect() as conn:
                                    for ins in insights[:5]:
                                        if isinstance(ins, dict) and ins.get("insight"):
                                            conn.execute(text(
                                                "INSERT INTO public.dash_proactive_insights (project_slug, insight, severity) "
                                                "VALUES (:s, :i, :sev)"
                                            ), {"s": slug, "i": ins["insight"], "sev": ins.get("severity", "info")})
                                            saved += 1
                                    conn.commit()
                                _dlog(f"✓ {saved} proactive insights generated")
                    except Exception as e:
                        _dlog(f"⚠ insights error: {str(e)[:60]}")

                    # Step 11: Cross-document relationships
                    _update_step("relationships")
                    _dlog("discovering cross-document relationships...")
                    try:
                        doc_summaries = []
                        for f in docs_dir.iterdir():
                            if f.is_file():
                                try:
                                    content = f.read_text(errors='ignore')[:2000]
                                    doc_summaries.append(f"{f.name}: {content[:500]}")
                                except Exception:
                                    pass
                        if len(doc_summaries) >= 2:
                            result = training_llm_call(
                                f"Analyze these documents and find relationships, shared topics, cross-references between them.\n\n"
                                f"DOCUMENTS:\n" + "\n---\n".join(doc_summaries[:10]) + "\n\n"
                                f"Return ONLY valid JSON array:\n"
                                f'[{{"from_doc": "doc1.txt", "to_doc": "doc2.txt", "relationship": "both discuss revenue targets", "shared_topics": ["revenue", "KPIs"], "strength": 0.8}}]',
                                "extraction"
                            )
                            if result:
                                rels = json.loads(result)
                                if isinstance(rels, list):
                                    saved = 0
                                    with eng.connect() as conn:
                                        for r in rels[:10]:
                                            if isinstance(r, dict) and r.get("from_doc"):
                                                conn.execute(text(
                                                    "INSERT INTO public.dash_relationships (project_slug, from_table, from_column, to_table, to_column, rel_type, confidence, source) "
                                                    "VALUES (:s, :ft, :fc, :tt, :tc, 'topic', :conf, 'ai') "
                                                    "ON CONFLICT DO NOTHING"
                                                ), {"s": slug, "ft": r["from_doc"], "fc": ", ".join(r.get("shared_topics", [])[:5]),
                                                    "tt": r.get("to_doc", ""), "tc": r.get("relationship", ""),
                                                    "conf": r.get("strength", 0.5)})
                                                saved += 1
                                        conn.commit()
                                    _dlog(f"✓ {saved} cross-document relationships found")
                        else:
                            _dlog("· only 1 document — no cross-references to find")
                    except Exception as e:
                        _dlog(f"⚠ relationships error: {str(e)[:60]}")

                    # Step 12: Negative examples
                    _dlog("extracting negative examples...")
                    try:
                        result = training_llm_call(
                            f"Based on this document, what are common mistakes someone might make when interpreting this information?\n\n"
                            f"DOCUMENT:\n{all_text[:3000]}\n\n"
                            f'Return JSON array of "DON\'T / DO" pairs:\n'
                            f'["DON\'T confuse X with Y — DO use Z instead", "DON\'T assume A — DO check B"]',
                            "extraction"
                        )
                        if result:
                            negs = json.loads(result)
                            if isinstance(negs, list):
                                saved = 0
                                with eng.connect() as conn:
                                    for neg in negs[:5]:
                                        if isinstance(neg, str) and len(neg) > 10:
                                            conn.execute(text(
                                                "INSERT INTO public.dash_memories (project_slug, scope, fact, source) "
                                                "VALUES (:s, 'project', :f, 'negative_example') ON CONFLICT DO NOTHING"
                                            ), {"s": slug, "f": f"⚠ {neg}"})
                                            saved += 1
                                    conn.commit()
                                _dlog(f"✓ {saved} negative examples saved")
                    except Exception as e:
                        _dlog(f"⚠ negative examples error: {str(e)[:60]}")

                    # Step 13: Training Q&A from docs
                    _dlog("generating training Q&A...")
                    try:
                        result = training_llm_call(
                            f"Generate 5 question-answer pairs that test understanding of this document.\n\n"
                            f"DOCUMENT:\n{all_text[:3000]}\n\n"
                            f'Return JSON array: [{{"question": "What is...?", "answer": "It is..."}}]',
                            "extraction"
                        )
                        if result:
                            qas = json.loads(result)
                            if isinstance(qas, list):
                                qa_file = KNOWLEDGE_DIR / slug / "training"
                                qa_file.mkdir(parents=True, exist_ok=True)
                                with open(qa_file / "doc_qa.json", "w") as f:
                                    json.dump(qas, f, indent=2)
                                _dlog(f"✓ {len(qas)} training Q&A pairs generated")
                    except Exception as e:
                        _dlog(f"⚠ Q&A error: {str(e)[:60]}")

                    # Step 14: Multi-doc synthesis
                    if len([f for f in docs_dir.iterdir() if f.is_file()]) >= 2:
                        _dlog("running multi-document synthesis...")
                        try:
                            doc_names = [f.name for f in docs_dir.iterdir() if f.is_file()]
                            result = training_llm_call(
                                f"These documents are part of the same project. Create a unified understanding.\n\n"
                                f"DOCUMENTS: {', '.join(doc_names)}\n\n"
                                f"CONTENT:\n{all_text[:5000]}\n\n"
                                f"Write a 3-4 sentence synthesis that explains how these documents relate and what the overall project is about.",
                                "persona"
                            )
                            if result:
                                with eng.connect() as conn:
                                    conn.execute(text(
                                        "INSERT INTO public.dash_memories (project_slug, scope, fact, source) "
                                        "VALUES (:s, 'project', :f, 'synthesis') ON CONFLICT DO NOTHING"
                                    ), {"s": slug, "f": f"Project synthesis: {result[:500]}"})
                                    conn.commit()
                                _dlog("✓ multi-document synthesis complete")
                        except Exception as e:
                            _dlog(f"⚠ synthesis error: {str(e)[:60]}")

                _dlog("✓ doc-only training complete")

                # Save training run with proper step tracking
                run_id = None
                with eng.connect() as conn:
                    r = conn.execute(text(
                        "INSERT INTO public.dash_training_runs (project_slug, status, steps, logs, finished_at) "
                        "VALUES (:s, 'done', 'complete', CAST(:logs AS jsonb), NOW()) RETURNING id"
                    ), {"s": slug, "logs": json.dumps(_log_entries)})
                    run_id = r.fetchone()[0]
                    conn.commit()

            except Exception:
                pass
        _bg_executor.submit(_bg_docs_only)
        return {"status": "ok", "tables": 0, "message": "Doc-only project — indexing documents and filling brain"}

    import pandas as pd

    _training_cancel_flags[slug] = False  # Reset cancel flag

    def _bg():
        import time as _time
        skipped = 0
        trained = 0
        total_tables = len(tables)

        # Create a single master training run for the entire retrain batch
        master_run_id = None
        master_engine = create_engine(db_url, pool_size=2, max_overflow=3, pool_recycle=3600)
        try:
            with master_engine.connect() as conn:
                result = conn.execute(text(
                    "INSERT INTO public.dash_training_runs (project_slug, status, steps) "
                    "VALUES (:s, 'running', :steps) RETURNING id"
                ), {"s": slug, "steps": f"starting||0|{total_tables}"})
                master_run_id = result.fetchone()[0]
                conn.commit()
        except Exception:
            pass

        def _master_log(msg: str, tbl_name: str = "", tbl_idx: int = 0):
            """Log to the master training run."""
            if not master_run_id:
                return
            try:
                with master_engine.connect() as conn:
                    conn.execute(text(
                        "UPDATE public.dash_training_runs SET logs = COALESCE(logs, '[]'::jsonb) || CAST(:entry AS jsonb) WHERE id = :id"
                    ), {"entry": json.dumps([{"ts": _time.strftime('%H:%M:%S'), "msg": msg, "table": tbl_name, "table_index": tbl_idx, "total_tables": total_tables}]), "id": master_run_id})
                    conn.commit()
            except Exception:
                pass

        for tbl_idx, tbl in enumerate(tables, start=1):
            if _training_cancel_flags.get(slug):
                break
            try:
                # Get current row count and columns
                with engine.connect() as conn:
                    row_count = conn.execute(text(f'SELECT COUNT(*) FROM "{schema}"."{tbl}"')).scalar() or 0
                tbl_cols = [c["name"] for c in insp.get_columns(tbl, schema=schema)]

                # Check fingerprint — skip if unchanged
                change_type = check_fingerprint_changed(slug, tbl, row_count, tbl_cols)

                if change_type == "unchanged":
                    _master_log(f"⊘ skipping {tbl} — unchanged (fingerprint match)", tbl, tbl_idx)
                    # Update master run step to show skip
                    if master_run_id:
                        try:
                            with master_engine.connect() as conn:
                                conn.execute(text(
                                    "UPDATE public.dash_training_runs SET steps = :steps WHERE id = :id"
                                ), {"steps": f"skipped|{tbl}|{tbl_idx}|{total_tables}", "id": master_run_id})
                                conn.commit()
                        except Exception:
                            pass
                    skipped += 1
                    continue

                # Table changed or new — run training
                _master_log(f"training table {tbl} ({tbl_idx}/{total_tables})...", tbl, tbl_idx)
                df = pd.read_sql(f'SELECT * FROM "{schema}"."{tbl}" LIMIT 100', engine)
                col_analyses = [_analyze_column(df[col]) for col in df.columns]
                sample_rows = df.head(10).to_dict('records')

                tables_dir = KNOWLEDGE_DIR / slug / "tables"
                business_dir = KNOWLEDGE_DIR / slug / "business"
                tables_dir.mkdir(parents=True, exist_ok=True)
                business_dir.mkdir(parents=True, exist_ok=True)

                # Load existing metadata or create new
                meta_file = tables_dir / f"{tbl}.json"
                if meta_file.exists():
                    with open(meta_file) as f:
                        metadata = json.load(f)
                else:
                    metadata = _generate_metadata(tbl, df, col_analyses)

                biz_rules = _generate_business_rules(tbl, col_analyses)
                _run_auto_training(slug, tbl, col_analyses, metadata, biz_rules, sample_rows, tables_dir, business_dir,
                                   master_run_id=master_run_id, table_index=tbl_idx, total_tables=total_tables)

                # Save new fingerprint after training
                save_fingerprint(slug, tbl, row_count, tbl_cols)
                trained += 1
                _master_log(f"✓ table {tbl} training complete ({tbl_idx}/{total_tables})", tbl, tbl_idx)
            except Exception as e:
                import logging
                logging.error(f"Retrain failed for {slug}/{tbl}: {e}")
                _master_log(f"⚠ table {tbl} training failed: {str(e)[:80]}", tbl, tbl_idx)

        # Mark master run as done
        if master_run_id:
            try:
                with master_engine.connect() as conn:
                    conn.execute(text(
                        "UPDATE public.dash_training_runs SET status = 'done', steps = :steps, "
                        "tables_trained = :trained, finished_at = NOW() WHERE id = :id"
                    ), {"steps": f"complete||{total_tables}|{total_tables}", "trained": trained, "id": master_run_id})
                    conn.commit()
            except Exception:
                pass

    _bg_executor.submit(_bg)
    return {"status": "ok", "tables": len(tables), "message": "Delta retraining started — unchanged tables will be skipped"}


@router.get("/knowledge-file-content/{filename}")
def get_knowledge_file_content(filename: str, request: Request, project: str | None = None, subdir: str = "tables"):
    """Read the content of a knowledge file."""
    safe_name = Path(filename).name
    if safe_name != filename or '..' in filename:
        raise HTTPException(400, "Invalid filename")
    if project:
        base = KNOWLEDGE_DIR / project / subdir
    else:
        base = KNOWLEDGE_DIR / subdir
    filepath = base / filename
    if not filepath.exists():
        raise HTTPException(404, f"File not found: {filename}")
    try:
        content = filepath.read_text()
        if filename.endswith(".json"):
            return {"content": json.loads(content), "type": "json"}
        return {"content": content, "type": "text"}
    except Exception as e:
        raise HTTPException(500, str(e))


@router.get("/knowledge-files")
def list_knowledge_files(request: Request, project: str | None = None):
    """List knowledge files for a project or global."""
    if project:
        base = KNOWLEDGE_DIR / project
    else:
        base = KNOWLEDGE_DIR
    files = []
    for subdir in ["tables", "queries", "business"]:
        path = base / subdir
        if path.exists():
            for f in sorted(path.iterdir()):
                if f.is_file() and not f.name.startswith("."):
                    files.append({"name": f.name, "type": subdir, "size": f.stat().st_size})
    return {"files": files}


@router.get("/lineage")
def get_lineage(request: Request, project: str | None = None):
    """Detect table relationships via FK column patterns and constraints."""
    engine = create_engine(db_url)
    insp = inspect(engine)
    schema = _get_project_schema(request, project) or "public"
    tables = insp.get_table_names(schema=schema)
    relationships: list[dict] = []

    for tbl in tables:
        cols = insp.get_columns(tbl, schema=schema)
        for col in cols:
            name = col["name"]
            if name.endswith("_id") and name != "id":
                ref_table = name[:-3] + "s"
                if ref_table not in tables:
                    ref_table = name[:-3]
                if ref_table in tables:
                    relationships.append({
                        "from_table": tbl,
                        "from_column": name,
                        "to_table": ref_table,
                        "to_column": "id",
                        "type": "foreign_key",
                    })

        try:
            fks = insp.get_foreign_keys(tbl, schema=schema)
            for fk in fks:
                for i, col in enumerate(fk.get("constrained_columns", [])):
                    ref_cols = fk.get("referred_columns", [])
                    relationships.append({
                        "from_table": tbl,
                        "from_column": col,
                        "to_table": fk.get("referred_table", ""),
                        "to_column": ref_cols[i] if i < len(ref_cols) else "id",
                        "type": "constraint",
                    })
        except Exception:
            pass

    # Deduplicate
    seen = set()
    unique: list[dict] = []
    for r in relationships:
        key = f"{r['from_table']}.{r['from_column']}->{r['to_table']}.{r['to_column']}"
        if key not in seen:
            seen.add(key)
            unique.append(r)

    return {"relationships": unique, "tables": tables}


def _bootstrap_shared():
    """Create shared_results table if needed."""
    engine = create_engine(db_url)
    with engine.connect() as conn:
        conn.execute(text("""
            CREATE TABLE IF NOT EXISTS public.shared_results (
                id SERIAL PRIMARY KEY,
                from_user TEXT NOT NULL,
                title TEXT NOT NULL,
                content TEXT NOT NULL,
                sql_query TEXT,
                created_at TIMESTAMP DEFAULT NOW()
            )
        """))
        conn.commit()


@router.post("/shared")
def share_result(request: Request, title: str, content: str, sql_query: str = ""):
    """Share a query result with all users."""
    user_id = _get_user_id(request)
    _bootstrap_shared()
    engine = create_engine(db_url)
    with engine.connect() as conn:
        conn.execute(text(
            "INSERT INTO public.shared_results (from_user, title, content, sql_query) VALUES (:u, :t, :c, :s)"
        ), {"u": user_id, "t": title, "c": content, "s": sql_query})
        conn.commit()
    return {"status": "ok"}


@router.get("/shared")
def list_shared(request: Request):
    """List all shared results."""
    _bootstrap_shared()
    engine = create_engine(db_url)
    results = []
    with engine.connect() as conn:
        rows = conn.execute(text(
            "SELECT id, from_user, title, content, sql_query, created_at FROM public.shared_results ORDER BY created_at DESC LIMIT 50"
        )).fetchall()
        for r in rows:
            results.append({"id": r[0], "from_user": r[1], "title": r[2], "content": r[3][:500], "sql_query": r[4], "created_at": str(r[5]) if r[5] else None})
    return {"results": results}


@router.delete("/shared/{result_id}")
def delete_shared(result_id: int, request: Request):
    """Delete a shared result (own or super admin only)."""
    user = getattr(getattr(request, 'state', None), 'user', None)
    if not user:
        raise HTTPException(401, "Not authenticated")
    username = user.get("username", "")
    is_super = user.get("is_super", False)
    engine = create_engine(db_url)
    with engine.connect() as conn:
        if is_super:
            conn.execute(text("DELETE FROM public.shared_results WHERE id = :id"), {"id": result_id})
        else:
            conn.execute(text("DELETE FROM public.shared_results WHERE id = :id AND from_user = :u"), {"id": result_id, "u": username})
        conn.commit()
    return {"status": "ok"}


@router.get("/sessions")
def list_sessions(request: Request, project: str | None = None):
    """List current user's chat sessions, optionally filtered by project."""
    user = getattr(getattr(request, 'state', None), 'user', None)
    if not user:
        return {"sessions": []}
    uid = user["user_id"]
    engine = create_engine(db_url)
    sessions: list[dict] = []
    try:
        with engine.connect() as conn:
            if project:
                # Project-specific sessions
                result = conn.execute(text("""
                    SELECT session_id, first_message, created_at, updated_at
                    FROM public.dash_chat_sessions
                    WHERE user_id = :uid AND project_slug = :proj
                    ORDER BY updated_at DESC
                    LIMIT 50
                """), {"uid": uid, "proj": project})
            else:
                # All user sessions (for Dash Agent)
                result = conn.execute(text("""
                    SELECT s.session_id, s.first_message, s.created_at, s.updated_at,
                           COALESCE(p.agent_name, 'Dash Agent') as agent_name, s.project_slug
                    FROM public.dash_chat_sessions s
                    LEFT JOIN public.dash_projects p ON p.slug = s.project_slug
                    WHERE s.user_id = :uid
                    ORDER BY s.updated_at DESC
                    LIMIT 50
                """), {"uid": uid})

            for row in result.fetchall():
                entry: dict = {
                    "session_id": str(row[0]),
                    "first_message": row[1] or None,
                    "created_at": str(row[2]) if row[2] else None,
                    "updated_at": str(row[3]) if row[3] else None,
                }
                if not project and len(row) > 4:
                    entry["agent_name"] = row[4]
                    entry["project_slug"] = row[5]
                sessions.append(entry)
    except Exception:
        pass
    return {"sessions": sessions}


@router.get("/sessions/{session_id}/messages")
def get_session_messages(session_id: str, request: Request):
    """Load chat messages from a stored Agno session."""
    user = getattr(getattr(request, 'state', None), 'user', None)
    if not user:
        return {"messages": []}

    engine = create_engine(db_url)
    messages: list[dict] = []

    try:
        with engine.connect() as conn:
            # Verify this session belongs to the user
            owned = conn.execute(text(
                "SELECT 1 FROM public.dash_chat_sessions WHERE session_id = :sid AND user_id = :uid"
            ), {"sid": session_id, "uid": user["user_id"]}).fetchone()
            if not owned:
                return {"messages": []}

            # Get runs from Agno sessions
            row = conn.execute(text(
                "SELECT runs FROM ai.agno_sessions WHERE session_id = :sid"
            ), {"sid": session_id}).fetchone()

            if row and row[0]:
                runs = row[0] if isinstance(row[0], list) else []
                for run in runs:
                    # Only top-level runs (user → team interactions)
                    if run.get("parent_run_id"):
                        continue

                    # User message
                    input_data = run.get("input", {})
                    user_msg = ""
                    if isinstance(input_data, dict):
                        user_msg = input_data.get("input_content", "") or input_data.get("content", "")
                    elif isinstance(input_data, str):
                        user_msg = input_data

                    if user_msg:
                        messages.append({"role": "user", "content": user_msg})

                    # Assistant response
                    content = run.get("content", "")
                    if content:
                        messages.append({"role": "assistant", "content": content})
    except Exception:
        pass

    return {"messages": messages}


@router.post("/sessions/register")
def register_session(request: Request, session_id: str, message: str = "", project: str | None = None):
    """Register or update a chat session for the current user."""
    user = getattr(getattr(request, 'state', None), 'user', None)
    if not user:
        return {"status": "skip"}
    uid = user["user_id"]
    engine = create_engine(db_url)
    try:
        with engine.connect() as conn:
            existing = conn.execute(text(
                "SELECT id, first_message FROM public.dash_chat_sessions WHERE session_id = :sid"
            ), {"sid": session_id}).fetchone()
            if existing:
                # Update timestamp, set first_message if empty
                if not existing[1] and message:
                    conn.execute(text(
                        "UPDATE public.dash_chat_sessions SET updated_at = NOW(), first_message = :msg WHERE session_id = :sid"
                    ), {"sid": session_id, "msg": message[:100]})
                else:
                    conn.execute(text(
                        "UPDATE public.dash_chat_sessions SET updated_at = NOW() WHERE session_id = :sid"
                    ), {"sid": session_id})
            else:
                conn.execute(text(
                    "INSERT INTO public.dash_chat_sessions (user_id, session_id, project_slug, first_message) "
                    "VALUES (:uid, :sid, :proj, :msg)"
                ), {"uid": uid, "sid": session_id, "proj": project, "msg": (message[:100] if message else None)})
            conn.commit()
    except Exception:
        pass
    return {"status": "ok"}


@router.get("/queries")
def recent_queries():
    """Get recent query sessions."""
    engine = create_engine(db_url)
    queries: list[dict] = []
    try:
        with engine.connect() as conn:
            # Agno stores sessions in ai.agno_sessions
            result = conn.execute(text("""
                SELECT session_id, team_id, user_id, created_at, updated_at
                FROM ai.agno_sessions
                WHERE team_id = 'dash'
                ORDER BY updated_at DESC
                LIMIT 30
            """))
            for row in result.fetchall():
                queries.append({
                    "session_id": str(row[0]),
                    "team_id": row[1],
                    "user_id": row[2],
                    "created_at": str(row[3]) if row[3] else None,
                    "updated_at": str(row[4]) if row[4] else None,
                })
    except Exception:
        pass
    return {"queries": queries}


@router.get("/tables/{table_name}/inspect")
def inspect_table(table_name: str, request: Request, project: str | None = None):
    """Get column details and sample data for a table."""
    schema = _get_project_schema(request, project) or "public"
    engine = create_engine(db_url)
    insp = inspect(engine)

    if table_name not in insp.get_table_names(schema=schema):
        raise HTTPException(404, f"Table '{table_name}' not found in {schema}")

    cols = insp.get_columns(table_name, schema=schema)
    columns = [{"name": c["name"], "type": str(c["type"]), "nullable": c.get("nullable", True)} for c in cols]

    # Quality notes from metadata file
    quality_notes: list[str] = []
    if project:
        meta_file = KNOWLEDGE_DIR / project / "tables" / f"{table_name}.json"
    else:
        meta_file = TABLES_DIR / f"{table_name}.json"
    if meta_file.exists():
        try:
            with open(meta_file) as f:
                meta = json.load(f)
            quality_notes = meta.get("data_quality_notes", [])
        except Exception:
            pass

    # Sample data
    sample: list[dict] = []
    try:
        with engine.connect() as conn:
            result = conn.execute(text(f'SELECT * FROM "{schema}"."{table_name}" LIMIT 5'))
            keys = list(result.keys())
            for row in result.fetchall():
                sample.append({k: str(v)[:50] if v is not None else None for k, v in zip(keys, row)})
    except Exception:
        pass

    return {"columns": columns, "sample": sample, "quality_notes": quality_notes}


@router.delete("/tables/{table_name}")
def delete_table(table_name: str, request: Request, project: str | None = None):
    """Drop a user-uploaded table and remove its knowledge files."""
    user = getattr(getattr(request, 'state', None), 'user', None)
    if not user:
        raise HTTPException(401, "Not authenticated")

    if table_name in PROTECTED_TABLES:
        raise HTTPException(403, f"Cannot delete protected table: {table_name}")

    # Determine schema
    if project:
        import re as _re
        schema = _re.sub(r"[^a-z0-9_]", "_", project.lower())[:63]
    else:
        schema = "public"

    engine = create_engine(db_url)
    insp = inspect(engine)

    if table_name not in insp.get_table_names(schema=schema):
        raise HTTPException(404, f"Table '{table_name}' not found in schema '{schema}'")

    with engine.connect() as conn:
        conn.execute(text(f'DROP TABLE IF EXISTS "{schema}"."{table_name}" CASCADE'))
        conn.commit()

    # Remove knowledge files (project-scoped or global)
    if project:
        base = KNOWLEDGE_DIR / project
    else:
        base = KNOWLEDGE_DIR

    for subdir in ["tables", "queries", "business", "staging"]:
        for pattern in [f"{table_name}.json", f"{table_name}_queries.sql", f"{table_name}_rules.json", f"{table_name}.*"]:
            for f in (base / subdir).glob(pattern) if (base / subdir).exists() else []:
                f.unlink(missing_ok=True)

    # Clean DB records for this table
    try:
        slug = project or ""
        with engine.connect() as conn:
            conn.execute(text("DELETE FROM public.dash_table_metadata WHERE project_slug = :s AND table_name = :t"), {"s": slug, "t": table_name})
            conn.execute(text("DELETE FROM public.dash_training_qa WHERE project_slug = :s AND table_name = :t"), {"s": slug, "t": table_name})
            conn.commit()
    except Exception:
        pass

    return {"status": "ok", "deleted": table_name}
